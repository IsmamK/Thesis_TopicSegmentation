"""
fix_slides_v10.py — Rebuild broken slides with correct 10" x 5.625" coordinates.

Fixes slides: 6 (Pipeline), 7 (How Seg Works), 8 (N1-N4), 13 (Oracle), 14 (Granularity), 22 (Final Verdict)
Keeps all other slides from v6 unchanged (including slide 18 External Validation).

Output: defense/LecSeg_Defense_Final_v10.pptx
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

INPUT  = "defense/LecSeg_Defense_Final_v6.pptx"
OUTPUT = "defense/LecSeg_Defense_Final_v10b.pptx"
FIG    = "defense/figures"

prs = Presentation(INPUT)

# Actual slide dimensions
W = prs.slide_width  / 914400   # 10.0 inches
H = prs.slide_height / 914400   # 5.625 inches
M = 0.2                          # margin

TITLE_T = 0.12
TITLE_H = 0.50
RULE_Y  = TITLE_T + TITLE_H + 0.03   # 0.65
BODY_T  = RULE_Y + 0.08              # 0.73
BODY_H  = H - BODY_T - 0.08          # 4.815
BODY_W  = W - 2 * M                  # 9.6

# Colours
BG      = RGBColor( 10,  15,  30)
C_WHITE = RGBColor(255, 255, 255)
C_GRY   = RGBColor(190, 195, 215)
C_HL    = RGBColor(255, 210,  70)
C_GRN   = RGBColor( 80, 200, 120)
C_BLU   = RGBColor(100, 180, 255)
C_RED   = RGBColor(240, 100,  80)
C_SUB   = RGBColor(180, 200, 255)

# ─────────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────────

def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_title(slide, text, size=20):
    txb = slide.shapes.add_textbox(
        Inches(M), Inches(TITLE_T), Inches(BODY_W), Inches(TITLE_H))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

def add_rule(slide):
    from pptx.oxml import parse_xml
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvSpPr><p:cNvPr id="997" name="rule"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{int(Inches(M))}" y="{int(Inches(RULE_Y))}"/>'
        f'<a:ext cx="{int(Inches(BODY_W))}" cy="0"/></a:xfrm>'
        '<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        '<a:ln w="9525"><a:solidFill><a:srgbClr val="3a3f5c"/></a:solidFill></a:ln>'
        '</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    try: slide.shapes._spTree.append(parse_xml(xml))
    except Exception: pass

def rect(slide, x, y, w, h, fill_rgb, line_rgb=(50,60,100), rounded=False):
    shape_type = 5 if rounded else 1
    shp = slide.shapes.add_shape(shape_type,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shp.line.color.rgb = RGBColor(*line_rgb)
    shp.line.width = Pt(0.75)
    shp.text_frame.word_wrap = True
    return shp

def rect_text(slide, x, y, w, h, fill_rgb, text, size=9, bold=False,
              color=C_WHITE, line_rgb=(50,60,100), rounded=False, align=PP_ALIGN.CENTER):
    shp = rect(slide, x, y, w, h, fill_rgb, line_rgb, rounded)
    tf = shp.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shp

def txbox(slide, text, x, y, w, h, size=10, bold=False,
          color=C_GRY, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def multilines(slide, items, x, y, w, h, ds=10, dc=None):
    if dc is None: dc = C_GRY
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            txt, sz, bd, col = item, ds, False, dc
        else:
            txt = item[0]
            sz  = item[1] if len(item) > 1 else ds
            bd  = item[2] if len(item) > 2 else False
            col = item[3] if len(item) > 3 else dc
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = col
        p.space_after = Pt(1)

def img(slide, path, x, y, w, h):
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def arrow(slide, x, y, color=C_GRY):
    txbox(slide, "▶", x, y, 0.18, 0.25, size=10, color=color, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 6 (idx 5) — PIPELINE
# Slide is 10 x 5.625 inches
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides[5]
clear_slide(s); set_bg(s)
add_title(s, "End-to-End Pipeline — Video In, Chapter Timestamps Out")
add_rule(s)

# ── 7 pipeline step boxes (y=0.76, h=0.78) ─────────────────────────────
PY = 0.76
PH = 0.78
BW = 1.18   # box width
GA = 0.10   # gap between boxes
pipe_total = 7*BW + 6*GA   # 8.26+0.6=8.86
px0 = (W - pipe_total) / 2  # ~0.57

pipe_steps = [
    ("VIDEO\nInput",                  (25,  35,  70)),
    ("WHISPER\nTranscription",        (18,  55, 100)),
    ("spaCy\nSentences",              (18,  55, 100)),
    ("BGE + E5\nEmbeddings",          (15,  80, 110)),
    ("Cosine Score\nCurve",           (20,  95,  60)),
    ("Candidate\nPool (~40/hr)",      (75,  35, 100)),
    ("Output\nTimestamps",            (110, 65,  15)),
]
for i, (lbl, col) in enumerate(pipe_steps):
    bx = px0 + i*(BW+GA)
    rect_text(s, bx, PY, BW, PH, col, lbl, size=8, bold=False)
    if i < 6:
        arrow(s, bx+BW+GA/2-0.09, PY+PH/2-0.12, color=C_GRY)

# ── Method Selector banner (y=1.62, h=0.38) ────────────────────────────
rect_text(s, M, 1.62, BODY_W, 0.38, (0,60,60),
    "METHOD SELECTOR  —  runs BEFORE segmentation  —  reads video features → picks algorithm config → THEN pipeline runs",
    size=9, bold=True, color=RGBColor(0,200,180), line_rgb=(0,160,140))

# ── N1-N4 boxes (y=2.10, h=1.22) ───────────────────────────────────────
NY = 2.10
NH = 1.22
NBW = 2.28
NBG = 0.10
n_total = 4*NBW + 3*NBG  # 9.12+0.3=9.42
nx0 = (W - n_total) / 2  # 0.29

n_data = [
    ("N1 — Two-Stage Predictor",
     "Stage 1: low threshold → 40 candidates, 96.8% recall\nStage 2: BGE AND E5 agree + min 11 sentences",
     (90,25,45), (110,30,55)),
    ("N2 — Entropy-Weighted Fusion",
     "Combines text+CLIP+OCR+prosody\nWeight = exp(−entropy): noisy = low weight",
     (110,80,10), (130,95,15)),
    ("N3 — Hierarchical Nesting",
     "Runs predictor TWICE:\nChapters (5-20min) ⊃ Subtopics (2-5min)",
     (55,25,100), (70,35,120)),
    ("N4 — Offline LLM Titling",
     "Llama 3.1-8B via Ollama (fully offline)\nReads transcript → chapter title",
     (100,55,10), (120,70,15)),
]
for i, (hdr, body, hdr_col, bdr_col) in enumerate(n_data):
    nx = nx0 + i*(NBW+NBG)
    # Background
    rect(s, nx, NY, NBW, NH, (10,16,35), line_rgb=bdr_col)
    # Header bar
    rect_text(s, nx, NY, NBW, 0.32, hdr_col, hdr, size=8, bold=True)
    # Body
    txbox(s, body, nx+0.08, NY+0.35, NBW-0.16, NH-0.40, size=8, color=C_GRY)

# ── Bottom note ─────────────────────────────────────────────────────────
txbox(s, "BGE-divisive = baseline (no N1-N4). Cross-model = BGE×E5 agreement only. Selector = ExtraTrees picks best config per video.",
      M, 3.38, BODY_W, 0.32, size=8, color=RGBColor(130,140,165), align=PP_ALIGN.CENTER)

print(f"Slide 6 done  (W={W}, H={H})")

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 7 (idx 6) — HOW SEGMENTATION WORKS  (3 columns)
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides[6]
clear_slide(s); set_bg(s)
add_title(s, "How Segmentation Works — From Sentences to Timestamps")
add_rule(s)

CW = 3.06   # column width
CG = 0.14   # column gap
CY = BODY_T
CH = BODY_H  # 4.815

col_fills = [(14,24,55), (14,44,34), (34,20,58)]
col_bdr   = [(30,55,110),(25,90,60),(75,40,110)]

for i in range(3):
    cx = M + i*(CW+CG)
    rect(s, cx, CY, CW, CH, col_fills[i], col_bdr[i])

# Column 1 content
cx1 = M
rect_text(s, cx1, CY, CW, 0.36, (30,58,120), "Step 1 — Sentence Embedding", size=9, bold=True)
multilines(s, [
    ("Each sentence → 1024-dim vector", 9, True, C_BLU),
    ("Model: BGE-large-en-v1.5 + E5-large-v2", 8),
    ("", 3),
    ("Same topic → vectors CLOSE together", 9, True, C_GRN),
    ("  e.g. 'eigenvalue' ↔ 'diagonal matrix'", 8),
    ("", 3),
    ("Diff. topic → vectors FAR apart", 9, True, C_RED),
    ("  e.g. 'eigenvalue' ↔ 'natural selection'", 8),
    ("", 5),
    ("Step 2 — Sliding Window Score", 9, True, C_BLU),
    ("At position i:", 8),
    ("  Before = avg(3 sentences before i)", 8),
    ("  After  = avg(3 sentences after i)", 8),
    ("  Score  = 1 − cosine_sim(Before, After)", 8),
    ("", 3),
    ("HIGH score → meaning changed", 9, True, C_HL),
    ("LOW score  → same topic continues", 8),
], cx1+0.1, CY+0.40, CW-0.2, CH-0.46)

# Column 2 content — score curve image + text
cx2 = M + CW + CG
rect_text(s, cx2, CY, CW, 0.36, (20,80,55), "Step 3 — Peak Detection", size=9, bold=True)
# Score curve image in the top portion of body
img(s, f"{FIG}/how_seg_works_v2.png", cx2+0.05, CY+0.40, CW-0.1, 2.35)
multilines(s, [
    ("Peaks in score curve", 9, True, C_GRN),
    ("= boundary candidates", 8),
    ("", 3),
    ("Filter rules:", 9, True, C_SUB),
    ("  • BGE AND E5 must agree", 8),
    ("  • Min 11 sentences apart", 8),
    ("  • (~90 sec minimum gap)", 8),
    ("", 3),
    ("Result: ~40 candidates", 9, True, C_HL),
    ("96.8% recall of true boundaries", 8),
], cx2+0.1, CY+2.82, CW-0.2, CH-2.88)

# Column 3 content
cx3 = M + 2*(CW+CG)
rect_text(s, cx3, CY, CW, 0.36, (75,38,18), "Step 4 — Timestamps + Titles", size=9, bold=True)
multilines(s, [
    ("Boundary sentence index", 9, True, C_HL),
    ("→ Whisper word timestamp", 8),
    ("→ HH:MM:SS format", 8),
    ("", 5),
    ("Example output:", 9, True, C_SUB),
    ("[00:00]  Introduction", 9, False, C_GRN),
    ("[15:30]  Eigenvalues", 9, False, C_GRN),
    ("[28:45]  Diagonalization", 9, False, C_GRN),
    ("[42:00]  Applications", 9, False, C_GRN),
    ("", 5),
    ("N4 (Llama 3.1-8B offline):", 9, True, C_SUB),
    ("Reads transcript text near", 8),
    ("each boundary → generates", 8),
    ("the chapter title shown above.", 8),
    ("", 5),
    ("Paste into YouTube →", 9, True, C_HL),
    ("instant chapter navigation", 8),
], cx3+0.1, CY+0.40, CW-0.2, CH-0.46)

print("Slide 7 done.")

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 8 (idx 7) — NOVEL METHODS N1-N4  (2×2 grid)
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides[7]
clear_slide(s); set_bg(s)
add_title(s, "Our Four Novel Contributions — N1 to N4")
add_rule(s)

BW8 = 4.75   # box width
BH8 = 2.28   # box height
BG8 = 0.10   # gap

# 2×2 grid: x = [M, M+BW8+BG8],  y = [BODY_T, BODY_T+BH8+BG8]
# Total: 2*4.75+0.10=9.60 wide, 2*2.28+0.10=4.66 tall → ends at 0.73+4.66=5.39 < 5.625 ✓
grid = [
    (0, 0, "N1 — Two-Stage Predictor",     (25,50,105),  (18,38,82), [
        ("Problem: one threshold misses real boundaries (too high) or adds noise (too low).", 9),
        ("", 3),
        ("Stage 1 — low threshold → large candidate pool:", 9, True, C_BLU),
        ("  ~40 candidates per hour of video. 96.8% of real boundaries captured.", 8),
        ("  Goal is RECALL — get every real boundary into the pool.", 8),
        ("", 3),
        ("Stage 2 — cross-model filter:", 9, True, C_BLU),
        ("  BGE AND E5 must both agree on the same peak.", 8),
        ("  Minimum 11 sentences (≈90 sec) between any two boundaries.", 8),
        ("", 3),
        ("Why novel: prior unsupervised work uses a single threshold.", 9, True, C_HL),
        ("Oracle Pk=0.0172 confirms Stage 1 pool is excellent.", 8),
    ]),
    (1, 0, "N2 — Entropy-Weighted Fusion",  (20,95,52),   (15,72,38), [
        ("Problem: how to combine text, visual, audio without labeled training data?", 9),
        ("", 3),
        ("Weight = exp(−entropy(signal))", 9, True, C_GRN),
        ("  Confident signal → low entropy → HIGH weight (e.g. CLIP on slide-heavy video)", 8),
        ("  Noisy signal    → high entropy → LOW weight (auto-discarded)", 8),
        ("  System adapts PER VIDEO — no fixed weights, no training data needed.", 8),
        ("", 3),
        ("Signals: BGE text  +  CLIP visual  +  OCR slide text  +  prosody pauses", 8),
        ("", 3),
        ("Why novel: prior fusion work requires a labeled dataset to learn weights.", 9, True, C_HL),
        ("We operate in zero-labeled-data regime — this is the only viable approach.", 8),
    ]),
    (0, 1, "N3 — Hierarchical Nesting",     (62,28,108),  (46,20,82), [
        ("Problem: all prior unsupervised systems output ONE level of segmentation.", 9),
        ("", 3),
        ("Solution — run the predictor TWICE:", 9, True, RGBColor(200,150,255)),
        ("  Run 1 (coarse):  chapter level  — boundaries every 5-20 minutes", 8),
        ("  Run 2 (fine):    subtopic level — boundaries every 2-5 minutes", 8),
        ("", 3),
        ("Constraint: every chapter boundary IS a subtopic boundary (enforced).", 8),
        ("  Hierarchically valid — a major topic change is also a subtopic change.", 8),
        ("", 3),
        ("Why novel: prior work is single-level. Our subtopic annotation (Kappa=0.71)", 9, True, C_HL),
        ("is what made this contribution possible.", 8),
    ]),
    (1, 1, "N4 — Offline LLM Titling",      (108,58,12),  (82,44,8), [
        ("Problem: timestamps alone are not human-readable navigation labels.", 9),
        ("", 3),
        ("Solution — Llama 3.1-8B via Ollama:", 9, True, C_HL),
        ("  Reads 10-15 sentences around each boundary.", 8),
        ("  Generates a chapter title: [15:30] Eigenvalue Decomposition", 8),
        ("  Fully OFFLINE — no internet, no API key, runs on any laptop.", 8),
        ("", 3),
        ("Pk is UNCHANGED — titles do not affect boundary positions.", 8),
        ("  Proves the rest of the pipeline is independently sound.", 8),
        ("", 3),
        ("Why novel: prior systems output raw timestamps only.", 9, True, C_HL),
        ("Offline deployment makes this viable for low-resource universities.", 8),
    ]),
]
for (ci, ri, title8, hdr_col, bdr_col, lines) in grid:
    bx = M + ci*(BW8+BG8)
    by = BODY_T + ri*(BH8+BG8)
    rect(s, bx, by, BW8, BH8, (10,16,35), bdr_col)
    rect_text(s, bx, by, BW8, 0.33, hdr_col, title8, size=9, bold=True)
    multilines(s, lines, bx+0.1, by+0.37, BW8-0.2, BH8-0.43)

print("Slide 8 done.")

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 13 (idx 12 in v6 — Oracle)
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides[12]
clear_slide(s); set_bg(s)
add_title(s, "The Oracle Experiment — What Is the Upper Ceiling?")
add_rule(s)

# Oracle diagram image — top 58% of body area
IMG_H = BODY_H * 0.58   # ~2.79
img(s, f"{FIG}/oracle_diagram_v2.png", M, BODY_T, BODY_W, IMG_H)

# 3 bullets below image — no overlap
BUL_Y = BODY_T + IMG_H + 0.08
BUL_H = BODY_H - IMG_H - 0.12

multilines(s, [
    ("①  96.8% recall — the correct timestamps ARE in the candidate pool. Detection is NOT the bottleneck.", 9, True, C_GRN),
    ("②  Oracle Pk = 0.0172 — near-perfect segmentation IS achievable from our pipeline architecture.", 9, True, C_HL),
    ("③  Gap ΔPk = 0.34 — entirely a selection / ranking problem. Next step: supervised boundary ranker.", 9, True, C_BLU),
], M, BUL_Y, BODY_W, BUL_H)

print("Slide 13 (Oracle) done.")

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 14 (idx 13 after reorder — Granularity Mismatch)
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides[13]
clear_slide(s); set_bg(s)
add_title(s, "Granularity Mismatch — Why Fine-Grained Signals Fail")
add_rule(s)

# Left: 3 level boxes  (width 4.5")
LW = 4.50
LX = M

txbox(s, "Three levels of topic change:", LX, BODY_T, LW, 0.28,
      size=9, bold=True, color=C_HL)

# Each level: header strip + body box
# Total height = BODY_H = 4.815, minus 0.30 for label = 4.515
# 3 levels: each gets 4.515/3 = 1.505
LVL_H  = (BODY_H - 0.32) / 3   # 1.498 each
LVL_TH = 0.33                   # header height within each level

levels = [
    ("FINE-GRAINED  (30-120 sec)",
     "Prosody pauses, discourse markers,\nBERTopic, GPT-2 surprisal\n→ Fires too often for chapter GT\n→ Over-segments → HIGH Pk error",
     (105,25,35), (80,20,28)),
    ("MID-GRAINED  (2-5 min)",
     "CLIP visual scene changes\n→ Closer to chapter scale\n→ Reduces over-segmentation\n→ HELPS Pk",
     (105,85,10), (80,65,8)),
    ("COARSE  (5-20 min) ← OUR TARGET",
     "YouTube creator chapter markers\n= our ground truth\nBGE cosine dissimilarity operates\nat this scale → BEST Pk",
     (18,90,50), (14,68,38)),
]
for i, (hdr, body, hcol, bcol) in enumerate(levels):
    ly = BODY_T + 0.32 + i*LVL_H
    rect_text(s, LX, ly, LW, LVL_TH, hcol, hdr, size=8, bold=True)
    rect(s, LX, ly+LVL_TH, LW, LVL_H-LVL_TH-0.04, (10,16,35), bcol)
    txbox(s, body, LX+0.1, ly+LVL_TH+0.06, LW-0.2, LVL_H-LVL_TH-0.12, size=8, color=C_GRY)

# Right: bar chart image — starts right after left panel
RX = LX + LW + 0.18
RW = W - RX - M
img(s, f"{FIG}/granularity_chart.png", RX, BODY_T, RW, BODY_H)

print("Slide 14 (Granularity) done.")

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 22 (idx 21 — Final Verdict)  3-column layout
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides[21]
clear_slide(s); set_bg(s)
add_title(s, "Final Verdict — What We Achieved and What Comes Next")
add_rule(s)

# 3 equal columns
CW22 = (BODY_W - 2*0.12) / 3   # (9.6 - 0.24)/3 = 3.12
CG22 = 0.12
CY22 = BODY_T
# Reserve bottom 0.32" for closing note
CH22 = BODY_H - 0.38

col_defs = [
    (M,                   "✅  What We Achieved",    (18,75,42),  (12,55,30), [
        ("Pk = 0.3588  /  WD = 0.3739", 10, True, C_HL),
        ("p = 0.025 vs all baselines", 9, False, C_GRN),
        ("", 3),
        ("Matches TreeSeg Pk=0.367", 9, False, C_GRY),
        ("with ZERO training labels", 8),
        ("", 3),
        ("External lecture video:", 9, True, C_GRN),
        ("Pk=0.3974 ≈ benchmark", 8),
        ("(confirms signal generalizes", 8),
        ("within lecture format)", 8),
        ("", 3),
        ("Fully offline, no API", 8),
        ("First public benchmark", 8),
    ]),
    (M+CW22+CG22,          "⚠️  Honest Limitations",  (88,50,14),  (65,37,10), [
        ("Selector leave-domain-out:", 9, True, RGBColor(240,160,50)),
        ("Pk = 0.4012", 10, True, C_HL),
        ("(below BGE-div baseline)", 8),
        ("", 3),
        ("30 videos is modest:", 9, False, C_GRY),
        ("bootstrap CIs are wide", 8),
        ("", 3),
        ("LLM titles work;", 9, False, C_GRY),
        ("positions do NOT improve", 8),
        ("", 3),
        ("Not comparable to large", 9, False, C_GRY),
        ("supervised systems", 8),
        ("(different scope/data)", 8),
    ]),
    (M+2*(CW22+CG22),      "🚀  What Comes Next",     (18,48,108), (12,36,80), [
        ("Supervised boundary ranker", 9, True, C_BLU),
        ("→ close ΔPk=0.34 gap", 8),
        ("(oracle shows this is the", 8),
        ("selection bottleneck)", 8),
        ("", 3),
        ("CLIP-centered fusion:", 9, True, C_BLU),
        ("already at right granularity", 8),
        ("", 3),
        ("200+ lectures → robust", 9, True, C_BLU),
        ("selector training", 8),
        ("", 3),
        ("Cross-lingual: LaBSE +", 9, True, C_BLU),
        ("multilingual Whisper", 8),
    ]),
]

for (cx, hdr22, hcol, bcol, lines) in col_defs:
    rect(s, cx, CY22, CW22, CH22, (10,16,35), bcol)
    rect_text(s, cx, CY22, CW22, 0.33, hcol, hdr22, size=9, bold=True)
    multilines(s, lines, cx+0.1, CY22+0.37, CW22-0.2, CH22-0.42)

txbox(s,
    "LecSeg-30 proves: zero-shot lecture segmentation IS achievable. "
    "We built it, measured it, and diagnosed the bottleneck (ΔPk=0.34 — a selection problem).",
    M, CY22+CH22+0.06, BODY_W, 0.28, size=9, bold=True, color=C_HL, align=PP_ALIGN.CENTER)

print("Slide 22 (Final Verdict) done.")

# ─────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\n✓ Saved: {OUTPUT}  ({len(prs.slides)} slides, {W:.1f}\" x {H:.3f}\")")
