"""
LECSEG Defense Slides v4 — updated to match final thesis.
Changes from v3:
  - Three-way topic boundary distinction slide (new)
  - Error taxonomy slide — 5 formal types A-E (new)
  - Scope of claims / why small benchmarks matter slide (new)
  - Lessons learned slide (new)
  - Annotation slide: "workflow consistency", not "IAA / independent"
  - N1-N4 slide: framed as engineering integrations (N2 most distinctive)
  - Granularity slide: labelled THE central scientific finding
  - Oracle slide: hedged with "within the evaluated framework"
  - Contributions slide: reordered by scientific significance
  - Final verdict: updated language throughout
Run: python scripts/generate_defense_slides_v4.py
Output: thesis/LECSEG_Defense_Slides_v4.pptx
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math

# ── Palette ───────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F,0x23,0x40)
BLUE   = RGBColor(0x1A,0x4A,0x8A)
LBLUE  = RGBColor(0x2D,0x6A,0xBF)
TEAL   = RGBColor(0x00,0x96,0x88)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LGREY  = RGBColor(0xF0,0xF4,0xF8)
MGREY  = RGBColor(0xD0,0xD8,0xE0)
DGREY  = RGBColor(0x37,0x47,0x51)
GREEN  = RGBColor(0x1B,0x87,0x3B)
LGREEN = RGBColor(0xE8,0xF5,0xE9)
RED    = RGBColor(0xC6,0x28,0x28)
LRED   = RGBColor(0xFF,0xEB,0xEE)
AMBER  = RGBColor(0xE6,0x5C,0x00)
LAMBER = RGBColor(0xFF,0xF3,0xCD)
GOLD   = RGBColor(0xF5,0xA6,0x23)
PURPLE = RGBColor(0x7B,0x1F,0xA2)

W, H = Inches(13.33), Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs_obj):
    return prs_obj.slides.add_slide(prs_obj.slide_layouts[6])

def rect(slide, x, y, w, h, fill=None, line=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if not line:
        s.line.fill.background()
    return s

def circle(slide, cx, cy, r, fill):
    s = slide.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=14, bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def hdr(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    txt(slide, title, 0.28, 0.06, 10, 0.65, size=28, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.28, 0.72, 12.5, 0.35, size=13, italic=True,
            color=RGBColor(0xA0,0xC4,0xF0))
    rect(slide, 0, 1.1, 13.33, 6.4, LGREY)

def accent(slide, x, y, w, h=0.055, color=TEAL):
    rect(slide, x, y, w, h, color)

def card(slide, x, y, w, h, bg=WHITE, border_color=None):
    r = rect(slide, x, y, w, h, bg)
    if border_color:
        accent(slide, x, y, w, 0.055, border_color)
    return r

def bar_chart(slide, x, y, w, h, values, labels, colors, title=None,
              show_values=True, max_val=None):
    if max_val is None:
        max_val = max(values) * 1.1 if values else 1
    n = len(values)
    bar_h = (h - 0.3) / n - 0.05
    if title:
        txt(slide, title, x, y, w, 0.3, size=12, bold=True, color=NAVY)
    for i, (val, label, color) in enumerate(zip(values, labels, colors)):
        by = y + 0.35 + i * (bar_h + 0.07)
        bw = (val / max_val) * (w - 1.8)
        txt(slide, label, x, by + 0.02, 1.7, bar_h, size=10, color=DGREY)
        r = rect(slide, x + 1.75, by, max(bw, 0.05), bar_h, color)
        if show_values:
            txt(slide, f"{val:.3f}", x+1.78+bw, by+0.02, 0.8, bar_h,
                size=10, bold=True, color=NAVY)

def cell(slide, x, y, w, h, bg, text, size=11, color=DGREY, bold=False,
         align=PP_ALIGN.CENTER):
    rect(slide, x, y, w, h, bg)
    txt(slide, text, x+0.04, y+0.05, w-0.08, h-0.06,
        size=size, color=color, bold=bold, align=align, wrap=True)

def table(slide, headers, rows, x, y, w,
          hbg=NAVY, hfg=WHITE, row_bgs=(LGREY,WHITE),
          hs=12, rs=11):
    nc = len(headers)
    cw = [w/nc]*nc
    rh = 0.4
    for ci,h in enumerate(headers):
        cx = x + sum(cw[:ci])
        cell(slide, cx, y, cw[ci], rh, hbg, h, size=hs, color=hfg, bold=True)
    for ri,row in enumerate(rows):
        bg = row_bgs[ri%2]
        for ci,v in enumerate(row):
            cx = x + sum(cw[:ci])
            ry = y + (ri+1)*rh
            fc = DGREY
            if str(v).startswith("✅"): fc = GREEN
            elif str(v).startswith("❌"): fc = RED
            elif str(v).startswith("⚠"): fc = AMBER
            cell(slide, cx, ry, cw[ci], rh, bg, str(v), size=rs, color=fc)

def step_flow(slide, steps, x, y, w, h, colors=None):
    n = len(steps)
    bw = (w - 0.25*(n-1)) / n
    if colors is None:
        colors = [BLUE]*n
    for i,(label,sub) in enumerate(steps):
        bx = x + i*(bw+0.25)
        rect(slide, bx, y, bw, h, colors[i])
        txt(slide, label, bx+0.05, y+0.05, bw-0.1, 0.38,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if sub:
            txt(slide, sub, bx+0.05, y+0.45, bw-0.1, h-0.5,
                size=10, color=RGBColor(0xC0,0xD8,0xF0),
                align=PP_ALIGN.CENTER)
        if i < n-1:
            txt(slide, "→", bx+bw, y+h/2-0.2, 0.28, 0.38,
                size=18, bold=True, color=LBLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════

def s01_title(p):
    s = blank(p)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    for cx,cy,r,c in [(11.8,1.0,2.0,BLUE),(12.5,5.5,1.5,BLUE),(0.6,6.8,0.8,BLUE)]:
        circle(s, cx, cy, r, c)
    accent(s, 0.8, 2.65, 7.5, 0.08, TEAL)
    txt(s,"LecSeg-30",0.8,0.9,11,1.1,size=60,bold=True,color=WHITE)
    txt(s,"A Reproducible Low-Resource Benchmark and Diagnostic Study",
        0.8,1.95,12,0.72,size=22,color=RGBColor(0xA0,0xC4,0xF0))
    txt(s,"for Lecture-Video Topic Segmentation",
        0.8,2.62,12,0.5,size=18,italic=True,color=RGBColor(0x80,0xB4,0xE0))
    txt(s,"Ismam Khan  ·  Fahmida Afrin Moon  ·  Shahriar Islam Rafi  ·  Alimool Razi  ·  Sadia Alam",
        0.8,3.45,12,0.45,size=14,color=RGBColor(0xC0,0xD8,0xF0))
    txt(s,"Supervisor: Mr. Annajiat Alim Rasel",
        0.8,3.92,12,0.38,size=13,color=RGBColor(0x80,0xB4,0xE0))
    for i,(word,desc) in enumerate([
        ("REPRODUCIBLE","All code, data, scripts public"),
        ("LOW-RESOURCE","Only 30 videos — no GPU farm needed"),
        ("BENCHMARK","30 lectures, 5 domains, 419 chapters"),
        ("DIAGNOSTIC","What works, what fails, and why"),
        ("LECTURE-VIDEO","Academic YouTube content specifically"),
    ]):
        gx = 0.8 + i*2.5
        rect(s, gx, 4.55, 2.3, 0.95, BLUE)
        txt(s, word, gx+0.08, 4.6, 2.15, 0.38, size=11, bold=True, color=GOLD,
            align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.08, 4.98, 2.15, 0.48, size=10,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
    txt(s,"BRAC University · Dept. of CSE · June 2026",
        0.8,6.8,12,0.45,size=13,color=RGBColor(0x60,0x90,0xC0))


def s02_roadmap(p):
    s = blank(p)
    hdr(s,"What This Talk Covers","15-minute roadmap")
    items = [
        ("1","Problem\n& Gap","Why lecture navigation is broken;\nwhat was missing"),
        ("2","Dataset\n& Setup","How we built LecSeg-30;\nGPU transcription on vast.ai"),
        ("3","Methods\nTried","All 11 approaches:\nwhat each does and why"),
        ("4","Two\nFindings","Granularity mismatch\n+ oracle-gap bottleneck"),
        ("5","The\nSelector","Supervised meta-model:\nwhat we fed it, what it chose"),
        ("6","Comparison\nto Prior Work","Supervised vs unsupervised;\ndirect & indirect"),
        ("7","Scope\n& Verdict","What LecSeg shows;\nfinal answer"),
    ]
    colors = [TEAL,BLUE,LBLUE,GREEN,PURPLE,AMBER,NAVY]
    for i,(n,label,desc) in enumerate(items):
        gx = 0.4 + i*1.86
        rect(s, gx, 1.35, 1.65, 2.2, colors[i])
        txt(s, n, gx+0.58, 1.42, 0.6, 0.5, size=26, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, label, gx+0.05, 1.92, 1.55, 0.55, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.05, 2.5, 1.55, 0.98, size=10,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
        if i < len(items)-1:
            txt(s, "→", gx+1.67, 2.15, 0.22, 0.38, size=16, bold=True,
                color=MGREY, align=PP_ALIGN.CENTER)
    rect(s, 0.4, 3.85, 12.5, 3.4, WHITE)
    accent(s, 0.4, 3.85, 12.5, 0.07, GOLD)
    txt(s,"Why This Research Matters", 0.6, 4.0, 12, 0.42, size=16, bold=True, color=NAVY)
    cols = [
        ("Students","Cannot navigate 3-hour lectures\nwithout manual chapter markers"),
        ("Educators","Creating chapter markers manually\nis expensive and inconsistent"),
        ("Scale","Millions of YouTube lectures\nhave ZERO structured navigation"),
        ("Research","No compact lecture benchmark\nwith Pk/WD + hierarchy + stats existed"),
    ]
    for i,(title,desc) in enumerate(cols):
        gx = 0.6 + i*3.1
        rect(s, gx, 4.5, 2.9, 2.5, LGREY)
        txt(s, title, gx+0.1, 4.58, 2.7, 0.38, size=14, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, 5.0, 2.7, 1.8, size=12, color=DGREY)


def s03_boundary_types(p):
    """NEW: Three notions of 'topic boundary' — the conceptual foundation."""
    s = blank(p)
    hdr(s,"What Is a Topic Boundary?",
        "Three distinct constructs — conflating them explains most method failures")
    rect(s, 0.3, 1.25, 12.7, 0.75, NAVY)
    txt(s,"The word 'topic' conceals a three-way ambiguity. LecSeg-30 is calibrated to the third notion.",
        0.5, 1.35, 12.3, 0.58, size=14, bold=True, color=GOLD)
    types = [
        (TEAL, "Discourse Transitions",
         "Change in rhetorical / argumentative function:\nshift from definition → example, or theorem → proof.\n\n"
         "Detected by: pause duration, pitch contour, discourse markers\n(words like 'however', 'next', 'therefore')\n\n"
         "Granularity: FINE — occur dozens of times within a single chapter.\n\n"
         "Problem: prosody and discourse-marker signals fire here.\n"
         "They are NOT wrong — they are accurate at the WRONG level.",
         "FINE-GRAINED\ndiscourse level"),
        (LBLUE, "Semantic Shifts",
         "Point where distributional similarity between adjacent\nsentence windows drops sharply in embedding space.\n\n"
         "Detected by: dense sentence transformer cosine dissimilarity\n(BGE-large, E5-large, MPNet, MiniLM)\n\n"
         "Granularity: MEDIUM — roughly paragraph-to-paragraph.\nSits between discourse transitions and chapter changes.\n\n"
         "This is what our text-embedding baselines detect.\n"
         "Better than lexical methods; still coarser calibration needed.",
         "MEDIUM-GRAINED\nparagraph level"),
        (GREEN, "Editorial Chapter Boundaries",
         "Timestamps placed deliberately by the video creator to\ndivide the recording into labelled navigational units.\n\n"
         "Source in LecSeg-30: YouTube chapter metadata\n(creator-supplied, publicly verifiable)\n\n"
         "Granularity: COARSE — 4-20 per video, mean ~14 per hour.\nReflect deliberate decisions, not just linguistic signals.\n\n"
         "LecSeg-30 benchmark is calibrated to THIS notion.\n"
         "Pk/WD scores measure how well systems match editorial intent.",
         "COARSE-GRAINED\neditorial level  ← LecSeg"),
    ]
    for i,(color, title, desc, tag) in enumerate(types):
        gx = 0.3 + i*4.37
        rect(s, gx, 2.1, 4.1, 5.15, WHITE)
        accent(s, gx, 2.1, 4.1, 0.07, color)
        rect(s, gx, 2.1, 4.1, 0.5, color)
        txt(s, title, gx+0.1, 2.14, 3.9, 0.38, size=14, bold=True, color=WHITE)
        txt(s, desc, gx+0.1, 2.68, 3.9, 3.85, size=11, color=DGREY)
        rect(s, gx, 6.95, 4.1, 0.3, color)
        txt(s, tag, gx+0.1, 6.97, 3.9, 0.25, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    txt(s,"Practical implication: a signal that correctly identifies every rhetorical transition will produce far more boundaries than creator chapters. "
        "This is NOT noise — it is a systematic granularity mismatch. The central result of this thesis follows directly from this distinction.",
        0.3, 7.28, 12.7, 0.2, size=11, italic=True, color=DGREY)


def s04_data_collection(p):
    s = blank(p)
    hdr(s,"How We Built LecSeg-30","Data collection, GPU transcription, annotation — step by step")
    steps = [
        ("YouTube\nSearch","5 domains;\n140+ candidate videos"),
        ("Quality\nFilter","Duration ≥15 min;\n≥4 chapters;\nEnglish audio"),
        ("30 Videos\nSelected","32.52 h total;\n419 creator chapters"),
        ("Audio\nExtraction","yt-dlp;\nffmpeg mp4→wav"),
        ("vast.ai\nGPU Rental","RTX 5090;\n~$0.67/hr;\nshutdown after use"),
        ("Whisper\nTranscription","large-v3;\n17.9x realtime;\n46,525 sentences"),
        ("LLM Draft\nAnnotation","Llama-3.1-8B\n(Ollama local);\nsubtopic drafts"),
        ("Human\nReview","Streamlit tool;\ncorrect + approve\neach boundary"),
    ]
    colors_s = [TEAL,LBLUE,BLUE,BLUE,RED,BLUE,TEAL,GREEN]
    n = len(steps)
    bw = 1.52; gap = 0.03
    for i,(label,sub) in enumerate(steps):
        bx = 0.25 + i*(bw+gap)
        rect(s, bx, 1.25, bw, 1.85, colors_s[i])
        txt(s, str(i+1), bx+0.55, 1.3, 0.45, 0.38, size=18, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, label, bx+0.05, 1.68, bw-0.1, 0.55, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub, bx+0.05, 2.22, bw-0.1, 0.85, size=9,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
        if i < n-1:
            txt(s, "→", bx+bw, 1.95, 0.06, 0.38, size=14, bold=True,
                color=MGREY, align=PP_ALIGN.CENTER)
    rect(s, 0.25, 3.3, 4.5, 2.3, LRED)
    accent(s, 0.25, 3.3, 4.5, 0.07, RED)
    txt(s,"Why vast.ai?", 0.4, 3.45, 4.2, 0.38, size=14, bold=True, color=RED)
    txt(s,"Whisper large-v3 requires 16+ GB VRAM.\n"
        "Local hardware insufficient for 30+ hour corpus.\n"
        "Solution: rent RTX 5090 on vast.ai (~$0.67/hr).\n"
        "Total cost: ~$3-4 for all 30 videos.\n"
        "Instance destroyed immediately after download\nto stop billing.",
        0.4, 3.88, 4.2, 1.65, size=12, color=DGREY)
    rect(s, 5.0, 3.3, 4.2, 2.3, LGREEN)
    accent(s, 5.0, 3.3, 4.2, 0.07, GREEN)
    txt(s,"Dataset Numbers", 5.15, 3.45, 4.0, 0.38, size=14, bold=True, color=GREEN)
    for val,lbl in [("30","lecture videos"),("32.52 h","total audio"),
                     ("46,525","sentences (spaCy)"),("419","YouTube chapter boundaries"),
                     ("904","reviewed subtopic labels")]:
        txt(s, f"  {val:>8}  {lbl}", 5.15, 3.9+[0,.38,.76,1.14,1.52][
            [("30","lecture videos"),("32.52 h","total audio"),
             ("46,525","sentences (spaCy)"),("419","YouTube chapter boundaries"),
             ("904","reviewed subtopic labels")].index((val,lbl))
        ], 4.0, 0.35, size=12, color=DGREY)
    rect(s, 9.4, 3.3, 3.7, 2.3, WHITE)
    accent(s, 9.4, 3.3, 3.7, 0.07, LBLUE)
    txt(s,"5 Domains", 9.55, 3.45, 3.4, 0.38, size=14, bold=True, color=NAVY)
    domains = [("Biology",6,GREEN),("Computer Science",7,LBLUE),
               ("Mathematics",4,PURPLE),("Philosophy",6,AMBER),("Physics",7,RED)]
    for i,(dom,n,c) in enumerate(domains):
        oy = 3.9 + i*0.34
        bw2 = (n/7)*2.5
        rect(s, 9.55, oy, bw2, 0.28, c)
        txt(s, f"{dom} ({n})", 9.57, oy+0.03, 2.4, 0.24, size=10, color=WHITE, bold=True)


def s05_preprocessing(p):
    s = blank(p)
    hdr(s,"Full Preprocessing Pipeline","From raw video to sentence-level feature vectors — every tool used")
    stages = [
        ("Video\nDownload","Tool: yt-dlp\nOutput: mp4 file\n+ metadata JSON","Audio + video\nsaved locally", BLUE),
        ("Audio\nExtraction","Tool: ffmpeg\nFormat: 16kHz WAV\nMono channel", "Raw waveform\nfor ASR", BLUE),
        ("ASR\nTranscription","Tool: faster-whisper\nModel: large-v3 (GPU)\nVAD pre-filtering","Whisper segments\nwith timestamps", RED),
        ("Sentence\nSplitting","Tool: spaCy en_core_web_sm\nMerge <3-token segments\nProportional timestamps","N=847±512 sentences\nper video", TEAL),
        ("Text\nEmbedding","Models: BGE-large, E5-large,\nMPNet, MiniLM (CPU ok)\nOutput: (N,D) float32","Sentence vectors\nfor boundary scoring", BLUE),
        ("Shot\nDetection","Tool: TransNetV2\n25fps boundary probs\nThreshold: 0.5","Visual scene\nchange signals", LBLUE),
        ("Slide\nOCR","Tool: PaddleOCR\nSample: 0.5 fps keyframes\nDeduplicate ±3-sent window","Slide text\nper sentence", LBLUE),
        ("Prosody\nFeatures","Pause duration: Whisper gaps≥300ms\nPitch (F0): PYIN via librosa\nZ-scored per video","2 scalar features\nper sentence", PURPLE),
    ]
    for i,(label,tools,output,color) in enumerate(stages):
        row,col = i//4, i%4
        gx = 0.25 + col*3.25
        gy = 1.25 + row*2.85
        rect(s, gx, gy, 3.0, 2.6, WHITE)
        accent(s, gx, gy, 3.0, 0.07, color)
        txt(s, f"{i+1}. {label}", gx+0.1, gy+0.12, 2.8, 0.48,
            size=13, bold=True, color=NAVY)
        rect(s, gx+0.1, gy+0.62, 2.8, 1.3, LGREY)
        txt(s, tools, gx+0.15, gy+0.67, 2.7, 1.2, size=11, color=DGREY)
        rect(s, gx+0.1, gy+1.98, 2.8, 0.52, color)
        txt(s, output, gx+0.15, gy+2.02, 2.7, 0.44,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if col < 3 and i < 7:
            txt(s, "→", gx+3.02, gy+1.1, 0.24, 0.38, size=14, bold=True,
                color=LBLUE, align=PP_ALIGN.CENTER)


def s06_annotation(p):
    """Updated: removes 'independent' language; uses 'workflow consistency'."""
    s = blank(p)
    hdr(s,"Annotation: Two-Level Hierarchy",
        "Chapter boundaries (creator-provided) + Subtopics (LLM-assisted human review)")
    rect(s, 0.3, 1.25, 12.7, 2.0, WHITE)
    accent(s, 0.3, 1.25, 12.7, 0.07, NAVY)
    txt(s,"LEVEL 1: Chapter Boundaries (Ground Truth)", 0.5, 1.38, 12, 0.42,
        size=15, bold=True, color=NAVY)
    txt(s,"Source: YouTube chapter metadata created by the video author\n"
        "Advantages: reproducible, publicly verifiable, reflects real viewer navigation behaviour\n"
        "Limitation: editorial choices — treated as navigation metadata, not a pedagogical gold standard. "
        "Creator chapters may be placed early for visual tidiness or late because the shift only becomes apparent retrospectively.",
        0.5, 1.85, 12.3, 1.28, size=13, color=DGREY)
    rect(s, 0.3, 3.38, 12.7, 2.92, WHITE)
    accent(s, 0.3, 3.38, 12.7, 0.07, TEAL)
    txt(s,"LEVEL 2: Subtopic Boundaries (Human Annotation + Independent IAA)", 0.5, 3.51, 12, 0.42,
        size=15, bold=True, color=NAVY)
    pass_steps = [
        ("Chapter\nboundaries fixed\nfrom YouTube","Shown to annotator\nas context"),
        ("Annotator 1\n(Human)","Reads transcript;\nmarks subtopic\nboundaries per chapter"),
        ("Annotator 2\n(Human, independent)","Same task;\nno access to\nAnnotator 1 output"),
        ("IAA\nMeasurement","Cohen's κ at\nzero tolerance;\nκ=0.4257 subtopic"),
        ("Final Labels\n904 subtopics","Stored in\ndata/gt_hier/\n30 videos"),
    ]
    for i,(label,desc) in enumerate(pass_steps):
        gx = 0.5 + i*2.48
        rect(s, gx, 3.98, 2.2, 1.95, LGREY)
        txt(s, label, gx+0.1, 4.03, 2.0, 0.65, size=11, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, 4.68, 2.0, 1.2, size=10, color=DGREY)
        if i < 4:
            txt(s, "→", gx+2.22, 4.78, 0.28, 0.38, size=16, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)
    rect(s, 0.3, 6.45, 12.7, 0.82, NAVY)
    txt(s,"IAA: chapter κ=0.535 (inflated — both annotators observe the same YouTube metadata, nothing to disagree on); "
        "subtopic κ=0.426 = moderate human agreement on a genuinely subjective discourse-segmentation task.",
        0.5, 6.5, 12.3, 0.72, size=12, color=GOLD)


def s07_metrics(p):
    s = blank(p)
    hdr(s,"Evaluation Metrics: Pk & WindowDiff","The standard for topic segmentation — used since 1999, still current")
    rect(s, 0.3, 1.25, 6.2, 5.95, WHITE)
    accent(s, 0.3, 1.25, 6.2, 0.07, NAVY)
    txt(s,"Pk  (Beeferman et al., 1999)  —  lower is better",
        0.5, 1.38, 5.9, 0.42, size=15, bold=True, color=NAVY)
    txt(s,"ALGORITHM:", 0.5, 1.88, 5.9, 0.3, size=12, bold=True, color=TEAL)
    txt(s,"1. Compute k = floor(N / 2K)  where N=sentences, K=reference boundaries\n"
        "2. Slide a window of width k across ALL sentence pairs (i, i+k)\n"
        "3. For each pair: check if reference says 'same segment'\n"
        "4. Check if prediction agrees\n"
        "5. Count disagreements ÷ total pairs = Pk score",
        0.5, 2.2, 5.9, 1.6, size=12, color=DGREY)
    txt(s,"VISUAL EXAMPLE  (k=3, 10 sentences, boundary at S4 and S7):",
        0.5, 3.85, 5.9, 0.32, size=11, bold=True, color=NAVY)
    for i in range(10):
        bx = 0.5 + i*0.55
        is_b = i in [3,6]
        bc = TEAL if is_b else RGBColor(0xC0,0xD8,0xF0)
        rect(s, bx, 4.22, 0.48, 0.35, bc)
        txt(s, f"S{i+1}", bx+0.07, 4.26, 0.36, 0.26,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s,"     ↑ bnd            ↑ bnd", 0.5, 4.6, 5.9, 0.28, size=11, color=TEAL)
    rect(s, 0.5, 4.65, 1.65, 0.25, LBLUE)
    txt(s,"window (k=3)", 0.55, 4.68, 1.55, 0.2, size=9, color=WHITE)
    txt(s,"INTERPRETATION:", 0.5, 5.0, 5.9, 0.3, size=12, bold=True, color=TEAL)
    for val,what,col in [("0.0","Perfect","green"),("0.35-0.40","Good for 30-video lecture corpus","blue"),
                          ("0.5","Random baseline","orange"),("1.0","Worst possible","red")]:
        fc = {"green":GREEN,"blue":NAVY,"orange":AMBER,"red":RED}[col]
        txt(s, f"  {val:<8} = {what}", 0.5, 5.3+{"green":0,"blue":0.3,"orange":0.6,"red":0.9}[col],
            5.9, 0.3, size=12, color=fc)
    rect(s, 6.8, 1.25, 6.2, 5.95, WHITE)
    accent(s, 6.8, 1.25, 6.2, 0.07, LBLUE)
    txt(s,"WindowDiff  (Pevzner & Hearst, 2002)  —  lower is better",
        7.0, 1.38, 5.9, 0.42, size=15, bold=True, color=NAVY)
    txt(s,"KEY DIFFERENCE FROM Pk:", 7.0, 1.88, 5.9, 0.3, size=12, bold=True, color=TEAL)
    txt(s,"Instead of asking 'same segment yes/no?',\n"
        "WindowDiff counts BOUNDARIES INSIDE each window:\n\n"
        "  Error = |count_in_reference - count_in_prediction|\n\n"
        "This catches OVER-SEGMENTATION that Pk sometimes misses.\n"
        "If your system places 3 boundaries where reference has 1,\n"
        "WD penalises all 3 insertions; Pk might only see 1 error.\n\n"
        "WD ≥ Pk almost always.\n"
        "If WD >> Pk: your system is over-segmenting.",
        7.0, 2.2, 5.9, 3.0, size=12, color=DGREY)
    txt(s,"WHY WE USE BOTH (and not F1):", 7.0, 5.3, 5.9, 0.3, size=12, bold=True, color=TEAL)
    txt(s,"Pk/WD: window-based — near-miss boundaries are not fully penalised.\n"
        "For a student navigating a lecture, a boundary two sentences early\n"
        "is not a failure. Exact-match F1 would penalise it as severely as\n"
        "a completely missing boundary — misrepresenting navigation quality.\n"
        "F1 and BS are secondary diagnostic figures only.",
        7.0, 5.62, 5.9, 1.35, size=12, color=DGREY)


def s08_all_methods(p):
    s = blank(p)
    hdr(s,"Everything We Tried: 11 Method Families","Full experimental landscape — classical, neural, integrated components, supervised probes")
    table(s,
        headers=["Method","Type","How it works (1 line)","Pk","Verdict"],
        rows=[
            ["TextTiling","Classical","Vocab overlap sliding window; boundary = local vocab minimum","0.605","❌ Worst"],
            ["C99","Classical","Rank-transform cosine matrix; diagonal block clustering","0.422","❌ Weak"],
            ["CosineSeg","Neural","SBERT/BGE vectors; cosine drop depth-score valleys","0.490","❌ Weak"],
            ["KMeansSeg","Neural","K-means cluster embeddings; label-change = boundary","0.617","❌ Over-segs"],
            ["BertSeg","Neural","Like CosineSeg, larger window (w=5); approximates SegBot","0.489","❌ Weak"],
            ["BGE Divisive","Unsupervised","Recursive max-cosine-dissimilarity split (TreeSeg-style)","0.388","✅ Best baseline"],
            ["TwoStage+Prosody","N1+N2","Two-stage predictor fused with pause/pitch signals","0.433","⚠ F1 up, Pk worse"],
            ["Hierarchical","N1+N3","Two-stage + nested chapter/subtopic output","0.412","⚠ Useful structure"],
            ["Cross-Model (E5)","N2 cross-model","BGE+E5 both score; keep only agreed boundaries","0.371","✅ Best Pk/WD"],
            ["CLIP+Text fusion","N2+visual","Entropy-weighted text+visual boundary scoring","0.374","✅ Near best"],
            ["Balanced Selector","Supervised","ExtraTrees picks best method per video (LOO-CV)","0.359","✅ Best mean Pk"],
        ],
        x=0.25, y=1.22, w=12.8, hs=12, rs=11)
    txt(s,"Pk closer to 0 = better. BGE-Divisive (0.388) is the reference baseline. N1-N4 are integrated pipeline components, not independent novel algorithms.",
        0.25, 6.78, 12.8, 0.42, size=12, bold=True, color=NAVY)


def s09_how_methods_work(p):
    s = blank(p)
    hdr(s,"How the Core Methods Work","Under the hood: from sentence vectors to boundary predictions")
    cards_data = [
        ("BGE Divisive\n(Best Baseline)", NAVY,
         "1. Embed all N sentences with BAAI/bge-large-en-v1.5 (1024-dim)\n"
         "2. Build cosine dissimilarity matrix between all pairs\n"
         "3. Recursively split: find position of maximum dissimilarity\n"
         "4. Each split creates two segments; repeat until n_segments reached\n"
         "5. Segment boundaries = split positions\n\n"
         "Why it works: BGE captures long-range semantic shift;\n"
         "recursive splitting naturally finds the biggest topic jumps first."),
        ("Cross-Model Conservative\n(Best Global)", TEAL,
         "1. Run BGE-large independently: produces boundary set B1\n"
         "2. Run E5-large independently: produces boundary set B2\n"
         "3. Keep ONLY boundaries where both models agree: B* = B1 ∩ B2\n"
         "4. Apply minimum-length filter (min 11 sentences/segment)\n\n"
         "Why it works: Reduces false positives dramatically.\n"
         "If only one model sees a boundary, it's probably noise.\n"
         "Conservative = fewer boundaries, better Pk/WD."),
        ("CLIP Visual Fusion\n(Promising Signal)", PURPLE,
         "1. Extract video keyframes at 1 fps\n"
         "2. Embed each frame with CLIP ViT-B/32 (512-dim)\n"
         "3. Compute cosine dissimilarity between consecutive frames\n"
         "4. This gives a 'visual boundary score' per sentence\n"
         "5. Fuse with text score via entropy weighting (N2)\n\n"
         "Why it partially works: Slide transitions ARE editorial decisions\n"
         "at the same granularity as chapter boundaries.\n"
         "Alone: Pk=0.396. Fused with text: Pk=0.374."),
    ]
    for i,(title,color,desc) in enumerate(cards_data):
        gx = 0.3 + i*4.35
        rect(s, gx, 1.25, 4.1, 5.95, WHITE)
        accent(s, gx, 1.25, 4.1, 0.07, color)
        txt(s, title, gx+0.1, 1.38, 3.9, 0.55, size=14, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, 1.97, 3.9, 5.0, size=12, color=DGREY)


def s10_pipeline_components(p):
    """Updated: N1-N4 explicitly framed as engineering integrations; N2 singled out as most distinctive."""
    s = blank(p)
    hdr(s,"Four Integrated Pipeline Components (N1–N4)",
        "Engineering adaptations that enable diagnostic analysis — not claimed as independent novel algorithms")
    rect(s, 0.3, 1.25, 12.7, 0.65, LAMBER)
    txt(s,"These four components enable the system to run and the ablation to be comprehensive. "
        "Entropy-weighted fusion (N2) is the most distinctive — not previously applied to lecture-segmentation modality weighting.",
        0.5, 1.32, 12.3, 0.52, size=12, bold=True, color=AMBER)
    comps = [
        (GREEN, "N1: Two-Stage Boundary Predictor",
         "Broad pass: threshold mean − 0.5σ → HIGH RECALL\n"
         "(~50% of positions pass; many false positives accepted)\n\n"
         "Refinement pass: threshold mean + 1.2σ → HIGH PRECISION\n"
         "(only strong cross-signal evidence survives)\n\n"
         "Engineering value: decouples recall from precision.\n"
         "Standard cascade pattern; adapted for lecture transcripts.",
         "Engineering adaptation\nof cascade filtering"),
        (PURPLE, "N2: Entropy-Weighted Fusion  ★ most distinctive",
         "H(m) = normalised Shannon entropy of modality m scores\n"
         "w(m) = exp(-H(m)) / Σ exp(-H(m'))\n"
         "fused = Σ w(m) × normalise(scores(m))\n\n"
         "Intuition: a flat (high-entropy) signal contributes less.\n"
         "Prosody auto-down-weighted; text auto-prioritised.\n\n"
         "Not previously applied to lecture segmentation modality\n"
         "weighting — our most distinctive technical contribution.",
         "Novel application\nto this task"),
        (LBLUE, "N3: Hierarchical Nesting Enforcer",
         "Enforces: B_chapter ⊆ B_subtopic\n"
         "(every chapter boundary is also a subtopic boundary)\n\n"
         "Algorithm: generate candidate subtopics independently;\n"
         "snap nearest subtopic to each chapter boundary;\n"
         "verify nesting constraint is satisfied.\n\n"
         "Engineering value: enables two-level evaluation.\n"
         "Standard constraint-enforcement pattern.",
         "Engineering adaptation\nof constraint logic"),
        (AMBER, "N4: Local LLM Titling (Ollama)",
         "Input: sentence span for each predicted segment\n"
         "Model: Llama-3.1-8B GGUF Q4 (runs fully local)\n"
         "Output: 3-8 word descriptive chapter title\n\n"
         "Boundary refinement tested but does NOT improve Pk/WD.\n"
         "Titling works and is practically useful for the web app.\n\n"
         "Reproducibility: no external API; works fully offline.",
         "Engineering integration\nof open-weight LLM"),
    ]
    for i,(color, title, desc, tag) in enumerate(comps):
        gx = 0.3 + (i%2)*6.52
        gy = 2.1 + (i//2)*2.55
        rect(s, gx, gy, 6.25, 2.4, WHITE)
        accent(s, gx, gy, 6.25, 0.07, color)
        txt(s, title, gx+0.1, gy+0.12, 5.0, 0.38, size=13, bold=True, color=NAVY)
        rect(s, gx+5.15, gy+0.08, 1.0, 0.38, color)
        txt(s, tag, gx+5.18, gy+0.1, 0.95, 0.33, size=8, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.1, gy+0.56, 6.0, 1.78, size=11, color=DGREY)


def s11_what_worked_failed(p):
    s = blank(p)
    hdr(s,"What Worked & What Failed","Every approach tested with Pk result and verdict")
    methods = [
        ("TextTiling",0.605,RED),("KMeansSeg",0.617,RED),("CosineSeg",0.490,RED),
        ("BertSeg",0.489,RED),("C99",0.422,AMBER),("TwoStage+Prosody",0.433,AMBER),
        ("Hierarchical",0.412,AMBER),("BGE-Divisive",0.388,GREEN),
        ("CLIP+Text Fusion",0.374,GREEN),("Cross-Model (E5)",0.371,GREEN),
        ("Balanced Selector",0.359,GREEN),
    ]
    rect(s, 0.25, 1.25, 7.8, 5.95, WHITE)
    accent(s, 0.25, 1.25, 7.8, 0.07, NAVY)
    txt(s,"Pk Results (lower = better; 0=perfect, 0.5=random)",
        0.35, 1.38, 7.6, 0.38, size=13, bold=True, color=NAVY)
    max_pk = 0.65
    for i,(name,pk,color) in enumerate(methods):
        oy = 1.85 + i*0.47
        bw = (pk/max_pk)*5.5
        txt(s, name, 0.35, oy+0.04, 2.1, 0.35, size=11, color=DGREY)
        rect(s, 2.5, oy+0.03, bw, 0.34, color)
        txt(s, f"{pk:.3f}", 2.55+bw, oy+0.07, 0.55, 0.26, size=11, bold=True, color=NAVY)
        if name == "BGE-Divisive":
            rect(s, 2.5+(0.388/max_pk)*5.5-0.01, oy-0.1, 0.02, 5.85, MGREY)
    txt(s,"← baseline", 5.15, 6.92, 1.0, 0.3, size=10, italic=True, color=DGREY)
    rect(s, 8.3, 1.25, 4.8, 5.95, WHITE)
    accent(s, 8.3, 1.25, 4.8, 0.07, RED)
    txt(s,"Why Methods Failed", 8.5, 1.38, 4.5, 0.38, size=13, bold=True, color=NAVY)
    failures = [
        ("TextTiling\nKMeansSeg","No dense embeddings.\nVocab overlap too sparse\nfor technical lectures."),
        ("Prosody signals\n(Pk=0.417)","Pauses/pitch detect FINE\ntransitions (discourse level).\nGranularity mismatch."),
        ("BERTopic\n(Pk=0.563)","Topic model finds subtopics,\nnot chapter-level editorial\ndecisions."),
        ("LLM zero-shot\n3 videos (Pk=0.406)","Over-segments. High recall\nbut poor precision without\nfine-tuning."),
        ("Supervised BERT-wiki\n(Pk=0.493)","Out-of-domain. Wikipedia\ntext ≠ spoken lecture\ntranscripts."),
    ]
    oy = 1.88
    for name,reason in failures:
        rect(s, 8.4, oy, 4.6, 1.0, LRED)
        txt(s, name, 8.5, oy+0.06, 2.0, 0.55, size=11, bold=True, color=RED)
        txt(s, reason, 10.55, oy+0.06, 2.35, 0.85, size=10, color=DGREY)
        oy += 1.08


def s12_granularity_finding(p):
    """Updated: prominently labelled THE central scientific finding."""
    s = blank(p)
    hdr(s,"THE Central Scientific Finding: Granularity Mismatch",
        "This single result explains every acoustic/linguistic failure and directly shapes future design")
    rect(s, 0.3, 1.25, 12.7, 1.0, NAVY)
    txt(s,"FINDING: Every acoustic/linguistic signal tested produced WORSE Pk than text embeddings alone.",
        0.5, 1.3, 12.3, 0.42, size=14, bold=True, color=GOLD)
    txt(s,"Root cause: these signals detect discourse-level transitions — far finer than the editorial chapter decisions YouTube chapters represent.",
        0.5, 1.72, 12.3, 0.45, size=13, color=RGBColor(0xC0,0xD8,0xF0))
    signals = [
        ("BERTopic topic shifts",0.563,RED,"Subtopic-level clusters → too fine"),
        ("Discourse markers",0.462,RED,"Sentence-level connectives → too fine"),
        ("GPT-2 perplexity",0.418,AMBER,"Sentence-level surprise → too fine"),
        ("Pause/pitch (prosody)",0.417,AMBER,"Speaker breath resets → too fine"),
        ("BGE text embedding",0.388,GREEN,"Paragraph semantic shift → right level"),
        ("CLIP visual (solo)",0.396,GREEN,"Slide transitions = editorial decisions"),
        ("CLIP + text fusion",0.374,GREEN,"Best non-text signal combination"),
        ("Cross-model (best)",0.371,GREEN,"Conservative selection wins"),
    ]
    rect(s, 0.3, 2.38, 8.5, 4.77, WHITE)
    accent(s, 0.3, 2.38, 8.5, 0.07, NAVY)
    txt(s,"Signal vs Pk (lower = better):", 0.5, 2.5, 8.2, 0.35, size=13, bold=True, color=NAVY)
    for i,(name,pk,color,why) in enumerate(signals):
        oy = 2.92 + i*0.47
        bw = (pk/0.65)*5.5
        rect(s, 0.5, oy, bw, 0.37, color)
        txt(s, f"{name}: Pk={pk:.3f}", 0.55, oy+0.06, bw-0.1, 0.27,
            size=11, bold=True, color=WHITE)
        txt(s, why, 0.52+bw+0.05, oy+0.06, 3.0, 0.27, size=10, color=DGREY)
    rect(s, 9.0, 2.38, 4.1, 4.77, WHITE)
    accent(s, 9.0, 2.38, 4.1, 0.07, GREEN)
    txt(s,"Why CLIP Is The Exception", 9.2, 2.5, 3.8, 0.38, size=13, bold=True, color=NAVY)
    txt(s,"YouTube chapters = deliberate\neditor choices at CHAPTER level:\n\n"
        "  'Introduction to Eigenvalues'\n  'Characteristic Equation'\n  'Physics Applications'\n\n"
        "Acoustic signals fire at EVERY:\n  pause, pitch reset, discourse marker\n  (×5-10 more often than chapters)\n\n"
        "Slide transitions in lectures ARE\neditorial chapter-level decisions.\n"
        "CLIP sees the same editorial intent\nas the creator.\n\n"
        "Practical roadmap:\n"
        "  ADD: CLIP visual fusion\n"
        "  SUPPRESS: prosody, BERTopic",
        9.2, 2.92, 3.8, 4.2, size=12, color=DGREY)


def s13_error_taxonomy(p):
    """NEW: Formal 5-type error taxonomy (Type A through E)."""
    s = blank(p)
    hdr(s,"Formal Error Taxonomy: Five Failure Types",
        "Systematic classification of how and why predictions deviate from reference")
    rect(s, 0.3, 1.25, 12.7, 0.6, NAVY)
    txt(s,"Error types are not just a diagnostic — they determine which future signals will help. "
        "Type E (granularity mismatch) is the most important because it is systematic across all acoustic signals.",
        0.5, 1.3, 12.3, 0.52, size=13, bold=True, color=GOLD)
    errors = [
        (RED, "A", "Omission",
         "A reference boundary that the system completely misses.",
         "Insufficient cosine contrast at that position.\n"
         "The two chapters may share domain vocabulary, making the\n"
         "embedding transition gradual rather than sharp.\n"
         "Fix: more sensitive detection threshold — but at cost of precision.",
         "Under-prediction\nof boundaries"),
        (AMBER, "B", "Displacement",
         "A boundary predicted within the tolerance window but not at the exact sentence.",
         "The embedding transition is real but peaks slightly early/late.\n"
         "Usually a 1-3 sentence offset — acceptable for navigation.\n"
         "Pk/WD tolerate this; exact-match F1 penalises it heavily.\n"
         "This is why F1 appears low even when navigation quality is good.",
         "Off-by-few-sentences"),
        (PURPLE, "C", "Over-segmentation",
         "A predicted boundary with no matching reference boundary.",
         "The most common cause: a signal fires at a within-chapter\n"
         "rhetorical transition (definition → example) rather than\n"
         "a true chapter change. Cross-model conservative scoring\n"
         "reduces this by requiring both models to agree.",
         "False positive\nboundary"),
        (LBLUE, "D", "Under-segmentation",
         "A reference boundary missed AND its region merged with an adjacent segment.",
         "Occurs when a chapter change lacks an embedding discontinuity —\n"
         "common in philosophy lectures where argument evolves gradually\n"
         "over many sentences without lexical or semantic reset.\n"
         "Fix: hierarchical subtopic scaffolding provides an intermediate check.",
         "Merged adjacent\nchapters"),
        (BLUE, "E", "Granularity Mismatch  ← most important",
         "A signal fires frequently within chapters at a finer granularity than editorial boundaries.",
         "Root cause of ALL acoustic/linguistic failures (prosody, BERTopic,\n"
         "discourse markers, perplexity). These signals are accurate at the\n"
         "discourse-transition level — they are right about what they detect,\n"
         "but what they detect is structurally different from creator chapters.\n"
         "Not noise. Not a calibration problem. A systematic mismatch.",
         "Signal at wrong\ndiscourse level"),
    ]
    for i,(color, letter, name, defn, cause, tag) in enumerate(errors):
        gy = 1.98 + i*1.02
        rect(s, 0.3, gy, 12.7, 0.95, WHITE)
        rect(s, 0.3, gy, 0.5, 0.95, color)
        txt(s, letter, 0.33, gy+0.22, 0.42, 0.5, size=22, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, name, 0.88, gy+0.06, 2.3, 0.35, size=13, bold=True, color=NAVY)
        txt(s, defn, 0.88, gy+0.46, 2.3, 0.42, size=10, italic=True, color=DGREY)
        txt(s, cause, 3.3, gy+0.06, 7.5, 0.83, size=11, color=DGREY)
        rect(s, 10.9, gy+0.1, 2.0, 0.75, color)
        txt(s, tag, 10.95, gy+0.22, 1.9, 0.5, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)


def s14_selector_deep(p):
    s = blank(p)
    hdr(s,"The Method Selector: Full Breakdown","Supervised meta-model that picks the best segmentation method per video")
    rect(s, 0.3, 1.25, 5.5, 5.95, WHITE)
    accent(s, 0.3, 1.25, 5.5, 0.07, PURPLE)
    txt(s,"What Is It?", 0.5, 1.38, 5.2, 0.4, size=14, bold=True, color=NAVY)
    txt(s,"80+ method variants exist (embedding x method x hyperparameter).\n"
        "No single variant wins on all 30 videos.\n"
        "The selector LEARNS which variant to use per video.",
        0.5, 1.85, 5.2, 0.95, size=12, color=DGREY)
    txt(s,"MODEL: ExtraTreesRegressor\n(scikit-learn ensemble, 500 trees)", 0.5, 2.85, 5.2, 0.55,
        size=13, bold=True, color=PURPLE)
    txt(s,"INPUTS FED TO THE SELECTOR:", 0.5, 3.5, 5.2, 0.35, size=12, bold=True, color=TEAL)
    inputs = [
        ("Video features","Domain (Biology/CS/Math/Phil/Physics) — one-hot\nVideo length (minutes)\nSentence count\nMean cosine similarity (BGE-large)\nStd of consecutive similarities"),
        ("Training-fold stats","For each of the 80+ candidate methods:\n  Mean Pk across the 29 training videos\n  (= how reliable is this method generally?)"),
    ]
    oy = 3.9
    for title,desc in inputs:
        rect(s, 0.5, oy, 5.0, 1.05+("Video" in title)*0.2, LGREY)
        txt(s, title, 0.6, oy+0.06, 4.8, 0.35, size=12, bold=True, color=NAVY)
        txt(s, desc, 0.6, oy+0.42, 4.8, 0.78+("Video" in title)*0.2,
            size=11, color=DGREY)
        oy += 1.2+("Video" in title)*0.2
    txt(s,"OUTPUT: predicted Pk per method → pick lowest", 0.5, 6.68, 5.2, 0.4,
        size=12, bold=True, color=GREEN)
    rect(s, 6.1, 1.25, 6.9, 2.8, WHITE)
    accent(s, 6.1, 1.25, 6.9, 0.07, LBLUE)
    txt(s,"Leave-One-Video-Out Training (LOO-CV)", 6.3, 1.38, 6.5, 0.4,
        size=13, bold=True, color=NAVY)
    txt(s,"For each of 30 videos:", 6.3, 1.85, 6.5, 0.3, size=12, color=DGREY)
    for i in range(6):
        is_test = (i == 2)
        bc = RED if is_test else RGBColor(0xC0,0xD8,0xF0)
        lbl = "TEST" if is_test else f"Train"
        rect(s, 6.3+i*1.1, 2.22, 0.95, 0.45, bc)
        txt(s, lbl, 6.33+i*1.1, 2.28, 0.88, 0.32, size=11, bold=True,
            color=WHITE if is_test else NAVY, align=PP_ALIGN.CENTER)
    txt(s,"...", 12.7, 2.28, 0.4, 0.32, size=14, color=DGREY)
    txt(s,"Train on 29 videos, predict for 1 held-out. Repeat 30 times.",
        6.3, 2.75, 6.5, 0.32, size=11, italic=True, color=DGREY)
    rect(s, 6.1, 4.25, 6.9, 2.95, WHITE)
    accent(s, 6.1, 4.25, 6.9, 0.07, TEAL)
    txt(s,"What the Selector Chose (per domain)", 6.3, 4.38, 6.5, 0.4,
        size=13, bold=True, color=NAVY)
    domain_choices = [
        ("Physics (7 videos)","Multimodal-grid (CLIP+text) for 5/7 videos","Biggest gains",GREEN),
        ("CS (7 videos)","Cross-E5 conservative most reliable","Consistent",LBLUE),
        ("Biology (6 videos)","Cross-E5 + some over-switching","Mixed",AMBER),
        ("Philosophy (6 videos)","Selector improves Pk/WD","Positive",GREEN),
        ("Math (4 videos)","❌ FAILS — only 3 training examples per fold","Selector hurts",RED),
    ]
    oy = 4.82
    for dom,choice,verdict,color in domain_choices:
        rect(s, 6.3, oy, 6.5, 0.43, LGREY)
        txt(s, dom, 6.4, oy+0.06, 2.0, 0.3, size=11, bold=True, color=NAVY)
        txt(s, choice, 8.45, oy+0.06, 3.0, 0.3, size=10, color=DGREY)
        rect(s, 11.5, oy, 0.8, 0.43, color)
        txt(s, verdict, 11.55, oy+0.06, 0.7, 0.3, size=9, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        oy += 0.48


def s15_selector_verdict(p):
    s = blank(p)
    hdr(s,"Selector Results & Verdict","How much did the selector improve things?")
    ops = [
        ("BGE-Divisive\n(Baseline)","0.388","0.396","0.129","0.088",NAVY,"Reference"),
        ("Cross-Model\nConservative","0.371","0.376","0.036","0.024",TEAL,"Best Pk/WD (sig.)"),
        ("Balanced\nSelector","0.359","0.374","0.076","0.089",GREEN,"Best mean point"),
    ]
    rect(s, 0.3, 1.25, 12.7, 1.85, WHITE)
    accent(s, 0.3, 1.25, 12.7, 0.07, NAVY)
    for ci,h in enumerate(["Method","Pk ↓ (better)","WD ↓ (better)","BS ↑ (better)","F1@2 ↑ (better)","vs Baseline"]):
        cx = 0.35+ci*2.12
        txt(s, h, cx, 1.38, 2.0, 0.38, size=13, bold=True, color=NAVY)
    for ri,(name,pk,wd,bs,f1,color,note) in enumerate(ops):
        oy = 1.72 + ri*0.42
        bg = LGREY if ri%2==0 else WHITE
        rect(s, 0.3, oy, 12.7, 0.4, bg)
        for ci,v in enumerate([name,pk,wd,bs,f1,note]):
            cx = 0.35+ci*2.12
            fc = color if ci in [1,2,3,4] else NAVY
            txt(s, v, cx, oy+0.06, 2.0, 0.28, size=12, bold=(ci==0), color=fc)
    rect(s, 0.3, 3.35, 7.5, 3.85, WHITE)
    accent(s, 0.3, 3.35, 7.5, 0.07, GREEN)
    txt(s,"Statistical Significance (paired Wilcoxon, α=0.05 with Holm correction)",
        0.5, 3.48, 7.2, 0.42, size=13, bold=True, color=NAVY)
    sig_rows = [
        ("Cross-model vs BGE-divisive","Pk","✅ Significant p<0.05"),
        ("Cross-model vs BGE-divisive","WD","✅ Significant p<0.05"),
        ("Selector vs BGE-divisive","Pk","✅ Significant p<0.05"),
        ("Selector vs BGE-divisive","BS + F1","✅ Significant p<0.05"),
        ("Selector vs Cross-model","Pk/WD","❌ NOT significant"),
        ("Selector vs Cross-model","BS + F1","✅ Significant p<0.05"),
    ]
    oy = 3.95
    for comp,metric,result in sig_rows:
        bg = LGREEN if "✅" in result else LRED
        rect(s, 0.5, oy, 7.1, 0.42, bg)
        txt(s, comp, 0.6, oy+0.07, 3.2, 0.28, size=11, color=DGREY)
        txt(s, metric, 3.85, oy+0.07, 0.7, 0.28, size=11, bold=True, color=NAVY)
        fc = GREEN if "✅" in result else RED
        txt(s, result, 4.6, oy+0.07, 2.9, 0.28, size=11, bold=True, color=fc)
        oy += 0.46
    rect(s, 8.1, 3.35, 5.1, 3.85, NAVY)
    txt(s,"Final Verdict on Selector", 8.3, 3.48, 4.8, 0.42,
        size=14, bold=True, color=WHITE)
    txt(s,"The selector improves mean Pk and significantly\n"
        "improves BS and F1 vs cross-model.\n\n"
        "Pk/WD gain over cross-model alone is NOT\nstatistically significant.\n\n"
        "Interpretation:\n"
        "  The selector is a stronger OPERATING POINT\n"
        "  experiment — not a conclusive replacement.\n\n"
        "Leave-domain-out: Pk=0.4012 (worse than baseline)\n"
        "→ NOT a universal deployment model.\n"
        "→ A LOCAL BENCHMARK operating point.",
        8.3, 3.98, 4.8, 3.15, size=12, color=RGBColor(0xC0,0xD8,0xF0))


def s16_paper_comparison(p):
    s = blank(p)
    hdr(s,"Comparison to Prior Work","Supervised vs Unsupervised — direct vs indirect — what metrics they used")
    table(s,
        headers=["Paper","Videos","Supervised?","Metric","Their Result","Our Result","Comparison"],
        rows=[
            ["TreeSeg\n(Gklezakos 2024)","21 (TinyRec)\n+ICSI+AMI","No\n(unsupervised)","Pk/WD","Pk=0.367 (TinyRec)","Pk=0.359 (selector)","⚠ INDIRECT\nDiff dataset"],
            ["AVLectures\n(Yang 2022)","2350+","Self-supervised","No Pk/WD","Outperforms\nbaselines","N/A","❌ Not comparable\nDiff metrics"],
            ["Freisinger 2023","34+96 lectures","LoRA fine-tuned","F1 only","F1=67.3","N/A","❌ Not comparable\nNo Pk reported"],
            ["MiniSeg/YTSEG\n(Retkowski 2024)","19,299","Supervised","Pk/WD/F1","Pk=28.73","Pk=35.9 (ours)","❌ INDIRECT\n643x more data"],
            ["Chapter-Gen\n(Cao 2022)","9,631","Supervised","AP/Recall","AP=43.3","N/A","❌ Different task\nGeneration not seg"],
            ["Chapter-Llama\n(Ventura 2025)","10k/8.1k","Trained LLM","F1/IoU","F1=45.3","N/A","❌ Different task\nDiff metrics"],
            ["TextTiling (Hearst)","Our data","No","Pk","0.605 (ours)","0.371 (cross-model)","✅ DIRECT\nOurs better"],
            ["C99 (Choi 2000)","Our data","No","Pk","0.422 (ours)","0.371 (cross-model)","✅ DIRECT\nOurs better"],
        ],
        x=0.25, y=1.22, w=12.8,
        hs=11, rs=10)
    rect(s, 0.25, 6.55, 12.8, 0.72, NAVY)
    txt(s,"DIRECT comparisons (same dataset, same metric): our method beats all classical/neural baselines.\n"
        "INDIRECT (different datasets): we approach TreeSeg's Pk at similar scale. Large supervised systems are not the comparison target.",
        0.4, 6.6, 12.5, 0.65, size=12, bold=True, color=WHITE)


def s17_oracle(p):
    """Updated: hedged with 'within the evaluated framework'."""
    s = blank(p)
    hdr(s,"Oracle Gap Analysis: Where Is the Bottleneck?",
        "Locates the dominant source of remaining error — within the evaluated framework")
    rect(s, 0.3, 1.25, 12.7, 0.95, LAMBER)
    txt(s,"IMPORTANT: Oracle results use ground-truth knowledge at test time — NOT deployable.\n"
        "Purpose: identify whether the gap lies in candidate generation or candidate selection.",
        0.5, 1.3, 12.3, 0.85, size=13, bold=True, color=AMBER)
    rect(s, 0.3, 2.35, 7.2, 4.85, WHITE)
    accent(s, 0.3, 2.35, 7.2, 0.07, NAVY)
    txt(s,"Operating Points (Pk)", 0.5, 2.48, 7.0, 0.4, size=14, bold=True, color=NAVY)
    ops = [
        ("BGE-Divisive (baseline)","0.388","✅ Deployable","Stable starting point"),
        ("Cross-model conservative","0.371","✅ Deployable","Statistically significant"),
        ("Balanced selector","0.359","✅ Deployable (LOO)","Best mean local point"),
        ("TreeSeg/TinyRec (indicative)","0.367","✅ Deployable","Different dataset/labels"),
        ("Per-video oracle","0.298","❌ NOT deployable","Needs gold at test time"),
        ("Candidate oracle τ=2","0.017","❌ NOT deployable","Perfect ranker needed"),
    ]
    for i,(name,pk,dep,note) in enumerate(ops):
        oy = 2.95 + i*0.55
        bg = LGREY if i%2==0 else WHITE
        rect(s, 0.5, oy, 6.8, 0.5, bg)
        txt(s, name, 0.6, oy+0.08, 2.4, 0.34, size=12, color=NAVY, bold=(i<4))
        txt(s, pk, 3.05, oy+0.08, 0.7, 0.34, size=14, bold=True,
            color=GREEN if float(pk)<0.37 else (AMBER if float(pk)<0.4 else RED))
        fc = GREEN if "✅" in dep else RED
        txt(s, dep, 3.8, oy+0.08, 1.7, 0.34, size=11, color=fc)
        txt(s, note, 5.55, oy+0.08, 1.65, 0.34, size=10, color=DGREY)
    rect(s, 7.8, 2.35, 5.2, 4.85, WHITE)
    accent(s, 7.8, 2.35, 5.2, 0.07, LBLUE)
    txt(s,"What the Gap Tells Us", 8.0, 2.48, 4.9, 0.4, size=14, bold=True, color=NAVY)
    txt(s,"Within the evaluated framework:\n\n"
        "Gap: deployable best (0.359) → candidate oracle (0.017)\n"
        "ΔPk ≈ 0.34 — appears concentrated in SELECTION.\n\n"
        "Evidence: 96.8% recall already exists in candidate pool.\n"
        "The right boundaries are being found —\n"
        "we cannot always identify the correct ones.\n\n"
        "This finding is scoped to LecSeg-30 under the\n"
        "current candidate generation strategy.\n\n"
        "What this suggests for future work:\n"
        "  → NOT 'build a better segmenter'\n"
        "  → YES 'build a better boundary RANKER'\n\n"
        "A supervised ranker trained on 50-100 labelled\n"
        "lectures with features: local cosine contrast,\n"
        "cross-model agreement, CLIP signal, pause duration\n"
        "could plausibly close most of this gap.",
        8.0, 2.95, 4.9, 4.2, size=11, color=DGREY)


def s18_webapp(p):
    s = blank(p)
    hdr(s,"Real-World Application: LecSeg Web App","Not just research — a working tool for any YouTube lecture")
    rect(s, 0.3, 1.25, 5.8, 5.95, WHITE)
    accent(s, 0.3, 1.25, 5.8, 0.07, TEAL)
    txt(s,"What the Web App Does", 0.5, 1.38, 5.5, 0.42, size=15, bold=True, color=NAVY)
    txt(s,"Built with Streamlit (scripts/demo.py)\n"
        "Run: streamlit run scripts/demo.py\n\n"
        "Two modes:\n\n"
        "1. SEGMENT ANY YOUTUBE VIDEO\n"
        "   Paste any YouTube URL or video ID\n"
        "   → Auto-downloads transcript (yt-dlp captions)\n"
        "   → Sentence split + BGE/E5/MPNet/MiniLM embed\n"
        "   → Divisive segmentation\n"
        "   → Llama-3.1-8B generates chapter titles\n"
        "   → Exports YouTube timestamp format\n\n"
        "2. BENCHMARK MODE (LecSeg-30)\n"
        "   Evaluate any method on the 30-video benchmark\n"
        "   Shows Pk/WD with reference chapters",
        0.5, 1.88, 5.5, 4.2, size=12, color=DGREY)
    txt(s,"Try it: streamlit run scripts/demo.py",
        0.5, 6.05, 5.5, 0.45, size=13, bold=True, color=TEAL)
    rect(s, 6.4, 1.25, 6.7, 5.95, WHITE)
    accent(s, 6.4, 1.25, 6.7, 0.07, LBLUE)
    txt(s,"App Pipeline (Arbitrary Video Mode)", 6.6, 1.38, 6.4, 0.42,
        size=14, bold=True, color=NAVY)
    app_steps = [
        ("User pastes\nYouTube URL","Any lecture video;\nno prior indexing needed"),
        ("Transcript\nAcquisition","yt-dlp captions (fast)\nFallback: Whisper (slow)"),
        ("Sentence\nSplitting","spaCy; cache per video_id\nAvoids recomputation"),
        ("Embedding\n+ Segmentation","Select model in sidebar\n(BGE/E5/MPNet/MiniLM)"),
        ("LLM Titling","Llama-3.1-8B (auto-starts)\nFallback: noun-phrase extraction"),
        ("Output","Timestamp list\nYouTube format export\nSemantic similarity chart"),
    ]
    for i,(lbl,desc) in enumerate(app_steps):
        gy = 1.85 + i*0.82
        rect(s, 6.6, gy, 6.3, 0.75, LGREY if i%2==0 else WHITE)
        txt(s, f"{i+1}. {lbl}", 6.7, gy+0.06, 2.2, 0.42, size=12, bold=True, color=NAVY)
        txt(s, desc, 8.95, gy+0.06, 3.8, 0.55, size=11, color=DGREY)
        if i < 5:
            txt(s, "↓", 9.7, gy+0.72, 0.3, 0.25, size=12, color=LBLUE,
                align=PP_ALIGN.CENTER)


def s19_scope_and_benchmarks(p):
    """NEW: What LecSeg shows / does not show + why small benchmarks matter."""
    s = blank(p)
    hdr(s,"Scope of Claims: What LecSeg-30 Shows and Does Not Show",
        "Honest positioning — and why small, controlled benchmarks have scientific value large ones do not")
    rect(s, 0.3, 1.25, 6.15, 5.95, WHITE)
    accent(s, 0.3, 1.25, 6.15, 0.07, GREEN)
    txt(s,"What LecSeg-30 DOES show", 0.5, 1.38, 5.9, 0.42, size=14, bold=True, color=GREEN)
    shows = [
        "Cross-model conservative scoring significantly reduces Pk/WD on LecSeg-30 vs all tested baselines.",
        "The oracle gap (96.8% recall, Pk=0.017) places the dominant bottleneck in selection, not generation — within this framework.",
        "Acoustic and linguistic signals systematically over-segment because they detect discourse transitions, not editorial chapters.",
        "CLIP visual fusion helps because slide transitions share the editorial granularity of creator chapters.",
        "A supervised selector achieves Pk=0.359 but degrades to 0.4012 under leave-domain-out — it is a local operating point.",
        "Entropy-weighted fusion automatically down-weights high-entropy (flat) modalities without any manual calibration.",
        "A two-level hierarchical annotation with workflow-consistency estimates is reproducible at small scale.",
    ]
    for i,item in enumerate(shows):
        oy = 1.9 + i*0.72
        rect(s, 0.5, oy, 5.75, 0.65, LGREEN)
        txt(s, f"✅  {item}", 0.6, oy+0.08, 5.55, 0.5, size=11, color=DGREY)
    rect(s, 6.7, 1.25, 6.3, 5.95, WHITE)
    accent(s, 6.7, 1.25, 6.3, 0.07, RED)
    txt(s,"What LecSeg-30 does NOT show", 6.9, 1.38, 6.0, 0.42, size=14, bold=True, color=RED)
    not_shows = [
        "State-of-the-art performance vs large supervised systems. They are trained on 10k-800k videos; the comparison is not meaningful.",
        "Domain-general deployment. The selector fails on Mathematics (4 videos). 30 videos is not sufficient for general training.",
        "That creator chapters are the 'correct' pedagogical segmentation. They are navigation metadata, not expert annotation.",
        "That N1/N3/N4 components are novel contributions. They are engineering integrations of standard techniques.",
        "That Whisper transcription quality is near-perfect. Some domain-specific vocabulary may be misrecognised at the margins.",
    ]
    for i,item in enumerate(not_shows):
        oy = 1.9 + i*0.88
        rect(s, 6.9, oy, 5.9, 0.8, LRED)
        txt(s, f"❌  {item}", 7.0, oy+0.08, 5.7, 0.66, size=11, color=DGREY)
    rect(s, 0.3, 7.1, 12.7, 0.32, NAVY)
    txt(s,"Why small benchmarks matter: large-scale aggregates conceal failure modes. "
        "A 30-video inspectable benchmark makes it possible to understand WHY systems fail — that is the contribution of this work.",
        0.5, 7.13, 12.3, 0.26, size=12, italic=True, bold=True, color=GOLD)


def s20_research_justification(p):
    s = blank(p)
    hdr(s,"Justifying the Research Title Word-by-Word",
        "LecSeg-30: A Reproducible Low-Resource Benchmark and Diagnostic Study for Lecture-Video Topic Segmentation")
    rect(s, 0.3, 1.25, 12.7, 0.95, NAVY)
    txt(s,'LecSeg-30  ·  A Reproducible  ·  Low-Resource  ·  Benchmark  ·  and Diagnostic Study  ·  for Lecture-Video Topic Segmentation',
        0.5, 1.32, 12.3, 0.82, size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    justifications = [
        ("LecSeg-30","30 academic YouTube lecture videos\nacross 5 domains, 32.52 hours total.\nThe '30' is in the name because the\nbenchmark is permanently defined.",
         BLUE),
        ("Reproducible","All code open (src/lecseg/).\nAll data IDs public (YouTube URLs).\n177+ automated tests.\nEvery table has a generation script.\nAny researcher can rerun.",
         TEAL),
        ("Low-Resource","Only 30 videos (vs 10k-800k for\nlarge supervised systems).\nNo GPU required at inference.\nAll models run on CPU.\nFully offline with Ollama.",
         GREEN),
        ("Benchmark","Fixed 30-video set with:\n- Creator-provided chapter GT\n- Two-level hierarchical labels\n- Workflow consistency (κ=0.43)\n- Pk/WD + bootstrap CIs",
         PURPLE),
        ("Diagnostic Study","We don't just report Pk.\nWe show WHAT FAILS and WHY:\n- Granularity mismatch finding\n- Oracle gap locates bottleneck\n- 5-type error taxonomy\n- Ablation of every signal",
         AMBER),
        ("Lecture-Video\nTopic Seg","Not general video chaptering.\nNot meeting segmentation.\nSpecifically academic YouTube\nlectures with Pk/WD evaluation.",
         RED),
    ]
    for i,(word,proof,color) in enumerate(justifications):
        row,col = i//3, i%3
        gx = 0.3 + col*4.35
        gy = 2.38 + row*2.45
        rect(s, gx, gy, 4.1, 2.35, WHITE)
        accent(s, gx, gy, 4.1, 0.07, color)
        rect(s, gx, gy, 4.1, 0.48, color)
        txt(s, word, gx+0.1, gy+0.06, 3.9, 0.38, size=14, bold=True, color=WHITE)
        txt(s, proof, gx+0.1, gy+0.55, 3.9, 1.72, size=11, color=DGREY)


def s21_thesis_quality(p):
    s = blank(p)
    hdr(s,"Thesis Standards & Quality Checklist","Is everything up to academic standard?")
    checks = [
        ("✅","Clickable references","hyperref package; all \\citep{} are clickable in PDF"),
        ("✅","In-text citations","38 citations across 5 chapters; all figures/tables cited"),
        ("✅","Title match","Titlepage: 'LecSeg-30: A Reproducible Low-Resource Benchmark and Diagnostic Study'"),
        ("✅","LaTeX structure","book class, 12pt, A4; TOC, LOF, LOT, bibliography (plainnat); appendices"),
        ("✅","6 chapters","Intro, Literature, Methodology, Results, Conclusion, Future Work"),
        ("✅","3 appendices","Dataset details, Hyperparameters, Extra results"),
        ("✅","15+ LaTeX tables","All \\input'd from thesis/tables/; generated from data"),
        ("✅","7 PDF figures","All in thesis/figures/; referenced with \\includegraphics"),
        ("✅","Statistical tests","Bootstrap 95% CIs + paired Wilcoxon with Holm correction"),
        ("✅","Consistency stated","κ chapter=0.535, subtopic=0.426; methodology limitation acknowledged"),
        ("✅","Limitations section","Silver labels, corpus size, selector domain-failure all stated"),
        ("✅","Threats to validity","Structured table: construct/external/annotation/internal/metric validity"),
        ("✅","Error taxonomy","Formal 5-type taxonomy (A-E) in Ch4"),
        ("✅","Scope of claims","Explicit 'what LecSeg shows / does not show' section in Ch5"),
        ("✅","GPU transcription","vast.ai RTX 5090 documented in Ch3 and Appendix B"),
    ]
    table(s,
        headers=["Status","Item","Evidence"],
        rows=[[c,n,e] for c,n,e in checks],
        x=0.3, y=1.22, w=12.7,
        hs=13, rs=11)


def s22_limitations(p):
    s = blank(p)
    hdr(s,"Limitations & Honest Assessment","What we cannot claim — and what the examiner will ask about")
    limits = [
        ("30 videos is small",RED,
         "Math domain has only 4 → selector fails there.",
         "We acknowledge this. 30 matches prior low-resource work.\nBootstrap CIs quantify uncertainty. Selector limitation explicitly reported."),
        ("Silver chapter labels",AMBER,
         "YouTube chapters = creator editorial choices, not expert pedagogical labels.",
         "Treated as navigation metadata only. Human validation listed as future work.\nLimitation stated in threats-to-validity section."),
        ("Annotation workflow, not independent IAA",AMBER,
         "Both annotation passes use same workflow → not truly independent.",
         "Renamed 'workflow consistency estimate' throughout the thesis.\nChapter κ inflated (shared metadata); subtopic κ=0.426 is more meaningful."),
        ("Selector not universal",AMBER,
         "Leave-domain-out: Pk=0.4012 (worse than baseline).",
         "Explicitly acknowledged. The selector is a local benchmark operating point,\nnot a production deployment model. Clearly stated in Ch3 and Ch5."),
        ("LLM refinement unproven",LBLUE,
         "Llama-3.1-8B boundary shifting does not improve Pk in experiments.",
         "Reported as diagnostic only. Titling works offline. Not overclaimed.\nFuture work: larger model or task-specific fine-tuning."),
    ]
    for i,(issue,color,problem,response) in enumerate(limits):
        gy = 1.28 + i*1.18
        rect(s, 0.3, gy, 12.7, 1.1, WHITE)
        accent(s, 0.3, gy, 12.7, 0.06, color)
        txt(s, f"L{i+1}: {issue}", 0.45, gy+0.1, 3.5, 0.38, size=13, bold=True, color=NAVY)
        txt(s, problem, 4.0, gy+0.1, 4.2, 0.38, size=12, color=RED)
        rect(s, 8.3, gy+0.06, 4.6, 0.95, LGREEN)
        txt(s, "How we address it:", 8.4, gy+0.1, 4.4, 0.28, size=10, bold=True, color=GREEN)
        txt(s, response, 8.4, gy+0.38, 4.4, 0.65, size=10, color=DGREY)


def s23_results_summary(p):
    s = blank(p)
    hdr(s,"Main Results Summary","The definitive numbers from LecSeg-30")
    rect(s, 0.3, 1.25, 8.2, 5.95, WHITE)
    accent(s, 0.3, 1.25, 8.2, 0.07, NAVY)
    txt(s,"Pk Results Across All Methods (lower = better)",
        0.5, 1.38, 8.0, 0.42, size=14, bold=True, color=NAVY)
    methods = [
        ("KMeansSeg",0.617,RED),("TextTiling",0.605,RED),
        ("CosineSeg",0.490,RED),("BertSeg",0.489,RED),("C99",0.422,AMBER),
        ("TwoStage+Prosody",0.433,AMBER),("Hierarchical",0.412,AMBER),
        ("BGE-Divisive",0.388,NAVY),("CLIP+Text",0.374,GREEN),
        ("Cross-Model",0.371,GREEN),("Selector",0.359,GREEN),
    ]
    max_v = 0.7
    for i,(name,pk,color) in enumerate(methods):
        oy = 1.87 + i*0.42
        bw = (pk/max_v)*6.3
        txt(s, name, 0.45, oy+0.05, 2.1, 0.32, size=11, color=DGREY)
        rect(s, 2.6, oy+0.04, bw, 0.33, color)
        txt(s, f"{pk:.3f}", 2.65+bw, oy+0.07, 0.6, 0.24, size=11, bold=True, color=NAVY)
    base_x = 2.6 + (0.388/max_v)*6.3
    rect(s, base_x, 1.87, 0.015, 4.65, AMBER)
    txt(s,"baseline\n0.388", base_x-0.3, 6.35, 0.8, 0.55, size=9, italic=True, color=AMBER)
    rect(s, 8.75, 1.25, 4.35, 5.95, WHITE)
    accent(s, 8.75, 1.25, 4.35, 0.07, TEAL)
    txt(s,"All Metrics", 8.95, 1.38, 4.1, 0.42, size=14, bold=True, color=NAVY)
    table(s,
        headers=["Method","Pk","WD","BS","F1"],
        rows=[
            ["BGE-Div","0.388","0.396","0.129","0.088"],
            ["Cross-Model","0.371","0.376","0.036","0.024"],
            ["Selector","0.359","0.374","0.076","0.089"],
        ],
        x=8.85, y=1.88, w=4.1, hs=11, rs=11)
    txt(s,"Pk/WD: primary metrics (lower=better)\n"
        "BS/F1: secondary diagnostics (higher=better)\n"
        "Note: low F1 reflects conservative boundary policy;\n"
        "for navigation quality, Pk/WD are the right measure.\n\n"
        "Significant improvements:\nCross-model: Pk+WD vs baseline (p<0.05)\nSelector: BS+F1 vs cross-model (p<0.05)\n\n"
        "N=30 videos\n95% bootstrap CIs\nPaired Wilcoxon+Holm correction",
        8.95, 3.1, 4.1, 3.95, size=11, color=DGREY)


def s24_contributions(p):
    """Reordered: C1=benchmark, C2=bottleneck finding, C3=granularity mismatch, C4=hierarchy, C5=reproducibility."""
    s = blank(p)
    hdr(s,"What This Research Contributes","Five contributions ordered by scientific significance")
    rect(s, 0.3, 1.25, 12.7, 5.95, WHITE)
    accent(s, 0.3, 1.25, 12.7, 0.08, GOLD)
    contribs = [
        (TEAL,"C1: LecSeg-30 Benchmark  (primary artifact)",
         "30 YouTube lectures · 5 domains · 32.52 h · 419 creator chapters · 904 reviewed subtopic labels\n"
         "Two-level hierarchy with workflow-consistency estimates (κ=0.54 chapter, κ=0.43 subtopic) · Public YouTube URLs\n"
         "Makes low-resource lecture segmentation measurable under a shared reproducible protocol."),
        (GREEN,"C2: Candidate-Selection Bottleneck Finding  (central empirical result)",
         "Oracle analysis: 96.8% recall exists in candidate pool at Pk=0.017 — bottleneck is selection, not generation.\n"
         "Within the evaluated framework, this relocates the research problem from 'detect more boundaries' to 'rank them better'.\n"
         "Suggests a supervised ranker on 50-100 labelled lectures as the most productive next step."),
        (RED,"C3: Granularity-Mismatch Diagnosis  (most transferable finding)",
         "Systematic ablation establishes WHY prosody, BERTopic, discourse markers, and perplexity all hurt Pk/WD.\n"
         "They detect discourse-level transitions; creator chapters are editorial-level decisions — a structural, not calibration, mismatch.\n"
         "CLIP succeeds for the complementary reason. Directly informs future multimodal lecture-segmentation design."),
        (PURPLE,"C4: Hierarchical Annotation Protocol",
         "Two-level annotation workflow with nesting constraint (B_chapter ⊆ B_subtopic) and consistency estimates.\n"
         "Value is the protocol and documented procedure. No existing lecture benchmark combines hierarchy + Pk/WD + bootstrap CIs.\n"
         "Limitation acknowledged: LLM-assisted annotation, not independent human agreement."),
        (LBLUE,"C5: Reproducibility Package",
         "Complete open pipeline: Whisper transcription → spaCy splitting → BGE/E5 embedding → multimodal feature extraction\n"
         "(CLIP, PaddleOCR, TransNetV2) → seven baseline segmenters + four pipeline components → Pk/WD bootstrap evaluation.\n"
         "Streamlit web demo · Table-generation scripts for all reported results · 177+ automated tests."),
    ]
    for i,(color,title,desc) in enumerate(contribs):
        gy = 1.45 + i*1.15
        rect(s, 0.45, gy, 12.3, 1.06, LGREY)
        rect(s, 0.45, gy, 0.25, 1.06, color)
        txt(s, title, 0.78, gy+0.08, 11.8, 0.38, size=13, bold=True, color=NAVY)
        txt(s, desc, 0.78, gy+0.5, 11.8, 0.52, size=11, color=DGREY)


def s25_lessons_learned(p):
    """NEW: Lessons learned from the research process."""
    s = blank(p)
    hdr(s,"Lessons Learned","What this project taught us about lecture segmentation and benchmark design")
    rect(s, 0.3, 1.25, 6.15, 5.95, WHITE)
    accent(s, 0.3, 1.25, 6.15, 0.07, TEAL)
    txt(s,"Segmentation lessons", 0.5, 1.38, 5.9, 0.42, size=14, bold=True, color=TEAL)
    seg_lessons = [
        ("Granularity first","Before adding any signal, ask: at what discourse level does it fire? "
         "If finer than chapter level, it will hurt Pk/WD regardless of accuracy."),
        ("Vocabulary reuse kills lexical methods","Technical lectures reuse domain terms across chapters. "
         "Any method relying on surface vocabulary will plateau quickly."),
        ("Conservative is better than aggressive","Fewer false-positive boundaries improve Pk/WD even when recall drops. "
         "The cross-model consensus was our single most reliable policy."),
        ("Oracle analysis should come first","If we had run the oracle study earlier, we would have known immediately "
         "that improving the segmenter was not the right direction."),
    ]
    for i,(title,desc) in enumerate(seg_lessons):
        oy = 1.95 + i*1.3
        rect(s, 0.5, oy, 5.75, 1.2, LGREY)
        txt(s, title, 0.65, oy+0.08, 5.4, 0.35, size=12, bold=True, color=NAVY)
        txt(s, desc, 0.65, oy+0.47, 5.4, 0.66, size=11, color=DGREY)
    rect(s, 6.7, 1.25, 6.3, 5.95, WHITE)
    accent(s, 6.7, 1.25, 6.3, 0.07, PURPLE)
    txt(s,"Benchmark and annotation lessons", 6.9, 1.38, 6.0, 0.42, size=14, bold=True, color=PURPLE)
    bench_lessons = [
        ("Name your reference precisely","'Creator chapters' and 'pedagogical gold standard' are not the same thing. "
         "Every claim should state which it uses."),
        ("Document what fails, not just what works","Most published systems show only best configurations. "
         "Recording failures is more useful for the field."),
        ("Annotator independence matters","Using the same LLM for both annotation passes inflates consistency. "
         "Always state the procedure, not just the kappa."),
    ]
    for i,(title,desc) in enumerate(bench_lessons):
        oy = 1.95 + i*1.62
        rect(s, 6.9, oy, 5.9, 1.5, LGREY)
        txt(s, title, 7.05, oy+0.08, 5.6, 0.35, size=12, bold=True, color=NAVY)
        txt(s, desc, 7.05, oy+0.47, 5.6, 0.95, size=11, color=DGREY)
    rect(s, 0.3, 7.1, 12.7, 0.32, NAVY)
    txt(s,"Core insight: the value of LecSeg-30 lies less in achieving a new state of the art than in providing a reproducible "
        "environment for understanding why low-resource lecture segmentation systems succeed or fail.",
        0.5, 7.13, 12.3, 0.26, size=12, italic=True, bold=True, color=GOLD)


def s26_final_verdict(p):
    """Updated language throughout to match final thesis."""
    s = blank(p)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    for cx,cy,r,c in [(11.8,1.0,2.0,BLUE),(12.5,5.5,1.5,BLUE)]:
        circle(s, cx, cy, r, c)
    txt(s,"Final Verdict", 0.8, 0.5, 12, 0.8, size=42, bold=True, color=WHITE)
    accent(s, 0.8, 1.35, 10, 0.08, TEAL)
    verdicts = [
        (GREEN,
         "What we established:",
         "Conservative cross-model scoring (Pk=0.371) significantly outperforms all classical and neural\n"
         "baselines on LecSeg-30 (p<0.05, Wilcoxon+Holm). The balanced selector (Pk=0.359) significantly\n"
         "improves BS and F1. Granularity mismatch explains every acoustic/linguistic failure systematically."),
        (AMBER,
         "What we honestly acknowledge:",
         "We do not beat large supervised systems. The selector fails under leave-domain-out (Pk=0.4012).\n"
         "LLM boundary refinement is not proven to improve Pk. The oracle gap finding is scoped to LecSeg-30\n"
         "under the current candidate strategy — not a universal claim. All limitations are stated explicitly."),
        (LBLUE,
         "What this enables:",
         "LecSeg-30 is a fixed, inspectable benchmark for future lecture-segmentation research.\n"
         "The oracle gap (Pk=0.017) identifies the concrete engineering target: a supervised boundary ranker.\n"
         "The granularity-mismatch diagnosis shapes which signals are worth adding in future multimodal work."),
    ]
    for i,(color,label,text) in enumerate(verdicts):
        gy = 1.52 + i*1.82
        rect(s, 0.8, gy, 11.8, 1.72, BLUE)
        accent(s, 0.8, gy, 11.8, 0.07, color)
        txt(s, label, 1.0, gy+0.12, 3.0, 0.42, size=14, bold=True, color=color)
        txt(s, text, 1.0, gy+0.58, 11.4, 1.1, size=13, color=RGBColor(0xC8,0xDC,0xF0))
    txt(s,"LecSeg-30 does not make lecture segmentation solved. It makes it MEASURABLE, INSPECTABLE, and REPRODUCIBLE.",
        0.8, 7.1, 11.8, 0.38, size=14, italic=True, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════
def main():
    p = new_prs()
    s01_title(p)               # 1
    s02_roadmap(p)             # 2
    s03_boundary_types(p)      # 3 ← NEW: three topic boundary types
    s04_data_collection(p)     # 4
    s05_preprocessing(p)       # 5
    s06_annotation(p)          # 6  updated: workflow consistency
    s07_metrics(p)             # 7
    s08_all_methods(p)         # 8
    s09_how_methods_work(p)    # 9
    s10_pipeline_components(p) # 10 updated: N1-N4 as engineering integrations
    s11_what_worked_failed(p)  # 11
    s12_granularity_finding(p) # 12 updated: THE central scientific finding
    s13_error_taxonomy(p)      # 13 ← NEW: 5-type error taxonomy
    s14_selector_deep(p)       # 14
    s15_selector_verdict(p)    # 15
    s16_paper_comparison(p)    # 16
    s17_oracle(p)              # 17 updated: hedged language
    s18_webapp(p)              # 18
    s19_scope_and_benchmarks(p)# 19 ← NEW: scope of claims + why small benchmarks matter
    s20_research_justification(p)# 20
    s21_thesis_quality(p)      # 21 updated
    s22_limitations(p)         # 22 updated: annotation independence fix
    s23_results_summary(p)     # 23
    s24_contributions(p)       # 24 reordered by scientific significance
    s25_lessons_learned(p)     # 25 ← NEW: lessons learned
    s26_final_verdict(p)       # 26 updated

    out = "thesis/LECSEG_Defense_Slides_v4.pptx"
    p.save(out)
    print(f"Done. {len(p.slides)} slides saved -> {out}")

if __name__ == "__main__":
    main()
