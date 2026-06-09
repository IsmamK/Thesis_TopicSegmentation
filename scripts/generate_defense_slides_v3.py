"""
LECSEG Defense Slides v3 — comprehensive 25-slide deck.
Run: python scripts/generate_defense_slides_v3.py
Output: thesis/LECSEG_Defense_Slides_v3.pptx
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math

# ── Palette ───────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F,0x23,0x40)
BLUE   = RGBColor(0x1A,0x4A,0x8A)
LBLUE  = RGBColor(0x2D,0x6A,0xBF)
TEAL   = RGBColor(0x00,0x96,0x88)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LGREY  = RGBColor(0xF0,0xF4,0xF8)
MGREY  = RGBColor(0xD0,0xD8,0xE0)
DGREY  = RGBColor(0x37,0x47,0x51)
GREEN  = RGBColor(0x1B,0x87,0x3B)
LGREEN = RGBColor(0xE8,0xF5,0xE9)
RED    = RGBColor(0xC6,0x28,0x28)
LRED   = RGBColor(0xFF,0xEB,0xEE)
AMBER  = RGBColor(0xE6,0x5C,0x00)
LAMBER = RGBColor(0xFF,0xF3,0xCD)
GOLD   = RGBColor(0xF5,0xA6,0x23)
PURPLE = RGBColor(0x7B,0x1F,0xA2)

W, H = Inches(13.33), Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs_obj):
    return prs_obj.slides.add_slide(prs_obj.slide_layouts[6])

def rect(slide, x, y, w, h, fill=None, line=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if not line:
        s.line.fill.background()
    return s

def circle(slide, cx, cy, r, fill):
    s = slide.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=14, bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def hdr(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    txt(slide, title, 0.28, 0.06, 10, 0.65, size=28, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.28, 0.72, 12.5, 0.35, size=13, italic=True,
            color=RGBColor(0xA0,0xC4,0xF0))
    rect(slide, 0, 1.1, 13.33, 6.4, LGREY)

def accent(slide, x, y, w, h=0.055, color=TEAL):
    rect(slide, x, y, w, h, color)

def card(slide, x, y, w, h, bg=WHITE, border_color=None):
    r = rect(slide, x, y, w, h, bg)
    if border_color:
        accent(slide, x, y, w, 0.055, border_color)
    return r

def bar_chart(slide, x, y, w, h, values, labels, colors, title=None,
              show_values=True, max_val=None):
    """Simple horizontal bar chart."""
    if max_val is None:
        max_val = max(values) * 1.1 if values else 1
    n = len(values)
    bar_h = (h - 0.3) / n - 0.05
    if title:
        txt(slide, title, x, y, w, 0.3, size=12, bold=True, color=NAVY)
    for i, (val, label, color) in enumerate(zip(values, labels, colors)):
        by = y + 0.35 + i * (bar_h + 0.07)
        bw = (val / max_val) * (w - 1.8)
        txt(slide, label, x, by + 0.02, 1.7, bar_h, size=10, color=DGREY)
        r = rect(slide, x + 1.75, by, max(bw, 0.05), bar_h, color)
        if show_values:
            txt(slide, f"{val:.3f}", x+1.78+bw, by+0.02, 0.8, bar_h,
                size=10, bold=True, color=NAVY)

def cell(slide, x, y, w, h, bg, text, size=11, color=DGREY, bold=False,
         align=PP_ALIGN.CENTER):
    rect(slide, x, y, w, h, bg)
    txt(slide, text, x+0.04, y+0.05, w-0.08, h-0.06,
        size=size, color=color, bold=bold, align=align, wrap=True)

def table(slide, headers, rows, x, y, w,
          hbg=NAVY, hfg=WHITE, row_bgs=(LGREY,WHITE),
          hs=12, rs=11):
    nc = len(headers)
    cw = [w/nc]*nc
    rh = 0.4
    for ci,h in enumerate(headers):
        cx = x + sum(cw[:ci])
        cell(slide, cx, y, cw[ci], rh, hbg, h, size=hs, color=hfg, bold=True)
    for ri,row in enumerate(rows):
        bg = row_bgs[ri%2]
        for ci,v in enumerate(row):
            cx = x + sum(cw[:ci])
            ry = y + (ri+1)*rh
            fc = DGREY
            if str(v).startswith("✅"): fc = GREEN
            elif str(v).startswith("❌"): fc = RED
            elif str(v).startswith("⚠"): fc = AMBER
            cell(slide, cx, ry, cw[ci], rh, bg, str(v), size=rs, color=fc)

def step_flow(slide, steps, x, y, w, h, colors=None):
    """Horizontal flow diagram: [box] -> [box] -> ..."""
    n = len(steps)
    bw = (w - 0.25*(n-1)) / n
    if colors is None:
        colors = [BLUE]*n
    for i,(label,sub) in enumerate(steps):
        bx = x + i*(bw+0.25)
        rect(slide, bx, y, bw, h, colors[i])
        txt(slide, label, bx+0.05, y+0.05, bw-0.1, 0.38,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if sub:
            txt(slide, sub, bx+0.05, y+0.45, bw-0.1, h-0.5,
                size=10, color=RGBColor(0xC0,0xD8,0xF0),
                align=PP_ALIGN.CENTER)
        if i < n-1:
            txt(slide, "→", bx+bw, y+h/2-0.2, 0.28, 0.38,
                size=18, bold=True, color=LBLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════

def s01_title(p):
    s = blank(p)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    for cx,cy,r,c in [(11.8,1.0,2.0,BLUE),(12.5,5.5,1.5,BLUE),(0.6,6.8,0.8,BLUE)]:
        circle(s, cx, cy, r, c)
    accent(s, 0.8, 2.65, 7.5, 0.08, TEAL)
    txt(s,"LecSeg-30",0.8,0.9,11,1.1,size=60,bold=True,color=WHITE)
    txt(s,"A Reproducible Low-Resource Benchmark and Diagnostic Study",
        0.8,1.95,12,0.72,size=22,color=RGBColor(0xA0,0xC4,0xF0))
    txt(s,"for Lecture-Video Topic Segmentation",
        0.8,2.62,12,0.5,size=18,italic=True,color=RGBColor(0x80,0xB4,0xE0))
    # Team
    txt(s,"Ismam Khan  ·  Fahmida Afrin Moon  ·  Shahriar Islam Rafi  ·  Alimool Razi  ·  Sadia Alam",
        0.8,3.45,12,0.45,size=14,color=RGBColor(0xC0,0xD8,0xF0))
    txt(s,"Supervisor: Mr. Annajiat Alim Rasel",
        0.8,3.92,12,0.38,size=13,color=RGBColor(0x80,0xB4,0xE0))
    # 5 keywords from title
    for i,(word,desc) in enumerate([
        ("REPRODUCIBLE","All code, data, scripts public"),
        ("LOW-RESOURCE","Only 30 videos — no GPU farm needed"),
        ("BENCHMARK","30 lectures, 5 domains, 419 chapters"),
        ("DIAGNOSTIC","What works, what fails, and why"),
        ("LECTURE-VIDEO","Academic YouTube content specifically"),
    ]):
        gx = 0.8 + i*2.5
        rect(s, gx, 4.55, 2.3, 0.95, BLUE)
        txt(s, word, gx+0.08, 4.6, 2.15, 0.38, size=11, bold=True, color=GOLD,
            align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.08, 4.98, 2.15, 0.48, size=10,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
    txt(s,"BRAC University · Dept. of CSE · June 2026",
        0.8,6.8,12,0.45,size=13,color=RGBColor(0x60,0x90,0xC0))


def s02_roadmap(p):
    s = blank(p)
    hdr(s,"What This Talk Covers","15-minute roadmap")
    items = [
        ("1","Problem\n& Gap","Why lecture navigation is broken;\nwhat was missing"),
        ("2","Dataset\n& Setup","How we built LecSeg-30;\nGPU transcription on vast.ai"),
        ("3","Methods\nTried","All 11 approaches:\nwhat each does and why"),
        ("4","What\nWorked","Results per method with verdict;\ngranularity finding"),
        ("5","The\nSelector","Supervised meta-model:\nwhat we fed it, what it chose"),
        ("6","Comparison\nto Prior Work","Supervised vs unsupervised;\ndirect & indirect"),
        ("7","Web App\n& Verdict","Real-world application;\nfinal answer"),
    ]
    colors = [TEAL,BLUE,LBLUE,GREEN,PURPLE,AMBER,NAVY]
    for i,(n,label,desc) in enumerate(items):
        gx = 0.4 + i*1.86
        rect(s, gx, 1.35, 1.65, 2.2, colors[i])
        txt(s, n, gx+0.58, 1.42, 0.6, 0.5, size=26, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, label, gx+0.05, 1.92, 1.55, 0.55, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.05, 2.5, 1.55, 0.98, size=10,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
        if i < len(items)-1:
            txt(s, "→", gx+1.67, 2.15, 0.22, 0.38, size=16, bold=True,
                color=MGREY, align=PP_ALIGN.CENTER)
    # why this matters
    rect(s, 0.4, 3.85, 12.5, 3.4, WHITE)
    accent(s, 0.4, 3.85, 12.5, 0.07, GOLD)
    txt(s,"Why This Research Matters", 0.6, 4.0, 12, 0.42, size=16, bold=True, color=NAVY)
    cols = [
        ("Students","Cannot navigate 3-hour lectures\nwithout manual chapter markers"),
        ("Educators","Creating chapter markers manually\nis expensive and inconsistent"),
        ("Scale","Millions of YouTube lectures\nhave ZERO structured navigation"),
        ("Research","No compact lecture benchmark\nwith Pk/WD + hierarchy + stats existed"),
    ]
    for i,(title,desc) in enumerate(cols):
        gx = 0.6 + i*3.1
        rect(s, gx, 4.5, 2.9, 2.5, LGREY)
        txt(s, title, gx+0.1, 4.58, 2.7, 0.38, size=14, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, 5.0, 2.7, 1.8, size=12, color=DGREY)


def s03_data_collection(p):
    s = blank(p)
    hdr(s,"How We Built LecSeg-30","Data collection, GPU transcription, annotation — step by step")
    # Timeline flow
    steps = [
        ("YouTube\nSearch","5 domains;\n140+ candidate videos"),
        ("Quality\nFilter","Duration ≥15 min;\n≥4 chapters;\nEnglish audio"),
        ("30 Videos\nSelected","32.52 h total;\n419 creator chapters"),
        ("Audio\nExtraction","yt-dlp;\nffmpeg mp4→wav"),
        ("vast.ai\nGPU Rental","RTX 5090;\n~\$0.67/hr;\nshutdown after use"),
        ("Whisper\nTranscription","large-v3;\n17.9x realtime;\n46,525 sentences"),
        ("LLM Draft\nAnnotation","Llama-3.1-8B\n(Ollama local);\nsubtopic drafts"),
        ("Human\nReview","Streamlit tool;\ncorrect + approve\neach boundary"),
    ]
    colors_s = [TEAL,LBLUE,BLUE,BLUE,RED,BLUE,TEAL,GREEN]
    n = len(steps)
    bw = 1.52; gap = 0.03
    for i,(label,sub) in enumerate(steps):
        bx = 0.25 + i*(bw+gap)
        rect(s, bx, 1.25, bw, 1.85, colors_s[i])
        txt(s, str(i+1), bx+0.55, 1.3, 0.45, 0.38, size=18, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, label, bx+0.05, 1.68, bw-0.1, 0.55, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub, bx+0.05, 2.22, bw-0.1, 0.85, size=9,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
        if i < n-1:
            txt(s, "→", bx+bw, 1.95, 0.06, 0.38, size=14, bold=True,
                color=MGREY, align=PP_ALIGN.CENTER)
    # vast.ai callout
    rect(s, 0.25, 3.3, 4.5, 2.3, LRED)
    accent(s, 0.25, 3.3, 4.5, 0.07, RED)
    txt(s,"Why vast.ai?", 0.4, 3.45, 4.2, 0.38, size=14, bold=True, color=RED)
    txt(s,"Whisper large-v3 requires 16+ GB VRAM.\n"
        "Local hardware insufficient for 30+ hour corpus.\n"
        "Solution: rent RTX 5090 on vast.ai (~\$0.67/hr).\n"
        "Total cost: ~\$3-4 for all 30 videos.\n"
        "Instance destroyed immediately after download\nto stop billing.",
        0.4, 3.88, 4.2, 1.65, size=12, color=DGREY)
    # dataset stats
    rect(s, 5.0, 3.3, 4.2, 2.3, LGREEN)
    accent(s, 5.0, 3.3, 4.2, 0.07, GREEN)
    txt(s,"Dataset Numbers", 5.15, 3.45, 4.0, 0.38, size=14, bold=True, color=GREEN)
    for val,lbl in [("30","lecture videos"),("32.52 h","total audio"),
                     ("46,525","sentences (spaCy)"),("419","YouTube chapter boundaries"),
                     ("904","reviewed subtopic labels")]:
        txt(s, f"  {val:>8}  {lbl}", 5.15, 3.9+[0,.38,.76,1.14,1.52][
            [("30","lecture videos"),("32.52 h","total audio"),
             ("46,525","sentences (spaCy)"),("419","YouTube chapter boundaries"),
             ("904","reviewed subtopic labels")].index((val,lbl))
        ], 4.0, 0.35, size=12, color=DGREY)
    # domain breakdown
    rect(s, 9.4, 3.3, 3.7, 2.3, WHITE)
    accent(s, 9.4, 3.3, 3.7, 0.07, LBLUE)
    txt(s,"5 Domains", 9.55, 3.45, 3.4, 0.38, size=14, bold=True, color=NAVY)
    domains = [("Biology",6,GREEN),("Computer Science",7,LBLUE),
               ("Mathematics",4,PURPLE),("Philosophy",6,AMBER),("Physics",7,RED)]
    for i,(dom,n,c) in enumerate(domains):
        oy = 3.9 + i*0.34
        bw2 = (n/7)*2.5
        rect(s, 9.55, oy, bw2, 0.28, c)
        txt(s, f"{dom} ({n})", 9.57, oy+0.03, 2.4, 0.24, size=10, color=WHITE, bold=True)


def s04_preprocessing(p):
    s = blank(p)
    hdr(s,"Full Preprocessing Pipeline","From raw video to sentence-level feature vectors — every tool used")
    # 6 stages in two rows
    stages = [
        ("Video\nDownload","Tool: yt-dlp\nOutput: mp4 file\n+ metadata JSON","Audio + video\nsaved locally", BLUE),
        ("Audio\nExtraction","Tool: ffmpeg\nFormat: 16kHz WAV\nMono channel", "Raw waveform\nfor ASR", BLUE),
        ("ASR\nTranscription","Tool: faster-whisper\nModel: large-v3 (GPU)\nVAD pre-filtering","Whisper segments\nwith timestamps", RED),
        ("Sentence\nSplitting","Tool: spaCy en_core_web_sm\nMerge <3-token segments\nProportional timestamps","N=847±512 sentences\nper video", TEAL),
        ("Text\nEmbedding","Models: BGE-large, E5-large,\nMPNet, MiniLM (CPU ok)\nOutput: (N,D) float32","Sentence vectors\nfor boundary scoring", BLUE),
        ("Shot\nDetection","Tool: TransNetV2\n25fps boundary probs\nThreshold: 0.5","Visual scene\nchange signals", LBLUE),
        ("Slide\nOCR","Tool: PaddleOCR\nSample: 0.5 fps keyframes\nDeduplicate ±3-sent window","Slide text\nper sentence", LBLUE),
        ("Prosody\nFeatures","Pause duration: Whisper gaps≥300ms\nPitch (F0): PYIN via librosa\nZ-scored per video","2 scalar features\nper sentence", PURPLE),
    ]
    for i,(label,tools,output,color) in enumerate(stages):
        row,col = i//4, i%4
        gx = 0.25 + col*3.25
        gy = 1.25 + row*2.85
        rect(s, gx, gy, 3.0, 2.6, WHITE)
        accent(s, gx, gy, 3.0, 0.07, color)
        txt(s, f"{i+1}. {label}", gx+0.1, gy+0.12, 2.8, 0.48,
            size=13, bold=True, color=NAVY)
        rect(s, gx+0.1, gy+0.62, 2.8, 1.3, LGREY)
        txt(s, tools, gx+0.15, gy+0.67, 2.7, 1.2, size=11, color=DGREY)
        rect(s, gx+0.1, gy+1.98, 2.8, 0.52, color)
        txt(s, output, gx+0.15, gy+2.02, 2.7, 0.44,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if col < 3 and i < 7:
            txt(s, "→", gx+3.02, gy+1.1, 0.24, 0.38, size=14, bold=True,
                color=LBLUE, align=PP_ALIGN.CENTER)


def s05_annotation(p):
    s = blank(p)
    hdr(s,"Annotation: Two-Level Hierarchy","Chapters (creator-provided) + Subtopics (LLM-assisted human review)")
    # Level 1
    rect(s, 0.3, 1.25, 12.7, 2.0, WHITE)
    accent(s, 0.3, 1.25, 12.7, 0.07, NAVY)
    txt(s,"LEVEL 1: Chapter Boundaries (Ground Truth)", 0.5, 1.38, 12, 0.42,
        size=15, bold=True, color=NAVY)
    txt(s,"Source: YouTube chapter metadata created by the video author\n"
        "Advantages: reproducible, publicly verifiable, reflects real viewer navigation behavior\n"
        "Limitation: editorial choices (may be early/late, inconsistent granularity) — treated as navigation metadata, not pedagogical gold standard",
        0.5, 1.85, 12.3, 1.2, size=13, color=DGREY)
    # Level 2 flow
    rect(s, 0.3, 3.45, 12.7, 2.85, WHITE)
    accent(s, 0.3, 3.45, 12.7, 0.07, TEAL)
    txt(s,"LEVEL 2: Subtopic Boundaries (LLM-Assisted Human Review)", 0.5, 3.58, 12, 0.42,
        size=15, bold=True, color=NAVY)
    pass_steps = [
        ("Chapter\nboundaries fixed\nfrom YouTube","Given to annotator\nas context"),
        ("Pass 1: LLM\n(Llama-3.1-8B\nOllama local)","Generates draft\nsubtopic boundaries\nfor each chapter"),
        ("Pass 2: Human\nReviewer","Uses Streamlit tool\n(scripts/annotate.py)\nto correct drafts"),
        ("IAA Check\n(Cohen's kappa)","Second LLM pass\nwithout seeing first\nκ=0.43 subtopic"),
        ("Final Labels\n904 subtopics","Stored in\ndata/gt_hier/\n30 videos"),
    ]
    for i,(label,desc) in enumerate(pass_steps):
        gx = 0.5 + i*2.48
        rect(s, gx, 4.05, 2.2, 1.95, LGREY)
        txt(s, label, gx+0.1, 4.1, 2.0, 0.65, size=11, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, 4.75, 2.0, 1.2, size=10, color=DGREY)
        if i < 4:
            txt(s, "→", gx+2.22, 4.85, 0.28, 0.38, size=16, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)
    # Hierarchy visual
    rect(s, 0.3, 6.48, 12.7, 0.8, NAVY)
    txt(s,"Chapter:   [Introduction]              [Methodology]              [Results]",
        0.5, 6.52, 12.3, 0.35, size=13, bold=True, color=GOLD)
    txt(s,"Subtopic:  [Intro][Problem][Setup]      [N1][N2][N3][N4]           [Results][Discussion]",
        0.5, 6.85, 12.3, 0.35, size=13, color=RGBColor(0xA0,0xC4,0xF0))


def s06_metrics(p):
    s = blank(p)
    hdr(s,"Evaluation Metrics: Pk & WindowDiff","The industry standard for topic segmentation — used since 1999, still current")
    # Pk
    rect(s, 0.3, 1.25, 6.2, 5.95, WHITE)
    accent(s, 0.3, 1.25, 6.2, 0.07, NAVY)
    txt(s,"Pk  (Beeferman et al., 1999)  —  lower is better",
        0.5, 1.38, 5.9, 0.42, size=15, bold=True, color=NAVY)
    txt(s,"ALGORITHM:", 0.5, 1.88, 5.9, 0.3, size=12, bold=True, color=TEAL)
    txt(s,"1. Compute k = floor(N / 2K)  where N=sentences, K=reference boundaries\n"
        "2. Slide a window of width k across ALL sentence pairs (i, i+k)\n"
        "3. For each pair: check if reference says 'same segment'\n"
        "4. Check if prediction agrees\n"
        "5. Count disagreements ÷ total pairs = Pk score",
        0.5, 2.2, 5.9, 1.6, size=12, color=DGREY)
    # Visual sentence diagram
    txt(s,"VISUAL EXAMPLE  (k=3, 10 sentences, boundary at S4 and S7):",
        0.5, 3.85, 5.9, 0.32, size=11, bold=True, color=NAVY)
    for i in range(10):
        bx = 0.5 + i*0.55
        is_b = i in [3,6]
        bc = TEAL if is_b else RGBColor(0xC0,0xD8,0xF0)
        rect(s, bx, 4.22, 0.48, 0.35, bc)
        txt(s, f"S{i+1}", bx+0.07, 4.26, 0.36, 0.26,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s,"     ↑ bnd            ↑ bnd", 0.5, 4.6, 5.9, 0.28, size=11, color=TEAL)
    rect(s, 0.5, 4.65, 1.65, 0.25, LBLUE)  # window indicator
    txt(s,"window (k=3)", 0.55, 4.68, 1.55, 0.2, size=9, color=WHITE)
    txt(s,"INTERPRETATION:", 0.5, 5.0, 5.9, 0.3, size=12, bold=True, color=TEAL)
    for val,what,col in [("0.0","Perfect","green"),("0.35-0.40","Good for 30-video lecture corpus","blue"),
                          ("0.5","Random baseline","orange"),("1.0","Worst possible","red")]:
        fc = {"green":GREEN,"blue":NAVY,"orange":AMBER,"red":RED}[col]
        txt(s, f"  {val:<8} = {what}", 0.5, 5.3+{"green":0,"blue":0.3,"orange":0.6,"red":0.9}[col],
            5.9, 0.3, size=12, color=fc)
    # WD
    rect(s, 6.8, 1.25, 6.2, 5.95, WHITE)
    accent(s, 6.8, 1.25, 6.2, 0.07, LBLUE)
    txt(s,"WindowDiff  (Pevzner & Hearst, 2002)  —  lower is better",
        7.0, 1.38, 5.9, 0.42, size=15, bold=True, color=NAVY)
    txt(s,"KEY DIFFERENCE FROM Pk:", 7.0, 1.88, 5.9, 0.3, size=12, bold=True, color=TEAL)
    txt(s,"Instead of asking 'same segment yes/no?',\n"
        "WindowDiff counts BOUNDARIES INSIDE each window:\n\n"
        "  Error = |count_in_reference - count_in_prediction|\n\n"
        "This catches OVER-SEGMENTATION that Pk sometimes misses.\n"
        "If your system places 3 boundaries where reference has 1,\n"
        "WD penalises all 3 insertions; Pk might only see 1 error.\n\n"
        "WD ≥ Pk almost always.\n"
        "If WD >> Pk: your system is over-segmenting.",
        7.0, 2.2, 5.9, 3.0, size=12, color=DGREY)
    txt(s,"WHY WE USE BOTH:", 7.0, 5.3, 5.9, 0.3, size=12, bold=True, color=TEAL)
    txt(s,"Pk: measures window-consistency; standard since 1999\n"
        "WD: penalises over-prediction more heavily\n"
        "Together they give a balanced picture of boundary quality.\n"
        "We also add: Boundary Similarity (BS) and Tolerance-F1",
        7.0, 5.62, 5.9, 1.35, size=12, color=DGREY)


def s07_all_methods(p):
    s = blank(p)
    hdr(s,"Everything We Tried: 11 Method Families","Full experimental landscape — classical, neural, novel, multimodal, supervised probes")
    table(s,
        headers=["Method","Type","How it works (1 line)","Pk","Verdict"],
        rows=[
            ["TextTiling","Classical","Vocab overlap sliding window; boundary = local vocab minimum","0.605","❌ Worst"],
            ["C99","Classical","Rank-transform cosine matrix; diagonal block clustering","0.422","❌ Weak"],
            ["CosineSeg","Neural","SBERT/BGE vectors; cosine drop depth-score valleys","0.490","❌ Weak"],
            ["KMeansSeg","Neural","K-means cluster embeddings; label-change = boundary","0.617","❌ Over-segs"],
            ["BertSeg","Neural","Like CosineSeg, larger window (w=5); approximates SegBot","0.489","❌ Weak"],
            ["BGE Divisive","Unsupervised","Recursive max-cosine-dissimilarity split (TreeSeg-style)","0.388","✅ Best baseline"],
            ["TwoStage+Prosody","N1+N2","Two-stage predictor fused with pause/pitch signals","0.433","⚠ F1 up, Pk worse"],
            ["Hierarchical","N1+N3","Two-stage + nested chapter/subtopic output","0.412","⚠ Useful structure"],
            ["Cross-Model (E5)","N2 cross-model","BGE+E5 both score; keep only agreed boundaries","0.371","✅ Best Pk/WD"],
            ["CLIP+Text fusion","N2+visual","Entropy-weighted text+visual boundary scoring","0.374","✅ Near best"],
            ["Balanced Selector","Supervised","ExtraTrees picks best method per video (LOO-CV)","0.359","✅ Best mean Pk"],
        ],
        x=0.25, y=1.22, w=12.8, hs=12, rs=11)
    txt(s,"Pk closer to 0 = better. BGE-Divisive (0.388) is the reference baseline. Everything above it is an improvement.",
        0.25, 6.78, 12.8, 0.42, size=12, bold=True, color=NAVY)


def s08_how_methods_work(p):
    s = blank(p)
    hdr(s,"How the Core Methods Work","Under the hood: from sentence vectors to boundary predictions")
    cards_data = [
        ("BGE Divisive\n(Best Baseline)", NAVY,
         "1. Embed all N sentences with BAAI/bge-large-en-v1.5 (1024-dim)\n"
         "2. Build cosine dissimilarity matrix between all pairs\n"
         "3. Recursively split: find position of maximum dissimilarity\n"
         "4. Each split creates two segments; repeat until n_segments reached\n"
         "5. Segment boundaries = split positions\n\n"
         "Why it works: BGE captures long-range semantic shift;\n"
         "recursive splitting naturally finds the biggest topic jumps first."),
        ("Cross-Model Conservative\n(Best Global)", TEAL,
         "1. Run BGE-large independently: produces boundary set B1\n"
         "2. Run E5-large independently: produces boundary set B2\n"
         "3. Keep ONLY boundaries where both models agree: B* = B1 ∩ B2\n"
         "4. Apply minimum-length filter (min 11 sentences/segment)\n\n"
         "Why it works: Reduces false positives dramatically.\n"
         "If only one model sees a boundary, it's probably noise.\n"
         "Conservative = fewer boundaries, better Pk/WD."),
        ("CLIP Visual Fusion\n(Promising Signal)", PURPLE,
         "1. Extract video keyframes at 1 fps\n"
         "2. Embed each frame with CLIP ViT-B/32 (512-dim)\n"
         "3. Compute cosine dissimilarity between consecutive frames\n"
         "4. This gives a 'visual boundary score' per sentence\n"
         "5. Fuse with text score via entropy weighting (N2)\n\n"
         "Why it partially works: Slide transitions ARE editorial decisions\n"
         "comparable to chapter granularity.\n"
         "Alone: Pk=0.396. Fused with text: Pk=0.374."),
    ]
    for i,(title,color,desc) in enumerate(cards_data):
        gx = 0.3 + i*4.35
        rect(s, gx, 1.25, 4.1, 5.95, WHITE)
        accent(s, gx, 1.25, 4.1, 0.07, color)
        txt(s, title, gx+0.1, 1.38, 3.9, 0.55, size=14, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, 1.97, 3.9, 5.0, size=12, color=DGREY)


def s09_novel_n1_n2(p):
    s = blank(p)
    hdr(s,"Novel Methods N1 & N2","Two-Stage Predictor + Reliability-Weighted Fusion")
    # N1
    rect(s, 0.3, 1.25, 6.2, 5.95, WHITE)
    accent(s, 0.3, 1.25, 6.2, 0.07, GREEN)
    txt(s,"N1: Two-Stage Boundary Predictor", 0.5, 1.38, 5.9, 0.42, size=14, bold=True, color=NAVY)
    txt(s,"PROBLEM: single threshold can't balance recall vs precision",
        0.5, 1.88, 5.9, 0.3, size=12, bold=True, color=RED)
    # Stage boxes
    for i,(stg,thresh,goal,bg) in enumerate([
        ("STAGE 1\nBroad Pass","Threshold: mean − 0.5σ\n(LOOSE — keeps ~50% of positions)","Goal: HIGH RECALL\nAlmost all real boundaries included\nMany false positives — that's OK",LBLUE),
        ("STAGE 2\nRefinement","Threshold: mean + 1.2σ\n(TIGHT — only strong evidence)","Goal: HIGH PRECISION\nFilter candidates by fused signal\nFalse positives removed",GREEN),
    ]):
        gy = 2.25 + i*1.85
        rect(s, 0.5, gy, 5.9, 1.7, bg)
        txt(s, stg, 0.65, gy+0.08, 1.7, 0.55, size=12, bold=True, color=WHITE)
        txt(s, thresh, 2.4, gy+0.08, 3.9, 0.55, size=11, color=WHITE)
        txt(s, goal, 0.65, gy+0.7, 5.6, 0.88, size=11, color=WHITE)
        if i == 0:
            txt(s, "↓ pass candidates", 0.5, 4.0, 5.9, 0.3,
                size=11, italic=True, color=LBLUE, align=PP_ALIGN.CENTER)
    txt(s,"Result: improves over single-pass CosineSeg (0.490 → used in 0.371 pipeline)",
        0.5, 6.05, 5.9, 0.5, size=12, bold=True, color=GREEN)
    # N2
    rect(s, 6.8, 1.25, 6.2, 5.95, WHITE)
    accent(s, 6.8, 1.25, 6.2, 0.07, PURPLE)
    txt(s,"N2: Reliability-Weighted Fusion", 7.0, 1.38, 5.9, 0.42, size=14, bold=True, color=NAVY)
    txt(s,"INPUT: multiple modality boundary score arrays (one per sentence gap)\n"
        "  text_scores   = cosine dissimilarity from BGE/E5 [shape: N-1]\n"
        "  visual_scores = CLIP frame dissimilarity     [shape: N-1]\n"
        "  prosody_scores= pause + pitch signal          [shape: N-1]",
        7.0, 1.88, 5.9, 1.3, size=12, color=DGREY)
    txt(s,"FORMULA:", 7.0, 3.25, 5.9, 0.3, size=12, bold=True, color=PURPLE)
    txt(s,"  H(m) = normalised Shannon entropy of modality m's scores\n"
        "  w(m) = exp(-H(m)) / sum of all exp(-H(m'))\n"
        "  fused = sum( w(m) * normalise(scores(m)) )",
        7.0, 3.58, 5.9, 1.0, size=12, bold=True, color=NAVY)
    txt(s,"INTUITION:", 7.0, 4.65, 5.9, 0.3, size=12, bold=True, color=PURPLE)
    for mod,H,w,verdict in [
        ("Text (BGE)","LOW entropy","HIGH weight","Sharp peaks at real boundaries"),
        ("Visual (CLIP)","MEDIUM entropy","MEDIUM weight","Slide changes match chapters"),
        ("Prosody","HIGH entropy","LOW weight","Too flat — auto down-weighted"),
    ]:
        oy = 5.0+[0,.48,.96][["Text (BGE)","Visual (CLIP)","Prosody"].index(mod)]
        rect(s, 7.0, oy, 5.9, 0.43, LGREY)
        txt(s, f"{mod}: H={H} → {w} → {verdict}",
            7.1, oy+0.07, 5.7, 0.3, size=11, color=DGREY)


def s10_what_worked_failed(p):
    s = blank(p)
    hdr(s,"What Worked & What Failed","Every approach tested with Pk result and verdict")
    # Pk bar chart (horizontal)
    methods = [
        ("TextTiling",0.605,RED),("KMeansSeg",0.617,RED),("CosineSeg",0.490,RED),
        ("BertSeg",0.489,RED),("C99",0.422,AMBER),("TwoStage+Prosody",0.433,AMBER),
        ("Hierarchical",0.412,AMBER),("BGE-Divisive",0.388,GREEN),
        ("CLIP+Text Fusion",0.374,GREEN),("Cross-Model (E5)",0.371,GREEN),
        ("Balanced Selector",0.359,GREEN),
    ]
    # Draw horizontal bars
    rect(s, 0.25, 1.25, 7.8, 5.95, WHITE)
    accent(s, 0.25, 1.25, 7.8, 0.07, NAVY)
    txt(s,"Pk Results (lower = better; 0=perfect, 0.5=random)",
        0.35, 1.38, 7.6, 0.38, size=13, bold=True, color=NAVY)
    max_pk = 0.65
    for i,(name,pk,color) in enumerate(methods):
        oy = 1.85 + i*0.47
        bw = (pk/max_pk)*5.5
        txt(s, name, 0.35, oy+0.04, 2.1, 0.35, size=11, color=DGREY)
        rect(s, 2.5, oy+0.03, bw, 0.34, color)
        txt(s, f"{pk:.3f}", 2.55+bw, oy+0.07, 0.55, 0.26, size=11, bold=True, color=NAVY)
        if name == "BGE-Divisive":
            rect(s, 2.5+(0.388/max_pk)*5.5-0.01, oy-0.1, 0.02, 5.85, MGREY)
    txt(s,"← baseline", 5.15, 6.92, 1.0, 0.3, size=10, italic=True, color=DGREY)
    # Key failure reasons
    rect(s, 8.3, 1.25, 4.8, 5.95, WHITE)
    accent(s, 8.3, 1.25, 4.8, 0.07, RED)
    txt(s,"Why Methods Failed", 8.5, 1.38, 4.5, 0.38, size=13, bold=True, color=NAVY)
    failures = [
        ("TextTiling\nKMeansSeg","No dense embeddings.\nVocab overlap too sparse\nfor technical lectures."),
        ("Prosody signals\n(Pk=0.417)","Pauses/pitch detect FINE\ntransitions (subtopic level).\nGranularity mismatch."),
        ("BERTopic\n(Pk=0.563)","Topic model finds subtopics,\nnot chapter-level editorial\ndecisions."),
        ("LLM zero-shot\n3 videos (Pk=0.406)","Over-segments. High recall\nbut poor precision without\nfine-tuning."),
        ("Supervised BERT-wiki\n(Pk=0.493)","Out-of-domain. Wikipedia\ntext ≠ spoken lecture\ntranscripts."),
    ]
    oy = 1.88
    for name,reason in failures:
        rect(s, 8.4, oy, 4.6, 1.0, LRED)
        txt(s, name, 8.5, oy+0.06, 2.0, 0.55, size=11, bold=True, color=RED)
        txt(s, reason, 10.55, oy+0.06, 2.35, 0.85, size=10, color=DGREY)
        oy += 1.08


def s11_granularity_finding(p):
    s = blank(p)
    hdr(s,"Key Finding: The Granularity Mismatch","The single most actionable result from this research")
    rect(s, 0.3, 1.25, 12.7, 0.9, LRED)
    txt(s,"Every acoustic/linguistic signal we tested produced WORSE Pk than text embeddings.\n"
        "Cause: these signals detect FINE-GRAINED transitions (subtopic level), not the COARSE editorial decisions YouTube chapters represent.",
        0.5, 1.32, 12.3, 0.8, size=13, bold=True, color=RED)
    # Signal spectrum chart
    signals = [
        ("BERTopic topic shifts",0.563,RED,"Subtopic-level clusters"),
        ("Discourse markers",0.462,RED,"Sentence-level connectives"),
        ("GPT-2 perplexity",0.418,AMBER,"Sentence-level surprise"),
        ("Pause/pitch (prosody)",0.417,AMBER,"Speaker breath resets"),
        ("BGE text embedding",0.388,GREEN,"Paragraph semantic shift"),
        ("CLIP visual (solo)",0.396,GREEN,"Slide scene changes"),
        ("CLIP + text fusion",0.374,GREEN,"Best non-text signal"),
        ("Cross-model (best)",0.371,GREEN,"Conservative selection"),
    ]
    rect(s, 0.3, 2.3, 8.5, 4.85, WHITE)
    accent(s, 0.3, 2.3, 8.5, 0.07, NAVY)
    txt(s,"Signal vs Pk (lower = better):", 0.5, 2.42, 8.2, 0.35, size=13, bold=True, color=NAVY)
    for i,(name,pk,color,why) in enumerate(signals):
        oy = 2.85 + i*0.48
        bw = (pk/0.65)*5.5
        rect(s, 0.5, oy, bw, 0.38, color)
        txt(s, f"{name}: Pk={pk:.3f}", 0.55, oy+0.06, bw-0.1, 0.28,
            size=11, bold=True, color=WHITE)
        txt(s, why, 0.52+bw+0.05, oy+0.06, 3.0, 0.28, size=10, color=DGREY)
    # Explanation
    rect(s, 9.0, 2.3, 4.1, 4.85, WHITE)
    accent(s, 9.0, 2.3, 4.1, 0.07, GREEN)
    txt(s,"Why CLIP Is Different", 9.2, 2.42, 3.8, 0.38, size=13, bold=True, color=NAVY)
    txt(s,"YouTube chapters = deliberate\neditor choices at CHAPTER level:\n\n"
        "  'Introduction to Eigenvalues'\n  'Characteristic Equation'\n  'Physics Applications'\n\n"
        "Acoustic signals fire at EVERY:\n  pause, pitch reset, discourse marker\n  (×5-10 more often than chapters)\n\n"
        "Slide transitions in lectures ARE\neditorial chapter-level decisions.\n"
        "CLIP sees the same editorial intent.\n\n"
        "Practical implication:\n"
        "Add CLIP + text, suppress prosody.",
        9.2, 2.85, 3.8, 4.2, size=12, color=DGREY)


def s12_selector_deep(p):
    s = blank(p)
    hdr(s,"The Method Selector: Full Breakdown","Supervised meta-model that picks the best segmentation method per video")
    # What is it
    rect(s, 0.3, 1.25, 5.5, 5.95, WHITE)
    accent(s, 0.3, 1.25, 5.5, 0.07, PURPLE)
    txt(s,"What Is It?", 0.5, 1.38, 5.2, 0.4, size=14, bold=True, color=NAVY)
    txt(s,"80+ method variants exist (embedding x method x hyperparameter).\n"
        "No single variant wins on all 30 videos.\n"
        "The selector LEARNS which variant to use per video.",
        0.5, 1.85, 5.2, 0.95, size=12, color=DGREY)
    txt(s,"MODEL: ExtraTreesRegressor\n(scikit-learn ensemble, 500 trees)", 0.5, 2.85, 5.2, 0.55,
        size=13, bold=True, color=PURPLE)
    txt(s,"INPUTS FED TO THE SELECTOR:", 0.5, 3.5, 5.2, 0.35, size=12, bold=True, color=TEAL)
    inputs = [
        ("Video features","Domain (Biology/CS/Math/Phil/Physics) — one-hot\nVideo length (minutes)\nSentence count\nMean cosine similarity (BGE-large)\nStd of consecutive similarities"),
        ("Training-fold stats","For each of the 80+ candidate methods:\n  Mean Pk across the 29 training videos\n  (= how reliable is this method generally?)"),
    ]
    oy = 3.9
    for title,desc in inputs:
        rect(s, 0.5, oy, 5.0, 1.05+("Video" in title)*0.2, LGREY)
        txt(s, title, 0.6, oy+0.06, 4.8, 0.35, size=12, bold=True, color=NAVY)
        txt(s, desc, 0.6, oy+0.42, 4.8, 0.78+("Video" in title)*0.2,
            size=11, color=DGREY)
        oy += 1.2+("Video" in title)*0.2
    txt(s,"OUTPUT: predicted Pk per method → pick lowest", 0.5, 6.68, 5.2, 0.4,
        size=12, bold=True, color=GREEN)
    # LOO diagram
    rect(s, 6.1, 1.25, 6.9, 2.8, WHITE)
    accent(s, 6.1, 1.25, 6.9, 0.07, LBLUE)
    txt(s,"Leave-One-Video-Out Training (LOO-CV)", 6.3, 1.38, 6.5, 0.4,
        size=13, bold=True, color=NAVY)
    txt(s,"For each of 30 videos:", 6.3, 1.85, 6.5, 0.3, size=12, color=DGREY)
    for i in range(6):
        is_test = (i == 2)
        bc = RED if is_test else RGBColor(0xC0,0xD8,0xF0)
        lbl = "TEST" if is_test else f"Train"
        rect(s, 6.3+i*1.1, 2.22, 0.95, 0.45, bc)
        txt(s, lbl, 6.33+i*1.1, 2.28, 0.88, 0.32, size=11, bold=True,
            color=WHITE if is_test else NAVY, align=PP_ALIGN.CENTER)
    txt(s,"...", 12.7, 2.28, 0.4, 0.32, size=14, color=DGREY)
    txt(s,"Train on 29 videos, predict for 1 held-out. Repeat 30 times.",
        6.3, 2.75, 6.5, 0.32, size=11, italic=True, color=DGREY)
    # What it chose per domain
    rect(s, 6.1, 4.25, 6.9, 2.95, WHITE)
    accent(s, 6.1, 4.25, 6.9, 0.07, TEAL)
    txt(s,"What the Selector Chose (per domain)", 6.3, 4.38, 6.5, 0.4,
        size=13, bold=True, color=NAVY)
    domain_choices = [
        ("Physics (7 videos)","Multimodal-grid (CLIP+text) for 5/7 videos","Biggest Pk gains",GREEN),
        ("CS (7 videos)","Cross-E5 conservative most reliable","Consistent",LBLUE),
        ("Biology (6 videos)","Cross-E5 + some over-switching","Mixed",AMBER),
        ("Philosophy (6 videos)","Selector improves Pk/WD","Positive",GREEN),
        ("Math (4 videos)","❌ FAILS — only 3 training examples per fold","Selector hurts",RED),
    ]
    oy = 4.82
    for dom,choice,verdict,color in domain_choices:
        rect(s, 6.3, oy, 6.5, 0.43, LGREY)
        txt(s, dom, 6.4, oy+0.06, 2.0, 0.3, size=11, bold=True, color=NAVY)
        txt(s, choice, 8.45, oy+0.06, 3.0, 0.3, size=10, color=DGREY)
        rect(s, 11.5, oy, 0.8, 0.43, color)
        txt(s, verdict, 11.55, oy+0.06, 0.7, 0.3, size=9, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        oy += 0.48


def s13_selector_verdict(p):
    s = blank(p)
    hdr(s,"Selector Results & Verdict","How much did the selector improve things?")
    # Three operating points comparison
    ops = [
        ("BGE-Divisive\n(Baseline)","0.388","0.396","0.129","0.088",NAVY,"Reference"),
        ("Cross-Model\nConservative","0.371","0.376","0.036","0.024",TEAL,"Best Pk/WD (sig.)"),
        ("Balanced\nSelector","0.359","0.374","0.076","0.089",GREEN,"Best mean point"),
    ]
    rect(s, 0.3, 1.25, 12.7, 1.85, WHITE)
    accent(s, 0.3, 1.25, 12.7, 0.07, NAVY)
    for ci,h in enumerate(["Method","Pk ↓ (better)","WD ↓ (better)","BS ↑ (better)","F1@2 ↑ (better)","vs Baseline"]):
        cx = 0.35+ci*2.12
        txt(s, h, cx, 1.38, 2.0, 0.38, size=13, bold=True, color=NAVY)
    for ri,(name,pk,wd,bs,f1,color,note) in enumerate(ops):
        oy = 1.72 + ri*0.42
        bg = LGREY if ri%2==0 else WHITE
        rect(s, 0.3, oy, 12.7, 0.4, bg)
        for ci,v in enumerate([name,pk,wd,bs,f1,note]):
            cx = 0.35+ci*2.12
            fc = color if ci in [1,2,3,4] else NAVY
            txt(s, v, cx, oy+0.06, 2.0, 0.28, size=12, bold=(ci==0), color=fc)
    # Significance table
    rect(s, 0.3, 3.35, 7.5, 3.85, WHITE)
    accent(s, 0.3, 3.35, 7.5, 0.07, GREEN)
    txt(s,"Statistical Significance (paired Wilcoxon, α=0.05 with Holm correction)",
        0.5, 3.48, 7.2, 0.42, size=13, bold=True, color=NAVY)
    sig_rows = [
        ("Cross-model vs BGE-divisive","Pk","✅ Significant p<0.05"),
        ("Cross-model vs BGE-divisive","WD","✅ Significant p<0.05"),
        ("Selector vs BGE-divisive","Pk","✅ Significant p<0.05"),
        ("Selector vs BGE-divisive","BS + F1","✅ Significant p<0.05"),
        ("Selector vs Cross-model","Pk/WD","❌ NOT significant"),
        ("Selector vs Cross-model","BS + F1","✅ Significant p<0.05"),
    ]
    oy = 3.95
    for comp,metric,result in sig_rows:
        bg = LGREEN if "✅" in result else LRED
        rect(s, 0.5, oy, 7.1, 0.42, bg)
        txt(s, comp, 0.6, oy+0.07, 3.2, 0.28, size=11, color=DGREY)
        txt(s, metric, 3.85, oy+0.07, 0.7, 0.28, size=11, bold=True, color=NAVY)
        fc = GREEN if "✅" in result else RED
        txt(s, result, 4.6, oy+0.07, 2.9, 0.28, size=11, bold=True, color=fc)
        oy += 0.46
    # Verdict box
    rect(s, 8.1, 3.35, 5.1, 3.85, NAVY)
    txt(s,"Final Verdict on Selector", 8.3, 3.48, 4.8, 0.42,
        size=14, bold=True, color=WHITE)
    txt(s,"The selector improves mean Pk and significantly\n"
        "improves BS and F1 vs cross-model.\n\n"
        "Pk/WD gain over cross-model alone is NOT\nstatistically significant.\n\n"
        "Interpretation:\n"
        "  The selector is a stronger OPERATING POINT\n"
        "  experiment — not a conclusive replacement.\n\n"
        "Leave-domain-out: Pk=0.4012 (worse than baseline)\n"
        "→ NOT a universal deployment model.\n"
        "→ A LOW-RESOURCE BENCHMARK operating point.",
        8.3, 3.98, 4.8, 3.15, size=12, color=RGBColor(0xC0,0xD8,0xF0))


def s14_paper_comparison(p):
    s = blank(p)
    hdr(s,"Comparison to Prior Work","Supervised vs Unsupervised — direct vs indirect — what metrics they used")
    table(s,
        headers=["Paper","Videos","Supervised?","Metric","Their Result","Our Result","Comparison"],
        rows=[
            ["TreeSeg\n(Gklezakos 2024)","21 (TinyRec)\n+ICSI+AMI","No\n(unsupervised)","Pk/WD","Pk=0.367 (TinyRec)","Pk=0.359 (selector)","⚠ INDIRECT\nDiff dataset"],
            ["AVLectures\n(Singh 2022)","2350+","Self-supervised","No Pk/WD","Outperforms\nbaselines","N/A","❌ Not comparable\nDiff metrics"],
            ["Freisinger 2023","34+96 lectures","LoRA fine-tuned","F1 only","F1=67.3","N/A","❌ Not comparable\nNo Pk reported"],
            ["MiniSeg/YTSEG\n(Retkowski 2024)","19,299","Supervised","Pk/WD/F1","Pk=28.73","Pk=35.9 (ours)","❌ INDIRECT\n643x more data"],
            ["Chapter-Gen\n(Cao 2022)","9,631","Supervised","AP/Recall","AP=43.3","N/A","❌ Different task\nGeneration not seg"],
            ["Chapter-Llama\n(Ventura 2025)","10k/8.1k","Trained LLM","F1/IoU","F1=45.3","N/A","❌ Different task\nDiff metrics"],
            ["TextTiling (Hearst)","Our data","No","Pk","0.605 (ours)","0.371 (cross-model)","✅ DIRECT\nOurs better"],
            ["C99 (Choi 2000)","Our data","No","Pk","0.422 (ours)","0.371 (cross-model)","✅ DIRECT\nOurs better"],
        ],
        x=0.25, y=1.22, w=12.8,
        hs=11, rs=10)
    rect(s, 0.25, 6.55, 12.8, 0.72, NAVY)
    txt(s,"DIRECT comparisons (same dataset, same metric): our method beats all classical/neural baselines.\n"
        "INDIRECT comparisons (different datasets): we match or approach TreeSeg on Pk, using similar scale. Large supervised systems are not comparable.",
        0.4, 6.6, 12.5, 0.65, size=12, bold=True, color=WHITE)


def s15_oracle(p):
    s = blank(p)
    hdr(s,"Oracle Gap Analysis: Where Is the Ceiling?",
        "Shown to locate the bottleneck — NOT a performance claim. Read carefully.")
    # Warning box
    rect(s, 0.3, 1.25, 12.7, 0.95, LAMBER)
    txt(s,"WARNING: The oracle results are NOT deployable. They use ground truth at test time to pick the best boundaries/methods.\n"
        "We show them ONLY to answer: 'Where exactly is the remaining gap, and what needs to be solved next?'",
        0.5, 1.3, 12.3, 0.85, size=13, bold=True, color=AMBER)
    # Operating points table
    rect(s, 0.3, 2.35, 7.2, 4.85, WHITE)
    accent(s, 0.3, 2.35, 7.2, 0.07, NAVY)
    txt(s,"Operating Points (Pk)", 0.5, 2.48, 7.0, 0.4, size=14, bold=True, color=NAVY)
    ops = [
        ("BGE-Divisive (baseline)","0.388","✅ Deployable","Our stable starting point"),
        ("Cross-model conservative","0.371","✅ Deployable","Statistically significant"),
        ("Balanced selector","0.359","✅ Deployable (LOO)","Best mean local point"),
        ("TreeSeg/TinyRec (indicative)","0.367","✅ Deployable","Different dataset"),
        ("Per-video oracle","0.298","❌ NOT deployable","Needs gold at test time"),
        ("Candidate oracle τ=2","0.017","❌ NOT deployable","Perfect selector needed"),
    ]
    for i,(name,pk,dep,note) in enumerate(ops):
        oy = 2.95 + i*0.55
        bg = LGREY if i%2==0 else WHITE
        rect(s, 0.5, oy, 6.8, 0.5, bg)
        txt(s, name, 0.6, oy+0.08, 2.4, 0.34, size=12, color=NAVY, bold=(i<4))
        txt(s, pk, 3.05, oy+0.08, 0.7, 0.34, size=14, bold=True,
            color=GREEN if float(pk)<0.37 else (AMBER if float(pk)<0.4 else RED))
        fc = GREEN if "✅" in dep else RED
        txt(s, dep, 3.8, oy+0.08, 1.7, 0.34, size=11, color=fc)
        txt(s, note, 5.55, oy+0.08, 1.65, 0.34, size=10, color=DGREY)
    # Gap explanation
    rect(s, 7.8, 2.35, 5.2, 4.85, WHITE)
    accent(s, 7.8, 2.35, 5.2, 0.07, LBLUE)
    txt(s,"What the Gap Tells Us", 8.0, 2.48, 4.9, 0.4, size=14, bold=True, color=NAVY)
    txt(s,"Gap: deployable best (Pk=0.359) to candidate oracle (Pk=0.017)\n"
        "ΔPk ≈ 0.34 — this entire gap is in SELECTION, not generation.\n\n"
        "Proof: 96.8% recall already exists in the candidate pool.\n"
        "The right boundaries are ALREADY being found —\n"
        "we just cannot always pick the correct ones.\n\n"
        "What this means for future work:\n"
        "  → NOT 'build a better segmenter'\n"
        "  → YES 'build a better candidate RANKER'\n\n"
        "A supervised ranker trained on:\n"
        "  - Local cosine contrast per candidate\n"
        "  - Cross-model agreement score\n"
        "  - Pause duration at candidate position\n"
        "  - Slide-change signal (CLIP)\n"
        "  ...with 50-100 labelled lectures\n"
        "could plausibly close most of this gap.",
        8.0, 2.95, 4.9, 4.2, size=12, color=DGREY)


def s16_webapp(p):
    s = blank(p)
    hdr(s,"Real-World Application: LecSeg Web App","Not just research — a working tool for any YouTube lecture")
    rect(s, 0.3, 1.25, 5.8, 5.95, WHITE)
    accent(s, 0.3, 1.25, 5.8, 0.07, TEAL)
    txt(s,"What the Web App Does", 0.5, 1.38, 5.5, 0.42, size=15, bold=True, color=NAVY)
    txt(s,"Built with Streamlit (scripts/demo.py)\n"
        "Run: streamlit run scripts/demo.py\n\n"
        "Two modes:\n\n"
        "1. SEGMENT ANY YOUTUBE VIDEO\n"
        "   Paste any YouTube URL or video ID\n"
        "   → Auto-downloads transcript (yt-dlp captions)\n"
        "   → Sentence split + BGE/E5/MPNet/MiniLM embed\n"
        "   → Divisive segmentation\n"
        "   → Llama-3.1-8B generates chapter titles\n"
        "   → Exports YouTube timestamp format\n\n"
        "2. BENCHMARK MODE (LecSeg-30)\n"
        "   Evaluate any method on the 30-video benchmark\n"
        "   Shows Pk/WD/BS/F1 with reference chapters",
        0.5, 1.88, 5.5, 4.2, size=12, color=DGREY)
    txt(s,"Try it: streamlit run scripts/demo.py",
        0.5, 6.05, 5.5, 0.45, size=13, bold=True, color=TEAL)
    # App pipeline visual
    rect(s, 6.4, 1.25, 6.7, 5.95, WHITE)
    accent(s, 6.4, 1.25, 6.7, 0.07, LBLUE)
    txt(s,"App Pipeline (Arbitrary Video Mode)", 6.6, 1.38, 6.4, 0.42,
        size=14, bold=True, color=NAVY)
    app_steps = [
        ("User pastes\nYouTube URL","Any lecture video;\nno prior indexing needed"),
        ("Transcript\nAcquisition","yt-dlp captions (fast)\nFallback: Whisper (slow)"),
        ("Sentence\nSplitting","spaCy; cache per video_id\nAvoids recomputation"),
        ("Embedding\n+ Segmentation","Select model in sidebar\n(BGE/E5/MPNet/MiniLM)"),
        ("LLM Titling","Llama-3.1-8B (auto-starts)\nFallback: noun-phrase extraction"),
        ("Output","Timestamp list\nYouTube format export\nSemantic similarity chart"),
    ]
    for i,(lbl,desc) in enumerate(app_steps):
        gy = 1.85 + i*0.82
        rect(s, 6.6, gy, 6.3, 0.75, LGREY if i%2==0 else WHITE)
        txt(s, f"{i+1}. {lbl}", 6.7, gy+0.06, 2.2, 0.42, size=12, bold=True, color=NAVY)
        txt(s, desc, 8.95, gy+0.06, 3.8, 0.55, size=11, color=DGREY)
        if i < 5:
            txt(s, "↓", 9.7, gy+0.72, 0.3, 0.25, size=12, color=LBLUE,
                align=PP_ALIGN.CENTER)


def s17_research_justification(p):
    s = blank(p)
    hdr(s,"Justifying the Research Title Word-by-Word",
        "LecSeg-30: A Reproducible Low-Resource Benchmark and Diagnostic Study for Lecture-Video Topic Segmentation")
    rect(s, 0.3, 1.25, 12.7, 0.95, NAVY)
    txt(s,'LecSeg-30  ·  A Reproducible  ·  Low-Resource  ·  Benchmark  ·  and Diagnostic Study  ·  for Lecture-Video Topic Segmentation',
        0.5, 1.32, 12.3, 0.82, size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    justifications = [
        ("LecSeg-30","30 academic YouTube lecture videos\nacross 5 domains, 32.52 hours total.\nThe '30' is in the name because the\nbenchmark is permanently defined.",
         BLUE),
        ("Reproducible","All code open (src/lecseg/).\nAll data IDs public (YouTube URLs).\n185 automated tests.\nEvery table has a generation script.\nAny researcher can rerun.",
         TEAL),
        ("Low-Resource","Only 30 videos (vs 10k-800k for\nlarge supervised systems).\nNo GPU required at inference.\nAll models run on CPU.\nFully offline with Ollama.",
         GREEN),
        ("Benchmark","Fixed 30-video set with:\n- Creator-provided chapter GT\n- Two-level hierarchical labels\n- IAA with Cohen's kappa\n- 5 evaluation metrics",
         PURPLE),
        ("Diagnostic Study","We don't just report Pk.\nWe show WHAT FAILS and WHY:\n- Granularity mismatch finding\n- Oracle gap locates bottleneck\n- Domain failure analysis\n- Ablation of every signal",
         AMBER),
        ("Lecture-Video\nTopic Seg","Not general video chaptering.\nNot meeting segmentation.\nSpecifically academic YouTube\nlectures with Pk/WD evaluation.",
         RED),
    ]
    for i,(word,proof,color) in enumerate(justifications):
        row,col = i//3, i%3
        gx = 0.3 + col*4.35
        gy = 2.38 + row*2.45
        rect(s, gx, gy, 4.1, 2.35, WHITE)
        accent(s, gx, gy, 4.1, 0.07, color)
        rect(s, gx, gy, 4.1, 0.48, color)
        txt(s, word, gx+0.1, gy+0.06, 3.9, 0.38, size=14, bold=True, color=WHITE)
        txt(s, proof, gx+0.1, gy+0.55, 3.9, 1.72, size=11, color=DGREY)


def s18_thesis_quality(p):
    s = blank(p)
    hdr(s,"Thesis Standards & Quality Checklist","Is everything up to academic standard?")
    checks = [
        ("✅","Clickable references","hyperref package with hidelinks+breaklinks; all \\citep{} are clickable in PDF"),
        ("✅","In-text citations","38 citations across 4 chapters; all figures/tables cited"),
        ("✅","Title match","Titlepage: 'LecSeg-30: A Reproducible Low-Resource Benchmark and Diagnostic Study'"),
        ("✅","LaTeX structure","book class, 12pt, A4; TOC, LOF, LOT, bibliography (plainnat); appendices"),
        ("✅","6 chapters","Intro, Literature, Methodology, Results, Conclusion, Future Work"),
        ("✅","3 appendices","Dataset details, Hyperparameters, Extra results"),
        ("✅","15 LaTeX tables","All \\input'd from thesis/tables/; generated from data"),
        ("✅","7 PDF figures","All in thesis/figures/; referenced with \\includegraphics"),
        ("✅","Statistical tests","Bootstrap 95% CIs + paired Wilcoxon with Holm correction"),
        ("✅","IAA reported","Cohen's kappa chapter=0.535, subtopic=0.426; limitation acknowledged"),
        ("✅","Limitations section","Silver labels, corpus size, strict F1, selector domain-failure all stated"),
        ("✅","Threats to validity","Separate section in Ch5: reference validity, annotation independence, small dataset"),
        ("✅","GPU transcription","vast.ai RTX 5090 documented in Appendix B (hardware section)"),
        ("✅","Preprocessing","Full pipeline: Whisper, spaCy, TransNetV2, PaddleOCR, librosa all documented in Ch3"),
        ("⚠️","Pipeline figure","Currently a text placeholder in Ch3 (fbox) — no actual PDF figure for pipeline diagram"),
    ]
    table(s,
        headers=["Status","Item","Evidence"],
        rows=[[c,n,e] for c,n,e in checks],
        x=0.3, y=1.22, w=12.7,
        hs=13, rs=11)


def s19_limitations(p):
    s = blank(p)
    hdr(s,"Limitations & Honest Assessment","What we cannot claim — and what the examiner will ask about")
    limits = [
        ("30 videos is small",RED,
         "Math domain has only 4 → selector fails there.",
         "We acknowledge this. 30 matches prior low-resource work.\nBootstrap CIs quantify uncertainty. Selector limitation explicitly reported."),
        ("Silver chapter labels",AMBER,
         "YouTube chapters = creator editorial choices, not expert pedagogical labels.",
         "Treated as navigation metadata only. Human validation listed as future work.\nLimitation stated in threats-to-validity section."),
        ("Low strict F1",AMBER,
         "Best Pk/WD method is very conservative → F1@2 = 0.024.",
         "Pk/WD and F1 measure different things. Selector recovers F1 to 0.089.\nBoth reported. The tradeoff is documented in the modern-metrics section."),
        ("Selector not universal",AMBER,
         "Leave-domain-out: Pk=0.4012 (worse than baseline).",
         "Explicitly acknowledged. The selector is a local benchmark operating point,\nnot a production deployment model. Clearly stated in Ch3 and Ch5."),
        ("LLM refinement unverified",LBLUE,
         "Llama-3.1-8B boundary shifting not proven to improve Pk significantly.",
         "Reported as diagnostic only. Titling works offline. Not overclaimed.\nFuture work: 70B model or task-specific fine-tuning."),
    ]
    for i,(issue,color,problem,response) in enumerate(limits):
        gy = 1.28 + i*1.18
        rect(s, 0.3, gy, 12.7, 1.1, WHITE)
        accent(s, 0.3, gy, 12.7, 0.06, color)
        txt(s, f"L{i+1}: {issue}", 0.45, gy+0.1, 3.5, 0.38, size=13, bold=True, color=NAVY)
        txt(s, problem, 4.0, gy+0.1, 4.2, 0.38, size=12, color=RED)
        rect(s, 8.3, gy+0.06, 4.6, 0.95, LGREEN)
        txt(s, "How we address it:", 8.4, gy+0.1, 4.4, 0.28, size=10, bold=True, color=GREEN)
        txt(s, response, 8.4, gy+0.38, 4.4, 0.65, size=10, color=DGREY)


def s20_results_summary(p):
    s = blank(p)
    hdr(s,"Main Results Summary","The definitive numbers from LecSeg-30")
    # Large Pk comparison chart
    rect(s, 0.3, 1.25, 8.2, 5.95, WHITE)
    accent(s, 0.3, 1.25, 8.2, 0.07, NAVY)
    txt(s,"Pk Results Across All Methods (lower = better)",
        0.5, 1.38, 8.0, 0.42, size=14, bold=True, color=NAVY)
    methods = [
        ("KMeansSeg",0.617,RED),("TextTiling",0.605,RED),
        ("CosineSeg",0.490,RED),("BertSeg",0.489,RED),("C99",0.422,AMBER),
        ("TwoStage+Prosody",0.433,AMBER),("Hierarchical",0.412,AMBER),
        ("BGE-Divisive",0.388,NAVY),("CLIP+Text",0.374,GREEN),
        ("Cross-Model",0.371,GREEN),("Selector",0.359,GREEN),
    ]
    max_v = 0.7
    for i,(name,pk,color) in enumerate(methods):
        oy = 1.87 + i*0.42
        bw = (pk/max_v)*6.3
        txt(s, name, 0.45, oy+0.05, 2.1, 0.32, size=11, color=DGREY)
        rect(s, 2.6, oy+0.04, bw, 0.33, color)
        txt(s, f"{pk:.3f}", 2.65+bw, oy+0.07, 0.6, 0.24, size=11, bold=True, color=NAVY)
    # baseline line
    base_x = 2.6 + (0.388/max_v)*6.3
    rect(s, base_x, 1.87, 0.015, 4.65, AMBER)
    txt(s,"baseline\n0.388", base_x-0.3, 6.35, 0.8, 0.55, size=9, italic=True, color=AMBER)
    # Full metrics table
    rect(s, 8.75, 1.25, 4.35, 5.95, WHITE)
    accent(s, 8.75, 1.25, 4.35, 0.07, TEAL)
    txt(s,"All Metrics", 8.95, 1.38, 4.1, 0.42, size=14, bold=True, color=NAVY)
    table(s,
        headers=["Method","Pk","WD","BS","F1"],
        rows=[
            ["BGE-Div","0.388","0.396","0.129","0.088"],
            ["Cross-Model","0.371","0.376","0.036","0.024"],
            ["Selector","0.359","0.374","0.076","0.089"],
        ],
        x=8.85, y=1.88, w=4.1, hs=11, rs=11)
    txt(s,"Pk/WD: lower is better\nBS/F1: higher is better\n\nSignificant improvements:\nCross-model: Pk+WD vs baseline\nSelector: BS+F1 vs cross-model\n\nN=30 videos\n95% bootstrap CIs\nPaired Wilcoxon+Holm",
        8.95, 3.1, 4.1, 3.95, size=12, color=DGREY)


def s21_contributions(p):
    s = blank(p)
    hdr(s,"What This Research Contributes","LecSeg-30: A Reproducible Low-Resource Benchmark and Diagnostic Study")
    rect(s, 0.3, 1.25, 12.7, 5.95, WHITE)
    accent(s, 0.3, 1.25, 12.7, 0.08, GOLD)
    contribs = [
        (TEAL,"C1: LecSeg-30 Benchmark",
         "30 YouTube lectures · 5 domains · 32.52 h · 419 creator chapters · 904 subtopics\n"
         "Two-level hierarchy with IAA (κ=0.54 chapter, κ=0.43 subtopic) · Public YouTube URLs"),
        (BLUE,"C2: Unified Evaluation Suite",
         "Project-local implementations of Pk, WD, BS, tolerance-F1, H-WD\n"
         "Bootstrap 95% CIs · Paired Wilcoxon+Holm · 185 automated tests · All scripts included"),
        (GREEN,"C3: Statistically Supported Improvement",
         "Cross-model conservative: Pk=0.371 — significantly better than all baselines (p<0.05)\n"
         "Selector: best mean Pk=0.359, significant F1/BS improvement over cross-model"),
        (PURPLE,"C4: Granularity Mismatch Finding",
         "Acoustic/linguistic signals (prosody, BERTopic, perplexity, discourse) all over-segment\n"
         "CLIP visual is the exception: matches chapter-level editorial granularity\n"
         "Practical roadmap: add CLIP, suppress prosody for chapter-level segmentation"),
        (AMBER,"C5: Oracle Gap Diagnosis",
         "96.8% recall in candidate pool — bottleneck is SELECTION, not generation\n"
         "ΔPk≈0.34 gap is entirely in ranking. Next step: supervised boundary ranker with 50-100 lectures"),
        (LBLUE,"C6: Reproducible Pipeline + Web App",
         "End-to-end from YouTube URL to titled chapters · CPU-friendly · Fully offline\n"
         "Streamlit web demo · vast.ai GPU transcription documented · All hyperparameters logged"),
    ]
    for i,(color,title,desc) in enumerate(contribs):
        row,col = i//2, i%2
        gx = 0.45 + col*6.35
        gy = 1.45 + row*1.88
        rect(s, gx, gy, 6.1, 1.78, LGREY)
        accent(s, gx, gy, 6.1, 0.07, color)
        txt(s, title, gx+0.1, gy+0.12, 5.9, 0.42, size=13, bold=True, color=NAVY)
        txt(s, desc, gx+0.1, gy+0.57, 5.9, 1.15, size=11, color=DGREY)


def s22_final_verdict(p):
    s = blank(p)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    for cx,cy,r,c in [(11.8,1.0,2.0,BLUE),(12.5,5.5,1.5,BLUE)]:
        circle(s, cx, cy, r, c)
    txt(s,"Final Verdict", 0.8, 0.5, 12, 0.8, size=42, bold=True, color=WHITE)
    accent(s, 0.8, 1.35, 10, 0.08, TEAL)
    verdicts = [
        (GREEN,
         "What we proved:",
         "Conservative cross-model scoring (Pk=0.371) is significantly better than\n"
         "every classical and neural baseline on LecSeg-30 (p<0.05, Wilcoxon+Holm).\n"
         "The balanced selector (Pk=0.359) significantly improves BS and F1.\n"
         "A systematic granularity mismatch explains why acoustic signals fail."),
        (AMBER,
         "What we honestly admit:",
         "We do not beat large supervised systems on external benchmarks.\n"
         "The selector is not domain-general (Math fails with 4 training examples).\n"
         "LLM boundary refinement is not proven to help Pk.\n"
         "These limitations are stated explicitly and are defensible."),
        (LBLUE,
         "What this enables:",
         "LecSeg-30 is a fixed, inspectable benchmark for future lecture-segmentation work.\n"
         "The oracle gap (Pk=0.017) identifies the next concrete engineering target.\n"
         "The web app makes the research immediately usable on any YouTube lecture.\n"
         "All code, data, scripts are reproducible by any researcher."),
    ]
    for i,(color,label,text) in enumerate(verdicts):
        gy = 1.52 + i*1.82
        rect(s, 0.8, gy, 11.8, 1.72, BLUE)
        accent(s, 0.8, gy, 11.8, 0.07, color)
        txt(s, label, 1.0, gy+0.12, 3.0, 0.42, size=14, bold=True, color=color)
        txt(s, text, 1.0, gy+0.58, 11.4, 1.1, size=13, color=RGBColor(0xC8,0xDC,0xF0))
    txt(s,"LecSeg-30 does not make lecture segmentation solved. It makes it MEASURABLE and REPRODUCIBLE.",
        0.8, 7.1, 11.8, 0.38, size=14, italic=True, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════
def main():
    p = new_prs()
    s01_title(p)
    s02_roadmap(p)
    s03_data_collection(p)
    s04_preprocessing(p)
    s05_annotation(p)
    s06_metrics(p)
    s07_all_methods(p)
    s08_how_methods_work(p)
    s09_novel_n1_n2(p)
    s10_what_worked_failed(p)
    s11_granularity_finding(p)
    s12_selector_deep(p)
    s13_selector_verdict(p)
    s14_paper_comparison(p)
    s15_oracle(p)
    s16_webapp(p)
    s17_research_justification(p)
    s18_thesis_quality(p)
    s19_limitations(p)
    s20_results_summary(p)
    s21_contributions(p)
    s22_final_verdict(p)
    out = "thesis/LECSEG_Defense_Slides_v3.pptx"
    p.save(out)
    print(f"Done. {len(p.slides)} slides saved -> {out}")

if __name__ == "__main__":
    main()
