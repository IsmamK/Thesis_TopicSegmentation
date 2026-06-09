"""
Generate comprehensive LECSEG defense slides (v2).
Run: python scripts/generate_defense_slides_v2.py
Output: thesis/LECSEG_Defense_Slides_v2.pptx
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x23, 0x40)
BLUE   = RGBColor(0x1A, 0x4A, 0x8A)
LBLUE  = RGBColor(0x2D, 0x6A, 0xBF)
TEAL   = RGBColor(0x00, 0x96, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF0, 0xF4, 0xF8)
DGREY  = RGBColor(0x37, 0x47, 0x51)
GREEN  = RGBColor(0x1B, 0x87, 0x3B)
RED    = RGBColor(0xC6, 0x28, 0x28)
AMBER  = RGBColor(0xE6, 0x5C, 0x00)
GOLD   = RGBColor(0xF5, 0xA6, 0x23)

W, H = Inches(13.33), Inches(7.5)   # 16:9


# ── Helpers ────────────────────────────────────────────────────────────────
def prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p


def blank(prs_obj):
    layout = prs_obj.slide_layouts[6]  # blank
    return prs_obj.slides.add_slide(layout)


def fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def box(slide, x, y, w, h, color: RGBColor | None = None, radius=0):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    if color:
        fill(s, color)
    else:
        s.fill.background()
    return s


def txt(slide, text, x, y, w, h,
        size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    """Dark navy header bar."""
    box(slide, 0, 0, 13.33, 1.15, NAVY)
    txt(slide, title, 0.3, 0.08, 10, 0.75, size=30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.8, 12, 0.4, size=14, color=RGBColor(0xA0, 0xC4, 0xF0))


def accent_bar(slide, x, y, w, h=0.06, color=TEAL):
    box(slide, x, y, w, h, color)


def bullet_block(slide, items, x, y, w, h, title=None,
                 bg=LGREY, title_color=NAVY, text_color=DGREY,
                 size=15, title_size=16):
    box(slide, x, y, w, h, bg)
    oy = y + 0.12
    if title:
        txt(slide, title, x+0.15, oy, w-0.3, 0.35,
            size=title_size, bold=True, color=title_color)
        oy += 0.38
    for item in items:
        txt(slide, f"▸  {item}", x+0.15, oy, w-0.3, 0.35,
            size=size, color=text_color)
        oy += 0.32
    return oy


def step_box(slide, x, y, w, h, number, label, desc, bg=BLUE):
    box(slide, x, y, w, h, bg)
    txt(slide, str(number), x+0.1, y+0.05, 0.5, 0.45,
        size=28, bold=True, color=WHITE)
    txt(slide, label, x+0.65, y+0.05, w-0.75, 0.38,
        size=16, bold=True, color=WHITE)
    txt(slide, desc, x+0.15, y+0.52, w-0.25, 0.55,
        size=12, color=RGBColor(0xC0, 0xD8, 0xF0))


def compare_table(slide, headers, rows, x, y, w,
                  hdr_bg=NAVY, hdr_color=WHITE,
                  row_bgs=(LGREY, WHITE), text_color=DGREY,
                  hdr_size=13, row_size=12):
    cols = len(headers)
    col_w = w / cols
    row_h = 0.38
    # header
    for ci, h_text in enumerate(headers):
        box(slide, x + ci*col_w, y, col_w, row_h, hdr_bg)
        txt(slide, h_text, x+ci*col_w+0.05, y+0.06, col_w-0.1, row_h-0.1,
            size=hdr_size, bold=True, color=hdr_color, align=PP_ALIGN.CENTER)
    # rows
    for ri, row in enumerate(rows):
        bg = row_bgs[ri % 2]
        for ci, cell in enumerate(row):
            box(slide, x+ci*col_w, y+(ri+1)*row_h, col_w, row_h, bg)
            color = text_color
            if cell.startswith("✅"):
                color = GREEN
            elif cell.startswith("⚠"):
                color = AMBER
            elif cell.startswith("❌"):
                color = RED
            txt(slide, cell, x+ci*col_w+0.04, y+(ri+1)*row_h+0.05,
                col_w-0.08, row_h-0.08, size=row_size, color=color,
                align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2):
    """Simple right-pointing arrow box."""
    mx = (x1+x2)/2 - 0.15
    box(slide, mx, (y1+y2)/2-0.08, 0.3, 0.16, LBLUE)
    txt(slide, "→", mx-0.02, (y1+y2)/2-0.12, 0.35, 0.28,
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(p):
    s = blank(p)
    # Full background
    box(s, 0, 0, 13.33, 7.5, NAVY)
    # Decorative circles
    for cx, cy, r, a in [(11.5,1.2,2.2,0.06),(12.8,5.8,1.8,0.04),(0.5,6.5,1.0,0.03)]:
        b = s.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x1A,0x4A,0x8A)
        b.line.fill.background()
    accent_bar(s, 0.8, 2.85, 6, 0.08, TEAL)
    txt(s, "LECSEG", 0.8, 1.2, 12, 1.0, size=58, bold=True, color=WHITE)
    txt(s, "AI-Powered Lecture Topic Segmentation", 0.8, 2.2, 12, 0.7,
        size=26, color=RGBColor(0xA0,0xC4,0xF0))
    txt(s, "Hierarchical Benchmark · 30 Lectures · 5 Domains · Reproducible Pipeline",
        0.8, 3.05, 12, 0.6, size=17, color=RGBColor(0x80,0xB4,0xE0), italic=True)
    txt(s, "Pre-Thesis Research  ·  T2520718", 0.8, 6.6, 8, 0.5,
        size=14, color=RGBColor(0x60,0x90,0xC0))
    txt(s, "BRAC University  ·  2026", 9.5, 6.6, 4, 0.5,
        size=14, color=RGBColor(0x60,0x90,0xC0), align=PP_ALIGN.RIGHT)


def slide_problem(p):
    s = blank(p)
    header(s, "The Problem", "Why can't students navigate lecture videos?")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Left: scenario
    box(s, 0.3, 1.35, 5.8, 4.2, WHITE)
    accent_bar(s, 0.3, 1.35, 5.8, 0.07, RED)
    txt(s, "😩  The Student's Nightmare", 0.45, 1.5, 5.5, 0.5,
        size=17, bold=True, color=NAVY)
    scenario = [
        "3-hour linear algebra lecture on YouTube",
        "Student needs the eigenvalue derivation",
        "No chapter markers → must scrub blindly",
        "Manual chapters are rare, inconsistent,",
        "   or missing entirely for older content",
        "",
        "At scale: millions of educational videos",
        "with ZERO structured navigation",
    ]
    oy = 2.1
    for line in scenario:
        txt(s, line, 0.5, oy, 5.4, 0.35, size=14, color=DGREY)
        oy += 0.33

    # Right: scale
    box(s, 6.4, 1.35, 6.6, 4.2, WHITE)
    accent_bar(s, 6.4, 1.35, 6.6, 0.07, LBLUE)
    txt(s, "📊  The Scale", 6.55, 1.5, 6.3, 0.5, size=17, bold=True, color=NAVY)
    stats = [
        ("Coursera + edX", "Tens of thousands of lectures"),
        ("YouTube Education", "Millions of videos"),
        ("MIT OpenCourseWare", "2,400+ courses"),
        ("Avg lecture duration", "~50 minutes"),
        ("Without chapters", "> 80% of content"),
    ]
    oy = 2.1
    for label, val in stats:
        box(s, 6.55, oy, 2.8, 0.32, RGBColor(0xE8,0xEF,0xF8))
        box(s, 9.4, oy, 3.5, 0.32, RGBColor(0xD0,0xE4,0xF8))
        txt(s, label, 6.6, oy+0.04, 2.7, 0.26, size=13, color=DGREY)
        txt(s, val, 9.45, oy+0.04, 3.4, 0.26, size=13, bold=True, color=NAVY)
        oy += 0.37

    # Bottom: solution
    box(s, 0.3, 5.7, 12.7, 1.15, NAVY)
    txt(s, "✅  Solution: Automatically predict chapter timestamps from the video's own signal",
        0.5, 5.82, 12.3, 0.55, size=16, bold=True, color=WHITE)
    txt(s, "No manual labelling at inference time · Works on any lecture video",
        0.5, 6.25, 12.3, 0.35, size=13, color=RGBColor(0xA0,0xC4,0xF0))


def slide_gap(p):
    s = blank(p)
    header(s, "Research Gap", "What was missing before LECSEG?")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    gaps = [
        ("🎓", "Lecture-specific benchmark",
         "No compact, reproducible benchmark combined authentic academic lectures\n"
         "with hierarchical annotations AND Pk/WD evaluation together"),
        ("📊", "Two-level hierarchy",
         "Prior datasets had chapters OR subtopics — never a nested\n"
         "chapter ⊂ subtopic structure with enforced nesting constraints"),
        ("🔬", "Honest low-resource baseline",
         "Large systems use 10k–800k videos. No work showed what's achievable\n"
         "in the 30-video regime with statistical significance tests"),
        ("🤖", "Locally reproducible pipeline",
         "Supervised/LLM methods require training data or proprietary APIs.\n"
         "LECSEG runs entirely offline: Whisper + sentence-transformers + Ollama"),
    ]

    for i, (icon, title, desc) in enumerate(gaps):
        row, col = i // 2, i % 2
        gx = 0.3 + col * 6.5
        gy = 1.4 + row * 2.8
        box(s, gx, gy, 6.2, 2.6, WHITE)
        accent_bar(s, gx, gy, 6.2, 0.07, [TEAL,LBLUE,GREEN,AMBER][i])
        txt(s, icon, gx+0.15, gy+0.2, 0.7, 0.6, size=30, color=NAVY)
        txt(s, title, gx+0.9, gy+0.2, 5.1, 0.45, size=16, bold=True, color=NAVY)
        txt(s, desc, gx+0.15, gy+0.75, 5.9, 1.5, size=13, color=DGREY)


def slide_dataset(p):
    s = blank(p)
    header(s, "LECSEG-30 Dataset", "30 academic YouTube lectures · 5 domains · 32.52 hours")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Domain distribution
    box(s, 0.3, 1.35, 5.5, 5.9, WHITE)
    accent_bar(s, 0.3, 1.35, 5.5, 0.07, TEAL)
    txt(s, "Domain Distribution", 0.45, 1.5, 5.2, 0.4, size=16, bold=True, color=NAVY)

    domains = [
        ("🧬 Biology",     6, "green",  0.30),
        ("💻 Computer Sci",7, "blue",   0.35),
        ("📐 Mathematics", 4, "purple", 0.20),
        ("📚 Philosophy",  6, "orange", 0.30),
        ("⚛️ Physics",     7, "red",    0.35),
    ]
    bar_colors = [GREEN, LBLUE, RGBColor(0x7B,0x1F,0xA2), AMBER, RED]
    oy = 2.05
    for (label, n, _, frac), bc in zip(domains, bar_colors):
        txt(s, label, 0.5, oy, 2.2, 0.32, size=13, color=DGREY)
        bw = frac * 3.0
        b = s.shapes.add_shape(1, Inches(2.8), Inches(oy+0.04),
                               Inches(bw), Inches(0.24))
        b.fill.solid(); b.fill.fore_color.rgb = bc; b.line.fill.background()
        txt(s, str(n), 2.8+bw+0.05, oy, 0.4, 0.32, size=13, bold=True, color=NAVY)
        oy += 0.6

    txt(s, "Selection criteria:", 0.5, 4.65, 5.0, 0.32, size=13, bold=True, color=NAVY)
    for cr in ["≥ 15 min duration", "Clean English audio",
                "≥ 4 creator YouTube chapters", "Open licence"]:
        txt(s, f"  ✓  {cr}", 0.5, oy, 4.8, 0.3, size=12, color=DGREY)
        oy += 0.3

    # Stats
    box(s, 6.1, 1.35, 6.9, 5.9, WHITE)
    accent_bar(s, 6.1, 1.35, 6.9, 0.07, LBLUE)
    txt(s, "Key Numbers", 6.25, 1.5, 6.5, 0.4, size=16, bold=True, color=NAVY)

    stats = [
        ("30", "lecture videos"),
        ("32.52 h", "total audio"),
        ("46,525", "sentences (spaCy)"),
        ("419", "chapter boundaries (YouTube)"),
        ("904", "subtopic labels (LLM-assisted)"),
        ("κ = 0.54", "chapter IAA (Cohen's κ)"),
        ("κ = 0.43", "subtopic IAA"),
        ("F₁ = 1.00", "chapter boundary recovery"),
    ]
    oy2 = 2.05
    for val, label in stats:
        box(s, 6.25, oy2, 2.2, 0.5, RGBColor(0xE8,0xEF,0xF8))
        txt(s, val, 6.3, oy2+0.06, 2.1, 0.38, size=18, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER)
        txt(s, label, 8.5, oy2+0.1, 4.3, 0.38, size=13, color=DGREY)
        oy2 += 0.62


def slide_annotation(p):
    s = blank(p)
    header(s, "Annotation Protocol", "Two-level hierarchy: Chapters → Subtopics")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Level 1
    box(s, 0.3, 1.4, 12.7, 2.3, WHITE)
    accent_bar(s, 0.3, 1.4, 12.7, 0.07, NAVY)
    txt(s, "Level 1 — Chapter Boundaries  (creator-provided YouTube chapters)",
        0.5, 1.55, 12, 0.42, size=16, bold=True, color=NAVY)
    txt(s, "Taken directly from YouTube metadata → reproducible · publicly available · real navigation behavior\n"
        "⚠️  Not unquestionable pedagogical truth: creators may place chapters early/late or omit fine boundaries",
        0.5, 2.0, 12.3, 0.9, size=13, color=DGREY)
    txt(s, "418 chapter boundaries across 30 videos  ·  mean 14.0 per video",
        0.5, 2.85, 12.3, 0.38, size=14, bold=True, color=LBLUE)

    # Level 2
    box(s, 0.3, 3.85, 12.7, 2.4, WHITE)
    accent_bar(s, 0.3, 3.85, 12.7, 0.07, TEAL)
    txt(s, "Level 2 — Subtopic Boundaries  (LLM-assisted human review)",
        0.5, 4.0, 12, 0.42, size=16, bold=True, color=NAVY)

    steps = [
        ("Pass 1", "Llama-3.1-8B (Ollama)", "Generates draft subtopic\nsegmentation per chapter"),
        ("Pass 2", "Human Reviewer", "Checks, corrects, and\napproves each draft boundary"),
        ("Result", "904 subtopic labels", "Nesting enforced:\nchapter ⊂ subtopic always"),
    ]
    for i, (step, who, what) in enumerate(steps):
        gx = 0.5 + i * 4.2
        box(s, gx, 4.5, 3.8, 1.5, LGREY)
        txt(s, step, gx+0.1, 4.55, 1.5, 0.35, size=13, bold=True, color=TEAL)
        txt(s, who, gx+0.1, 4.88, 3.5, 0.38, size=14, bold=True, color=NAVY)
        txt(s, what, gx+0.1, 5.28, 3.5, 0.65, size=12, color=DGREY)
        if i < 2:
            txt(s, "→", gx+3.85, 5.0, 0.4, 0.4, size=24, bold=True, color=LBLUE,
                align=PP_ALIGN.CENTER)

    txt(s, "Limitation: both passes use same LLM → agreement may be inflated vs. fully independent humans",
        0.5, 6.3, 12.5, 0.4, size=12, italic=True, color=AMBER)


def slide_metrics(p):
    s = blank(p)
    header(s, "How We Measure: Pk & WindowDiff", "The standard metrics for topic segmentation — lower is better")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Pk explanation
    box(s, 0.3, 1.35, 6.2, 5.9, WHITE)
    accent_bar(s, 0.3, 1.35, 6.2, 0.07, NAVY)
    txt(s, "Pk (Beeferman 1999)", 0.5, 1.5, 5.8, 0.42, size=16, bold=True, color=NAVY)

    txt(s, "HOW IT WORKS:", 0.5, 2.0, 5.8, 0.3, size=12, bold=True, color=TEAL)
    txt(s,
        "Slide a window of width k across the transcript.\n"
        "At each position i, check:\n"
        "  • Do sentences i and i+k belong to the SAME segment\n"
        "    in the reference (ground truth)?\n"
        "  • Does the PREDICTED segmentation agree?\n\n"
        "If they disagree → count as an error.\n"
        "Pk = errors / total window positions",
        0.5, 2.35, 5.8, 2.3, size=13, color=DGREY)

    # Visual diagram for Pk
    txt(s, "Example (k=2 window, 8 sentences):", 0.5, 4.65, 5.8, 0.3,
        size=12, bold=True, color=NAVY)
    # Sentence boxes
    for i in range(8):
        bx = 0.5 + i * 0.68
        bc = TEAL if i == 3 else (LBLUE if i == 6 else RGBColor(0xC0,0xD8,0xF0))
        b = s.shapes.add_shape(1, Inches(bx), Inches(5.0),
                               Inches(0.6), Inches(0.32))
        b.fill.solid(); b.fill.fore_color.rgb = bc; b.line.fill.background()
        txt(s, f"S{i+1}", bx+0.08, 5.02, 0.5, 0.26, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "↑ boundary", 2.55, 5.35, 1.2, 0.28, size=11, color=TEAL, bold=True)
    txt(s, "↑ boundary", 4.59, 5.35, 1.2, 0.28, size=11, color=TEAL, bold=True)

    txt(s, "INTERPRETATION:", 0.5, 5.7, 5.8, 0.3, size=12, bold=True, color=TEAL)
    txt(s, "0 = perfect  ·  0.5 = random  ·  Lower is ALWAYS better\n"
        "Typical good result on 30-video lecture corpus: Pk ≈ 0.35–0.40",
        0.5, 6.0, 5.8, 0.6, size=12, color=DGREY)

    # WindowDiff
    box(s, 6.8, 1.35, 6.2, 5.9, WHITE)
    accent_bar(s, 6.8, 1.35, 6.2, 0.07, LBLUE)
    txt(s, "WindowDiff (Pevzner & Hearst 2002)", 7.0, 1.5, 5.8, 0.42,
        size=16, bold=True, color=NAVY)

    txt(s, "HOW IT WORKS:", 7.0, 2.0, 5.8, 0.3, size=12, bold=True, color=TEAL)
    txt(s,
        "Same sliding window, BUT instead of asking\n"
        "'same segment or not?', it counts:\n"
        "  • Number of boundaries INSIDE the window\n"
        "    in the reference\n"
        "  • Number of boundaries INSIDE the window\n"
        "    in the prediction\n\n"
        "WD error = windows where these counts differ\n\n"
        "Key difference from Pk:\n"
        "  → Penalises over-segmentation MORE heavily\n"
        "  → Catches insertions of extra boundaries\n"
        "    that Pk sometimes misses",
        7.0, 2.35, 5.8, 3.5, size=13, color=DGREY)

    txt(s, "INTERPRETATION:", 7.0, 5.7, 5.8, 0.3, size=12, bold=True, color=TEAL)
    txt(s, "0 = perfect  ·  Lower is ALWAYS better\n"
        "WD ≥ Pk almost always. If WD >> Pk: system over-segments.",
        7.0, 6.0, 5.8, 0.6, size=12, color=DGREY)


def slide_pipeline(p):
    s = blank(p)
    header(s, "The LECSEG Pipeline", "Two independent parts: unsupervised segmentation + supervised selector")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Part 1 label
    box(s, 0.3, 1.35, 7.8, 0.45, NAVY)
    txt(s, "PART 1 — Unsupervised Segmentation (no training data needed)",
        0.45, 1.38, 7.6, 0.38, size=13, bold=True, color=WHITE)

    # Part 2 label
    box(s, 9.0, 1.35, 4.0, 0.45, TEAL)
    txt(s, "PART 2 — Supervised Selector",
        9.1, 1.38, 3.8, 0.38, size=13, bold=True, color=WHITE)

    steps_1 = [
        ("1", "YouTube Video", "Audio track extracted\nURL → yt-dlp → mp4"),
        ("2", "ASR (Whisper)", "large-v3 on GPU\n17.9× realtime speed"),
        ("3", "Sentence Split", "spaCy sentenciser\n~847 sentences/video"),
        ("4", "Text Embed", "SBERT/BGE/E5/MPNet\n384–1024 dim vectors"),
        ("5", "Boundary Score", "Cosine dissimilarity\nbetween sentence windows"),
        ("6", "Segmentation", "11 methods evaluated\n(see next slides)"),
    ]
    for i, (n, label, desc) in enumerate(steps_1):
        gx = 0.3 + i * 1.32
        gy = 1.95
        box(s, gx, gy, 1.2, 1.7, BLUE)
        txt(s, n, gx+0.42, gy+0.05, 0.4, 0.38, size=22, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, label, gx+0.05, gy+0.45, 1.1, 0.45, size=11, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.05, gy+0.9, 1.1, 0.75, size=10,
            color=RGBColor(0xC0,0xD8,0xF0), align=PP_ALIGN.CENTER)
        if i < 5:
            txt(s, "→", gx+1.22, gy+0.65, 0.12, 0.38, size=14, bold=True,
                color=LBLUE, align=PP_ALIGN.CENTER)

    # Part 2: selector box
    box(s, 9.0, 1.95, 4.0, 1.7, TEAL)
    txt(s, "7", 9.4, 2.0, 0.4, 0.38, size=22, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, "Method Selector", 9.05, 2.45, 3.85, 0.45, size=13, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, "ExtraTreesRegressor\nLeave-one-video-out\nPicks best method/video",
        9.05, 2.9, 3.85, 0.75, size=11, color=RGBColor(0xC0,0xF0,0xE8),
        align=PP_ALIGN.CENTER)

    # Outputs
    box(s, 0.3, 3.85, 12.7, 2.8, WHITE)
    accent_bar(s, 0.3, 3.85, 12.7, 0.07, GOLD)
    txt(s, "Outputs", 0.5, 4.0, 12, 0.38, size=15, bold=True, color=NAVY)

    outputs = [
        ("📍 Chapter timestamps", "List of boundary sentence indices\n→ converted to HH:MM:SS"),
        ("🏷️ Segment titles", "Llama-3.1-8B generates\n3–8 word topic names"),
        ("📊 Hierarchy", "Chapters ⊂ Subtopics\nnesting enforced by construction"),
        ("📈 Metrics", "Pk / WD / BS / F1\nper-video + aggregate"),
    ]
    for i, (label, desc) in enumerate(outputs):
        gx = 0.5 + i * 3.15
        txt(s, label, gx, 4.45, 3.0, 0.38, size=13, bold=True, color=NAVY)
        txt(s, desc, gx, 4.85, 3.0, 0.7, size=12, color=DGREY)


def slide_methods_unsupervised(p):
    s = blank(p)
    header(s, "Unsupervised Methods Evaluated", "We implemented and benchmarked 11 method families — no training labels needed")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    methods = [
        ("TextTiling", "Classical\nBaseline",
         "Lexical cohesion sliding window.\n"
         "Counts shared vocabulary between blocks.\n"
         "Boundary = local vocabulary minimum.",
         "Pk=0.605 ↑ weakest"),
        ("C99", "Classical\nBaseline",
         "Rank-transform cosine similarity matrix.\n"
         "Groups coherent blocks by diagonal dominance.\n"
         "More robust to vocab outliers than TextTiling.",
         "Pk=0.422"),
        ("CosineSeg", "Neural\nBaseline",
         "SBERT sentence vectors.\n"
         "Cosine similarity drops between adjacent windows.\n"
         "Boundary = local cosine valley (depth score).",
         "Pk=0.490"),
        ("BertSeg", "Neural\nBaseline",
         "Same as CosineSeg but with larger comparison\n"
         "window (w=5). Approximates SegBot inference\n"
         "using pre-computed BERT embeddings.",
         "Pk=0.489"),
        ("Divisive (BGE)", "★ Best\nBaseline",
         "BGE-large-en-v1.5 embeddings (1024-dim).\n"
         "Recursively splits on max cosine dissimilarity.\n"
         "Stable, strong — our main comparison point.",
         "Pk=0.388 ✅ best baseline"),
        ("Cross-Model\nConservative", "★ Best\nGlobal",
         "Two models (BGE + E5) independently score.\n"
         "Keep only boundaries BOTH agree on.\n"
         "Conservative → fewer FP, better Pk/WD.",
         "Pk=0.371 ✅ best global"),
    ]

    cols = 3
    for i, (name, tag, desc, result) in enumerate(methods):
        row, col = i // cols, i % cols
        gx = 0.3 + col * 4.3
        gy = 1.4 + row * 2.7
        w, h = 4.0, 2.5
        box(s, gx, gy, w, h, WHITE)
        tag_color = GREEN if "Best" in tag else (LBLUE if "Neural" in tag else AMBER)
        accent_bar(s, gx, gy, w, 0.07, tag_color)
        txt(s, name, gx+0.12, gy+0.12, w-0.6, 0.38, size=15, bold=True, color=NAVY)
        box(s, gx+w-1.3, gy+0.08, 1.2, 0.38, tag_color)
        txt(s, tag, gx+w-1.28, gy+0.1, 1.15, 0.34, size=9, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, desc, gx+0.12, gy+0.55, w-0.2, 1.1, size=12, color=DGREY)
        box(s, gx+0.12, gy+1.72, w-0.2, 0.45, LGREY)
        txt(s, result, gx+0.2, gy+1.78, w-0.3, 0.35, size=12, bold=True,
            color=GREEN if "✅" in result else NAVY)


def slide_novel_n1(p):
    s = blank(p)
    header(s, "N1: Two-Stage Boundary Predictor", "Separates recall (stage 1) from precision (stage 2)")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Left: what problem it solves
    box(s, 0.3, 1.4, 5.0, 5.8, WHITE)
    accent_bar(s, 0.3, 1.4, 5.0, 0.07, RED)
    txt(s, "The Problem with Single-Pass Methods",
        0.45, 1.55, 4.7, 0.42, size=14, bold=True, color=NAVY)
    txt(s,
        "A single threshold must balance:\n\n"
        "  Too HIGH → misses real boundaries (low recall)\n"
        "  Too LOW  → many false positives (hurts Pk)\n\n"
        "Per-video score distributions vary widely.\n"
        "A global threshold fails on diverse content.",
        0.45, 2.05, 4.7, 2.5, size=13, color=DGREY)
    txt(s, "Result: imprecise boundaries, high Pk error",
        0.45, 4.5, 4.7, 0.4, size=13, bold=True, color=RED)

    # Right: the two-stage solution
    box(s, 5.6, 1.4, 7.4, 5.8, WHITE)
    accent_bar(s, 5.6, 1.4, 7.4, 0.07, GREEN)
    txt(s, "Two-Stage Solution (N1)", 5.75, 1.55, 7.1, 0.42,
        size=14, bold=True, color=NAVY)

    # Stage 1 box
    box(s, 5.75, 2.05, 7.1, 1.6, LGREY)
    txt(s, "STAGE 1 — Broad Pass (High Recall)",
        5.9, 2.1, 6.8, 0.38, size=13, bold=True, color=NAVY)
    txt(s,
        "Threshold: mean - 0.5 × std  (LOOSE)\n"
        "→ Retains ~50% of sentence gaps as candidates\n"
        "→ Almost all real boundaries are included\n"
        "→ Many false positives still present (that's OK)",
        5.9, 2.52, 6.8, 1.1, size=12, color=DGREY)

    txt(s, "↓  pass candidates to Stage 2", 5.9, 3.7, 6.8, 0.35,
        size=12, italic=True, color=LBLUE, align=PP_ALIGN.CENTER)

    # Stage 2 box
    box(s, 5.75, 4.1, 7.1, 1.7, RGBColor(0xE8,0xF5,0xE9))
    txt(s, "STAGE 2 — Refinement (High Precision)",
        5.9, 4.15, 6.8, 0.38, size=13, bold=True, color=NAVY)
    txt(s,
        "Fused signal from all modalities (text + visual + prosody)\n"
        "Threshold: mean + 1.2 × std  (TIGHT)\n"
        "→ Keeps only strong, multi-modal evidence boundaries\n"
        "→ Reduces false positives dramatically",
        5.9, 4.57, 6.8, 1.2, size=12, color=DGREY)

    txt(s, "Result: Pk improves over single-pass CosineSeg (0.490 → part of best 0.371)",
        5.75, 5.9, 7.1, 0.4, size=12, bold=True, color=GREEN)


def slide_novel_n2(p):
    s = blank(p)
    header(s, "N2: Reliability-Weighted Fusion", "Weight each modality by how confident its signal is")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Left: the idea
    box(s, 0.3, 1.4, 6.2, 5.8, WHITE)
    accent_bar(s, 0.3, 1.4, 6.2, 0.07, LBLUE)
    txt(s, "The Core Idea", 0.5, 1.55, 5.9, 0.42, size=15, bold=True, color=NAVY)
    txt(s,
        "Different modalities have different reliability:\n\n"
        "  • Text (BGE embeddings): sharp boundary signals\n"
        "    → HIGH confidence, LOW entropy → high weight\n\n"
        "  • Prosody (pauses/pitch): noisier, flatter signal\n"
        "    → LOW confidence, HIGH entropy → low weight\n\n"
        "  • Visual (CLIP keyframes): varies by video type\n"
        "    → Medium weight for slide-based lectures",
        0.5, 2.05, 5.9, 2.8, size=13, color=DGREY)

    txt(s, "The Formula:", 0.5, 4.85, 5.9, 0.38, size=13, bold=True, color=NAVY)
    txt(s, "wₘ = exp(−H(sₘ)) / Σ exp(−H(sₘ′))",
        0.5, 5.22, 5.9, 0.42, size=14, bold=True, color=LBLUE)
    txt(s,
        "H(sₘ) = normalised Shannon entropy of modality m's boundary scores.\n"
        "Low entropy (peaked distribution) → high reliability → high weight.",
        0.5, 5.68, 5.9, 0.8, size=12, color=DGREY)

    # Right: visual of entropy concept
    box(s, 6.8, 1.4, 6.2, 5.8, WHITE)
    accent_bar(s, 6.8, 1.4, 6.2, 0.07, TEAL)
    txt(s, "Entropy as Confidence Proxy", 7.0, 1.55, 5.9, 0.42,
        size=15, bold=True, color=NAVY)

    # High confidence modality visual
    box(s, 7.0, 2.1, 5.8, 1.9, LGREY)
    txt(s, "✅ High Confidence Modality (Text):", 7.1, 2.15, 5.6, 0.32,
        size=13, bold=True, color=GREEN)
    txt(s, "Clear peaks at boundaries → low entropy → high weight",
        7.1, 2.48, 5.6, 0.28, size=12, color=DGREY)
    # bar chart approximation
    vals = [0.1, 0.2, 0.8, 0.15, 0.1, 0.9, 0.12, 0.18]
    for i, v in enumerate(vals):
        bh = v * 0.8
        b = s.shapes.add_shape(1, Inches(7.1+i*0.68), Inches(2.75+0.8-bh),
                               Inches(0.55), Inches(bh))
        b.fill.solid()
        b.fill.fore_color.rgb = GREEN if v > 0.5 else RGBColor(0xC0,0xD8,0xC0)
        b.line.fill.background()

    # Low confidence modality visual
    box(s, 7.0, 4.2, 5.8, 1.9, LGREY)
    txt(s, "⚠️  Low Confidence Modality (Prosody):", 7.1, 4.25, 5.6, 0.32,
        size=13, bold=True, color=AMBER)
    txt(s, "Flat distribution → high entropy → auto down-weighted",
        7.1, 4.58, 5.6, 0.28, size=12, color=DGREY)
    flat_vals = [0.4, 0.5, 0.45, 0.6, 0.38, 0.52, 0.48, 0.42]
    for i, v in enumerate(flat_vals):
        bh = v * 0.8
        b = s.shapes.add_shape(1, Inches(7.1+i*0.68), Inches(4.85+0.8-bh),
                               Inches(0.55), Inches(bh))
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(0xE0,0xC0,0x80)
        b.line.fill.background()

    txt(s, "Key finding: CLIP visual embeddings scored medium–high reliability for slide-based lectures",
        6.8, 6.3, 6.2, 0.45, size=12, italic=True, color=LBLUE)


def slide_novel_n3n4(p):
    s = blank(p)
    header(s, "N3 + N4: Hierarchy & LLM Titling", "Structured output + human-readable chapter names")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # N3
    box(s, 0.3, 1.4, 6.2, 5.8, WHITE)
    accent_bar(s, 0.3, 1.4, 6.2, 0.07, NAVY)
    txt(s, "N3 — Hierarchical Segmenter", 0.5, 1.55, 5.9, 0.42,
        size=15, bold=True, color=NAVY)
    txt(s,
        "Run the predictor TWICE:\n"
        "  1. Fine pass → n_subtopic boundaries\n"
        "  2. Coarse pass → n_chapter boundaries\n\n"
        "Enforce nesting by construction:\n"
        "  B_chapter ⊆ B_subtopic  (always guaranteed)\n\n"
        "No joint optimisation needed — just merge sets.\n\n"
        "Why this matters:\n"
        "  • Chapters for navigation (coarse)\n"
        "  • Subtopics for study (fine-grained)\n"
        "  • Any chapter boundary is also a subtopic boundary",
        0.5, 2.05, 5.9, 3.5, size=13, color=DGREY)

    # Visual hierarchy
    box(s, 0.5, 5.55, 5.8, 1.45, LGREY)
    txt(s, "Chapter:  [───────Intro────────]  [──Methods──]  [──Results──]",
        0.6, 5.6, 5.6, 0.4, size=12, color=NAVY, bold=True)
    txt(s, "Subtopic: [─Intro─][─Prob─][─Setup─]  [─N1─][─N2─]  [─Res─][─Disc─]",
        0.6, 6.0, 5.6, 0.4, size=12, color=TEAL)
    txt(s, "Every │ chapter boundary is also a │ subtopic boundary",
        0.6, 6.4, 5.6, 0.38, size=11, italic=True, color=DGREY)

    # N4
    box(s, 6.8, 1.4, 6.2, 5.8, WHITE)
    accent_bar(s, 6.8, 1.4, 6.2, 0.07, TEAL)
    txt(s, "N4 — Local LLM Refinement & Titling", 7.0, 1.55, 5.9, 0.42,
        size=15, bold=True, color=NAVY)
    txt(s,
        "Uses Llama-3.1-8B via Ollama (fully local, offline)\n\n"
        "Two functions:\n\n"
        "1. BOUNDARY REFINEMENT (diagnostic)\n"
        "   For each predicted boundary b:\n"
        "   → Show LLM sentences [b-2 … b+2]\n"
        "   → Ask: 'which position is a cleaner topic break?'\n"
        "   → Shift boundary by ±2 sentences if better\n\n"
        "2. SEGMENT TITLING (deployed)\n"
        "   For each segment:\n"
        "   → Show LLM first 5 sentences\n"
        "   → Ask: 'give a 3–8 word topic title'\n"
        "   → Fallback: noun-phrase extraction if LLM offline",
        7.0, 2.05, 5.9, 4.0, size=13, color=DGREY)

    txt(s, "⚠️  Boundary refinement: not claimed as main Pk driver.\n"
        "    Verified as diagnostic only — titles work offline.",
        7.0, 6.15, 5.9, 0.75, size=12, italic=True, color=AMBER)


def slide_selector(p):
    s = blank(p)
    header(s, "The Method Selector (Part 2)", "Supervised meta-model: picks the best method for each video")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Left: what it is
    box(s, 0.3, 1.4, 5.8, 5.8, WHITE)
    accent_bar(s, 0.3, 1.4, 5.8, 0.07, TEAL)
    txt(s, "What Is the Selector?", 0.5, 1.55, 5.5, 0.42, size=15, bold=True, color=NAVY)
    txt(s,
        "The unsupervised pipeline produces 80+ method\n"
        "variants (embedding × method × hyperparameter).\n\n"
        "For each video, a different variant may be best.\n"
        "The selector learns WHICH variant to use per video.\n\n"
        "Model: ExtraTreesRegressor (random forest variant)\n\n"
        "Trained on (per leave-one-video-out fold):\n"
        "  • Video features:\n"
        "    - Domain one-hot (Biology/CS/Math/Phil/Physics)\n"
        "    - Video length, sentence count\n"
        "    - Mean/std of embedding similarity scores\n"
        "  • Method statistics from training fold:\n"
        "    - Mean Pk of each candidate method\n"
        "      on the 29 training videos\n\n"
        "Predicts: expected Pk for each method on the new video\n"
        "Selects: method with lowest predicted Pk",
        0.5, 2.05, 5.5, 4.5, size=13, color=DGREY)

    # Right: LOO-CV explanation
    box(s, 6.4, 1.4, 6.6, 5.8, WHITE)
    accent_bar(s, 6.4, 1.4, 6.6, 0.07, LBLUE)
    txt(s, "Leave-One-Video-Out Cross-Validation",
        6.6, 1.55, 6.2, 0.42, size=14, bold=True, color=NAVY)

    # LOO diagram
    for i in range(5):
        is_held = (i == 2)
        bc = RED if is_held else RGBColor(0xC0,0xD8,0xF0)
        b = s.shapes.add_shape(1, Inches(6.6+i*1.25), Inches(2.1),
                               Inches(1.1), Inches(0.45))
        b.fill.solid(); b.fill.fore_color.rgb = bc; b.line.fill.background()
        txt(s, "Test" if is_held else f"Train", 6.65+i*1.25, 2.18, 1.0, 0.28,
            size=12, bold=True, color=WHITE if is_held else NAVY,
            align=PP_ALIGN.CENTER)
    txt(s, "Repeat for ALL 30 videos (each is held out once)",
        6.6, 2.62, 6.2, 0.32, size=12, italic=True, color=DGREY)

    txt(s, "What the selector found:", 6.6, 3.05, 6.2, 0.35, size=13,
        bold=True, color=NAVY)
    findings = [
        ("Physics",    "Multimodal-grid variants (CLIP+text) win"),
        ("CS",         "Cross-E5 conservative is most reliable"),
        ("Biology",    "Cross-E5 variants, some over-switching"),
        ("Philosophy", "Selector gains — improves Pk/WD"),
        ("Math",       "❌ Selector FAILS — only 4 training examples"),
    ]
    oy = 3.45
    for domain, finding in findings:
        box(s, 6.6, oy, 6.2, 0.42, LGREY)
        txt(s, domain, 6.7, oy+0.06, 1.4, 0.3, size=12, bold=True, color=NAVY)
        txt(s, finding, 8.15, oy+0.06, 4.5, 0.3, size=12,
            color=RED if "FAILS" in finding else DGREY)
        oy += 0.47

    txt(s, "Leave-domain-out result: Pk=0.4012 — worse than baseline\n"
        "→ Selector is a LOCAL operating point, not a universal deployment model",
        6.6, 6.08, 6.2, 0.7, size=12, italic=True, color=AMBER)


def slide_unsup_vs_sup(p):
    s = blank(p)
    header(s, "Unsupervised vs Supervised — Where LECSEG Sits",
           "The pipeline is mostly unsupervised; the selector adds a thin supervised layer")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    compare_table(
        s,
        headers=["Property", "LECSEG Segmentation\n(N1–N4)", "LECSEG Selector\n(Part 2)", "MiniSeg/Chapter-Llama\n(external)"],
        rows=[
            ["Requires labelled\ntraining videos", "❌ None needed", "⚠️ 29 videos (LOO)", "✅ 10k–19k videos"],
            ["Training objective", "N/A", "Predict Pk from\nvideo features", "Boundary F1 or\ngeneration loss"],
            ["Inference needs\ngold chapter count", "⚠️ Yes (controlled eval)", "⚠️ Yes", "❌ No"],
            ["Reproducible locally\n(no API / GPU farm)", "✅ Yes (CPU ok)", "✅ Yes", "❌ No"],
            ["Evaluation metric", "Pk / WD / BS / F1", "Pk / WD / BS / F1", "Different metrics\n(IoU, CIDEr…)"],
            ["Scale of dataset", "30 videos", "30 videos (LOO)", "10k–800k videos"],
        ],
        x=0.3, y=1.45, w=12.7,
        hdr_size=12, row_size=12,
    )

    txt(s, "Key claim: LECSEG is a LOW-RESOURCE UNSUPERVISED benchmark, not a competitor to large supervised systems.\n"
        "Its value is reproducibility, hierarchy, statistical testing, and the oracle-gap diagnosis — not external leaderboard dominance.",
        0.3, 6.15, 12.7, 0.7, size=13, bold=True, color=NAVY)


def slide_related_low_resource(p):
    s = blank(p)
    header(s, "Related Work — Low-Resource & Unsupervised Methods",
           "These are the closest comparable systems to LECSEG")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    compare_table(
        s,
        headers=["System", "Videos", "Supervision", "Modality", "Key Metric", "vs LECSEG"],
        rows=[
            ["TreeSeg\n(Gklezakos 2024)",
             "21 (TinyRec)\n+ICSI+AMI",
             "Unsupervised\ndivisive clustering",
             "Text only\n(overlapping blocks)",
             "Pk=0.367 (TinyRec)\nPk=0.310 (ICSI)",
             "Closest comparator.\nDiff dataset — no shared benchmark yet"],
            ["Freisinger 2023\nMultilingual",
             "34 (Videoaula)\n96 (LectureDE)",
             "Unsupervised rule-based\n+ LLM (multilingual)",
             "Text + speech pauses",
             "F1=67.3 (Videoaula)\n(no Pk/WD reported)",
             "Different metrics;\nnon-English focus"],
            ["AVLectures\n(Singh 2022)",
             "2350+",
             "Self-supervised\n(no chapter labels)",
             "Audio + visual\n(narration alignment)",
             "Outperforms baselines\n(no Pk reported)",
             "Much larger scale;\nno Pk/WD comparison"],
            ["Tuna 2015\nClassroom",
             "Small classroom set",
             "Unsupervised\nmultimodal",
             "Video + audio + slides",
             "Domain-specific\nsegmentation quality",
             "Similar setting;\npre-SBERT embeddings"],
            ["Chand & Ogul\n2021",
             "Small lecture set",
             "Unsupervised framework",
             "Lecture video\n(audio + visual)",
             "Domain-specific\nmetrics",
             "No reproducible\nbenchmark or Pk/WD"],
        ],
        x=0.3, y=1.4, w=12.7,
        hdr_size=12, row_size=11,
    )

    txt(s,
        "LECSEG advantage over these: explicit Pk/WD/BS evaluation · hierarchical labels · bootstrap CIs · oracle gap analysis · fully reproducible",
        0.3, 6.55, 12.7, 0.4, size=12, italic=True, color=LBLUE)


def slide_related_high_resource(p):
    s = blank(p)
    header(s, "Related Work — High-Resource Systems (Context Only)",
           "These are NOT direct competitors — different scale, metrics, and supervision")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    compare_table(
        s,
        headers=["System", "Videos", "Supervision", "Reported Result", "Why Not Directly Comparable"],
        rows=[
            ["MiniSeg / YTSEG\n(Retkowski 2024)",
             "19,299",
             "Supervised\nsmart-chaptering",
             "Pk=28.7, F1=43.4\n(YTSEG benchmark)",
             "643× more videos; different benchmark;\nsupervised vs unsupervised"],
            ["Chapter-Gen\n(Cao 2022)",
             "9,631",
             "Supervised\nmultimodal",
             "AP=43.3, Recall@3s=60.1",
             "321× more videos; title-generation task;\nAP/Recall not Pk/WD"],
            ["VidChapters-7M\n(Yang 2023)",
             "817,000",
             "Large-scale\npre-training",
             "SODA_c=11.4\n(generation metric)",
             "27,233× more videos;\ngeneration not segmentation"],
            ["Chapter-Llama\n(Ventura 2025, CVPR)",
             "10k train\n8.1k test",
             "Trained long-context\nLLM",
             "Chapter F1=45.3\n(vs Vid2Seq 26.7)",
             "333× more videos; requires\nlong-context LLM fine-tuning"],
        ],
        x=0.3, y=1.4, w=12.7,
        hdr_size=13, row_size=12,
    )

    box(s, 0.3, 5.9, 12.7, 1.35, NAVY)
    txt(s,
        "LECSEG's positioning:  \"A 30-video, fully reproducible benchmark proves that statistically supported improvements\n"
        "are achievable at lecture-specific scale — and precisely locates what a future larger system needs to solve.\"\n"
        "We do NOT claim to beat these systems. We show what's possible and what's left to do with limited resources.",
        0.5, 5.95, 12.3, 1.2, size=13, color=WHITE)


def slide_how_reproduced(p):
    s = blank(p)
    header(s, "How We Reproduced Baselines",
           "All 11 methods implemented from scratch in src/lecseg/ — no black-box imports")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    methods = [
        ("TextTiling", "NLTK reference implementation + custom fallback",
         "hearst1994multi", "185 passing tests"),
        ("C99", "Rank-transform cosine matrix, divisive clustering",
         "choi2000latent", "Validated: Pk≈0.51 on Choi (published 0.12 — our lite reimpl)"),
        ("CosineSeg", "Sentence-transformer vectors + depth-score valleys",
         "koshorek2018text", "BGE-large: Pk=0.490 on LECSEG-30"),
        ("KMeansSeg", "K-means clustering on embeddings, label-transition boundaries",
         "sklearn", "Pk=0.617 — over-segments"),
        ("BertSeg", "Same as CosineSeg, window=5 — approximates SegBot inference",
         "li2018segbot", "Pk=0.489 on LECSEG-30"),
        ("Divisive (BGE)", "Recursive max-cosine-dissimilarity split",
         "gklezakos2024_treeseg", "Pk=0.388 — best text-only baseline"),
    ]

    box(s, 0.3, 1.4, 12.7, 0.45, WHITE)
    txt(s, "Method", 0.4, 1.47, 2.5, 0.32, size=13, bold=True, color=NAVY)
    txt(s, "What we implemented", 3.0, 1.47, 3.8, 0.32, size=13, bold=True, color=NAVY)
    txt(s, "Based on paper", 6.9, 1.47, 2.5, 0.32, size=13, bold=True, color=NAVY)
    txt(s, "Validation", 9.5, 1.47, 3.6, 0.32, size=13, bold=True, color=NAVY)

    for i, (name, impl, paper, validation) in enumerate(methods):
        bg = LGREY if i % 2 == 0 else WHITE
        gy = 1.92 + i * 0.78
        box(s, 0.3, gy, 12.7, 0.75, bg)
        txt(s, name, 0.4, gy+0.15, 2.5, 0.45, size=13, bold=True, color=NAVY)
        txt(s, impl, 3.0, gy+0.08, 3.8, 0.58, size=12, color=DGREY)
        txt(s, paper, 6.9, gy+0.15, 2.5, 0.45, size=11, italic=True, color=LBLUE)
        txt(s, validation, 9.5, gy+0.08, 3.6, 0.58, size=12, color=DGREY)

    txt(s, "All implementations in src/lecseg/  ·  185 automated unit tests  ·  validated against published baselines where possible",
        0.3, 6.65, 12.7, 0.38, size=12, italic=True, color=LBLUE)


def slide_main_results(p):
    s = blank(p)
    header(s, "Main Results", "Best method reduces Pk by 0.0296 absolute over strong BGE baseline (p < 0.05)")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    compare_table(
        s,
        headers=["Method", "Type", "Pk ↓", "WD ↓", "BS ↑", "F1@2 ↑", "Sig. vs baseline"],
        rows=[
            ["TextTiling", "Classical", "0.605", "0.898", "0.147", "0.139", "❌ Much worse"],
            ["C99", "Classical", "0.422", "0.449", "0.035", "0.029", "❌ Worse"],
            ["CosineSeg", "Neural", "0.490", "0.539", "0.062", "0.086", "❌ Worse"],
            ["BertSeg", "Neural", "0.489", "0.540", "0.077", "0.091", "❌ Worse"],
            ["BGE-Divisive", "★ Baseline", "0.388", "0.396", "0.129", "0.088", "— Reference"],
            ["Cross-Model (E5)", "★ Best Global", "0.371", "0.376", "0.036", "0.024", "✅ p<0.05 Pk & WD"],
            ["Balanced Selector", "★ Best Point", "0.359", "0.374", "0.076", "0.089", "✅ p<0.05 BS & F1"],
        ],
        x=0.3, y=1.45, w=12.7,
        hdr_size=13, row_size=12,
    )

    txt(s, "Remember: Lower Pk/WD is better  ·  Higher BS/F1 is better  ·  0 = perfect  ·  0.5 = random guessing",
        0.3, 5.85, 12.7, 0.38, size=13, bold=True, color=NAVY)

    txt(s,
        "Cross-model conservative: best Pk/WD, but very conservative (low F1 — predicts fewer boundaries)\n"
        "Balanced selector: best mean operating point — improves F1 while staying competitive on Pk/WD",
        0.3, 6.25, 12.7, 0.6, size=12, color=DGREY)


def slide_granularity(p):
    s = blank(p)
    header(s, "Key Finding: Granularity Mismatch",
           "Why do acoustic & linguistic signals fail to beat text embeddings?")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Left: the finding
    box(s, 0.3, 1.4, 6.2, 5.8, WHITE)
    accent_bar(s, 0.3, 1.4, 6.2, 0.07, RED)
    txt(s, "The Granularity Problem", 0.5, 1.55, 5.9, 0.42, size=15, bold=True, color=NAVY)

    signals = [
        ("Discourse markers", "Pk=0.462", "Marks sentence-level transitions"),
        ("BERTopic", "Pk=0.563", "Finds subtopic clusters, not chapters"),
        ("GPT-2 perplexity", "Pk=0.418", "Detects sentence-level surprises"),
        ("Pause/pitch (prosody)", "Pk=0.417", "Speaker breath & stress resets"),
        ("BGE text embed", "Pk=0.388", "✅ Chapter-level semantic shift"),
        ("CLIP visual", "Pk=0.396", "✅ Slide scene changes ~ editorial"),
    ]
    oy = 2.05
    for sig, pk, reason in signals:
        is_good = pk in ["Pk=0.388", "Pk=0.396"]
        bg = RGBColor(0xE8,0xF5,0xE9) if is_good else RGBColor(0xFF,0xEE,0xEE)
        box(s, 0.5, oy, 5.8, 0.6, bg)
        txt(s, sig, 0.6, oy+0.08, 2.0, 0.4, size=13, color=NAVY)
        txt(s, pk, 2.65, oy+0.08, 1.0, 0.4, size=13, bold=True,
            color=GREEN if is_good else RED)
        txt(s, reason, 3.7, oy+0.08, 2.5, 0.4, size=11, color=DGREY)
        oy += 0.67

    # Right: explanation
    box(s, 6.8, 1.4, 6.2, 5.8, WHITE)
    accent_bar(s, 6.8, 1.4, 6.2, 0.07, NAVY)
    txt(s, "Why?", 7.0, 1.55, 5.9, 0.42, size=15, bold=True, color=NAVY)
    txt(s,
        "YouTube chapters = deliberate EDITORIAL choices\n"
        "made by instructors at a COARSE level:\n\n"
        "  'Introduction to Eigenvalues'\n"
        "  'The Characteristic Equation'\n"
        "  'Applications in Physics'\n\n"
        "Acoustic signals detect FINE-GRAINED transitions:\n"
        "  • Every pause > 300ms\n"
        "  • Every pitch reset\n"
        "  • Every discourse connector (\"So\", \"Now\", \"Next\")\n\n"
        "These happen 5–10× more often than chapter shifts.\n"
        "→ They over-segment relative to the reference.\n\n"
        "CLIP exception: slide scene changes ARE editorial\n"
        "decisions (new slide = potential chapter).\n"
        "→ CLIP matches chapter granularity better.",
        7.0, 2.05, 5.9, 4.5, size=13, color=DGREY)

    txt(s, "Practical implication: future systems should prioritise CLIP visual + text,\n"
        "not acoustic or linguistic sub-sentence signals for chapter-level segmentation.",
        6.8, 6.28, 6.2, 0.7, size=12, bold=True, color=LBLUE)


def slide_oracle(p):
    s = blank(p)
    header(s, "Oracle Gap Analysis — What's the Ceiling?",
           "Shown to locate the bottleneck, NOT as a performance claim")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    # Why show this
    box(s, 0.3, 1.4, 5.5, 2.8, RGBColor(0xFF,0xF3,0xCD))
    txt(s, "⚠️  Why We Show This", 0.5, 1.55, 5.2, 0.4, size=14, bold=True, color=AMBER)
    txt(s,
        "NOT a result we claim to achieve.\n"
        "Shown as a diagnostic: to answer\n"
        "'Where exactly is the remaining gap?'\n\n"
        "Answer: entirely in SELECTION,\n"
        "not in candidate generation.",
        0.5, 2.0, 5.2, 2.0, size=13, color=DGREY)

    # Table
    compare_table(
        s,
        headers=["Operating Point", "Pk ↓", "Deployable?"],
        rows=[
            ["BGE-Divisive (baseline)", "0.388", "✅ Yes"],
            ["Cross-model conservative", "0.371", "✅ Yes"],
            ["Balanced selector", "0.359", "✅ Yes (LOO)"],
            ["TreeSeg TinyRec (indicative)", "0.367", "✅ (diff dataset)"],
            ["Per-video oracle (best method known)", "0.298", "❌ Needs ground truth"],
            ["Candidate oracle τ=2 (perfect selector)", "0.017", "❌ Theoretical max"],
        ],
        x=6.1, y=1.45, w=6.9,
        hdr_size=12, row_size=12,
    )

    # Gap arrow
    box(s, 0.3, 4.4, 5.5, 2.85, WHITE)
    accent_bar(s, 0.3, 4.4, 5.5, 0.07, LBLUE)
    txt(s, "The Gap ΔPk ≈ 0.34", 0.5, 4.55, 5.2, 0.4, size=14, bold=True, color=NAVY)
    txt(s,
        "96.8% recall already exists in the candidate pool.\n"
        "That means the RIGHT boundaries are already found —\n"
        "we just can't always pick the correct ones.\n\n"
        "To close this gap:\n"
        "  → Supervised boundary ranker (need ~50 labelled lectures)\n"
        "  → Use CLIP + pause + OCR as ranker features\n"
        "  → NOT a better segmentation model",
        0.5, 4.98, 5.2, 2.0, size=13, color=DGREY)


def slide_limitations(p):
    s = blank(p)
    header(s, "Limitations & Honest Assessment", "What we cannot claim and what comes next")
    box(s, 0, 1.15, 13.33, 6.35, LGREY)

    limits = [
        ("Corpus size", "30 videos is small. Math domain has only 4 videos → selector fails there.",
         "Add 8–10 Math lectures; expand to 100+ for ranker training"),
        ("Silver chapter labels", "YouTube chapters = creator navigation choices, NOT expert pedagogical labels.\n"
         "Creators may place chapters late, skip boundaries, or be inconsistent.",
         "Independent human annotation study (future work)"),
        ("Low strict F1", "Best Pk/WD method predicts few boundaries → very low exact F1.\n"
         "Approximate segment structure ≠ exact boundary hit.",
         "Boundary ranker to improve precision without hurting Pk"),
        ("Selector is not universal", "Leave-domain-out Pk=0.4012 — worse than baseline.\n"
         "Only reliable when related domains are in training fold.",
         "More training data; domain-invariant features"),
        ("LLM refinement unverified", "Llama-3.1-8B boundary shifting not shown to improve Pk significantly.\n"
         "Titling works; refinement is diagnostic only.",
         "Larger LLM (70B) or task-specific fine-tuning"),
    ]

    for i, (title, problem, fix) in enumerate(limits):
        gy = 1.4 + i * 1.05
        box(s, 0.3, gy, 12.7, 0.98, WHITE)
        accent_bar(s, 0.3, gy, 12.7, 0.06, [RED, AMBER, AMBER, AMBER, LBLUE][i])
        txt(s, title, 0.5, gy+0.1, 2.3, 0.35, size=13, bold=True, color=NAVY)
        txt(s, problem, 2.8, gy+0.07, 6.5, 0.55, size=12, color=DGREY)
        box(s, 9.4, gy+0.07, 3.5, 0.55, LGREY)
        txt(s, "→  " + fix, 9.5, gy+0.1, 3.3, 0.48, size=11, color=TEAL)


def slide_conclusion(p):
    s = blank(p)
    box(s, 0, 0, 13.33, 7.5, NAVY)
    for cx, cy, r in [(11.5,1.2,2.2),(12.8,5.8,1.8),(0.5,6.5,1.0)]:
        b = s.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x1A,0x4A,0x8A); b.line.fill.background()

    txt(s, "What LECSEG Contributes", 0.8, 0.6, 12, 0.7, size=32, bold=True, color=WHITE)
    accent_bar(s, 0.8, 1.35, 8, 0.07, TEAL)

    contribs = [
        ("📦", "LECSEG-30 benchmark",
         "30 lectures · 5 domains · 32h · 419 chapters · 904 subtopics\nHierarchical labels with IAA · fully reproducible"),
        ("📊", "Rigorous evaluation",
         "5 metrics · bootstrap 95% CIs · paired Wilcoxon tests\nStatistically supported improvement: ΔPk=0.029 over BGE-divisive"),
        ("🔍", "Granularity insight",
         "Proven that acoustic/linguistic signals over-segment vs YouTube chapters\nCLIP visual is the most promising non-text signal"),
        ("🗺️", "Oracle gap diagnosis",
         "96.8% recall in candidate pool — bottleneck is SELECTION not generation\nClear roadmap: supervised boundary ranker needs ~50 labelled lectures"),
    ]
    oy = 1.55
    for icon, title, desc in contribs:
        txt(s, icon, 0.8, oy, 0.6, 0.7, size=26, color=WHITE)
        txt(s, title, 1.5, oy+0.02, 4.5, 0.38, size=17, bold=True, color=GOLD)
        txt(s, desc, 1.5, oy+0.4, 11, 0.55, size=13, color=RGBColor(0xC0,0xD8,0xF0))
        oy += 1.12

    txt(s,
        "LECSEG does not claim to beat large supervised systems. It turns low-resource lecture segmentation\n"
        "into a MEASURABLE, REPRODUCIBLE problem with clear next steps.",
        0.8, 6.2, 12, 0.9, size=14, italic=True, color=RGBColor(0x80,0xB4,0xE0))


# ── Assemble ───────────────────────────────────────────────────────────────

def main():
    p = prs()
    slide_title(p)
    slide_problem(p)
    slide_gap(p)
    slide_dataset(p)
    slide_annotation(p)
    slide_metrics(p)
    slide_pipeline(p)
    slide_methods_unsupervised(p)
    slide_novel_n1(p)
    slide_novel_n2(p)
    slide_novel_n3n4(p)
    slide_selector(p)
    slide_unsup_vs_sup(p)
    slide_related_low_resource(p)
    slide_related_high_resource(p)
    slide_how_reproduced(p)
    slide_main_results(p)
    slide_granularity(p)
    slide_oracle(p)
    slide_limitations(p)
    slide_conclusion(p)

    out = "thesis/LECSEG_Defense_Slides_v2.pptx"
    p.save(out)
    print(f"Done. Saved {len(p.slides)} slides -> {out}")


if __name__ == "__main__":
    main()
